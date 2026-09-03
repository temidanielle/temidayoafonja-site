# -*- coding: utf-8 -*-
"""Rewrite the speaker notes of both Video 5 decks to v3.0.

Only notesSlideN.xml parts may change. Every other part -- slide XML, media,
rels, theme, presentation.xml -- must be byte-identical afterwards."""
import sys, shutil, zipfile, hashlib, os, copy
sys.path.insert(0,"/tmp/v5v31")
from notes_v31 import NOTES, reveal_notes
from pptx import Presentation
from pptx.oxml.ns import qn

OUT="/home/user/temidayoafonja-site/deliverables/video-5-slides/out"
MAIN=os.path.join(OUT,"Video_5_Main_Slides.pptx")
REV=os.path.join(OUT,"Video_5_Reveal_Builds.pptx")

def parts(path):
    with zipfile.ZipFile(path) as z:
        return {i.filename: hashlib.sha256(z.read(i.filename)).hexdigest()
                for i in z.infolist()}

def set_notes(slide, text):
    """Replace the notes text, cloning the first run's properties so type
    treatment is preserved. Touches only the notesSlide part."""
    tf = slide.notes_slide.notes_text_frame
    src = None
    for p in tf.paragraphs:
        for r in p.runs:
            src = r._r.find(qn('a:rPr')); break
        if src is not None: break
    body = tf._txBody
    for p in body.findall(qn('a:p')): body.remove(p)
    for line in text.split("\n\n"):
        p = body.makeelement(qn('a:p'), {})
        body.append(p)
        r = p.makeelement(qn('a:r'), {}); p.append(r)
        if src is not None: r.append(copy.deepcopy(src))
        t = p.makeelement(qn('a:t'), {}); t.text = line.replace("\n", " ")
        r.append(t)

def run(path, notes):
    before = parts(path)
    pres = Presentation(path)
    assert len(pres.slides) == len(notes), (len(pres.slides), len(notes))
    for s, n in zip(pres.slides, notes):
        set_notes(s, n)
    pres.save(path)
    after = parts(path)
    changed = sorted(k for k in after if before.get(k) != after[k])
    added = sorted(set(after) - set(before))
    removed = sorted(set(before) - set(after))
    return changed, added, removed

for path, notes in ((MAIN, NOTES), (REV, reveal_notes())):
    ch, ad, rm = run(path, notes)
    print("\n==", os.path.basename(path))
    print("  changed parts (%d):" % len(ch))
    for c in ch: print("     ", c)
    print("  added:", ad, " removed:", rm)
    bad = [c for c in ch if not c.startswith("ppt/notesSlides/notesSlide")
           and c != "docProps/app.xml" and c != "docProps/core.xml"]
    print("  NON-NOTES PARTS CHANGED:", bad)
