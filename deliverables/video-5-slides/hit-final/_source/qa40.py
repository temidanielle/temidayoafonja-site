# -*- coding: utf-8 -*-
"""Videos 1-5 v4.0 QA: canonical, direct-address, hook, memory, factual,
slide/notes, document and checksum checks."""
import os, sys, re, json, zipfile, subprocess, glob
sys.path.insert(0,"/tmp/da")
from changereport import old_spoken, new_blocks, report, scan_detached, viewer_ratio
from docxkit import sha256, hype_scan
from docx import Document
from pptx import Presentation

R="/home/user/temidayoafonja-site/deliverables"
CFG={
1:{"base":"/tmp/v4/v1","prior":R+"/video-1-slides/hit-final/Video_1_HIT_FINAL/LONG_FORM/Video1TeleprompterScriptwithslidemarkers_HIT_v3.1.txt",
   "prior_ver":"v3.1","prior_zip":"17e881ea97774f0d4a9e080f2077b093b6367f6f3ce14e22fe119ceb17a793e6",
   "main":R+"/video-1-slides/out/Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pptx",
   "rev":R+"/video-1-slides/out/Video-1-Reveal-Builds_v2.4.pptx",
   "pdf":R+"/video-1-slides/out/Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pdf",
   "slides":13,"frames":22,"markers":13,
   "title":"How to Change Jobs Without Starting Your Career Over",
   "thumb":"DON’T START FROM ZERO","cta":"Free Career Evidence Starter",
   "next":"Is Your Job Making You Less Marketable?",
   "forbidden":["fieldkit","Keep the Proof","career-decisions","3 Cs"],
   "memory":["Look underneath the title","Explain what the work changed",
             "Keep evidence before you need it"]},
2:{"base":"/tmp/da/v2","prior":R+"/video-2-slides/hit-final/Video_2_HIT_FINAL/LONG_FORM/Video2TeleprompterScriptwithslidemarkers_HIT_v3.0_DIRECT_ADDRESS.txt",
   "prior_ver":"v3.0 direct address","prior_zip":"f8ebaa45f657d5fbd60440a54bde58127c343b15fd0b8956c84d6cf7701e18a9",
   "main":R+"/video-2-slides/out/Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.1.pptx",
   "rev":R+"/video-2-slides/out/Video-2-Reveal-Builds_v1.1.pptx",
   "pdf":R+"/video-2-slides/out/Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.1.pdf",
   "slides":13,"frames":23,"markers":13,
   "title":"Is Your Job Making You Less Marketable?",
   "thumb":"VALUABLE HERE. STUCK HERE?","cta":"Capability Formation Field Kit",
   "next":"3 Things to Do Before Quitting Your Job",
   "forbidden":["Keep the Proof","Career Evidence Starter","career-decisions",
                "YOUR SKILLS ARE STALLING"],
   "memory":["Remove the company nouns","Find outside-context evidence",
             "Read the last 90 days"]},
3:{"base":"/tmp/da/v3","prior":R+"/video-3-slides/hit-final/Video_3_HIT_FINAL/LONG_FORM/Video3TeleprompterScriptwithslidemarkers_HIT_v3.0_DIRECT_ADDRESS.txt",
   "prior_ver":"v3.0 direct address","prior_zip":"62af1ca5d1c2096a61309d7d0529e761c07af2237521c2e18ad51c821e890874",
   "main":R+"/video-3-slides/out/Video-3-3-Things-to-Do-Before-Quitting-Your-Job_v1.1.pptx",
   "rev":R+"/video-3-slides/out/Video-3-Reveal-Builds_v1.1.pptx",
   "pdf":R+"/video-3-slides/out/Video-3-3-Things-to-Do-Before-Quitting-Your-Job_v1.1.pdf",
   "slides":13,"frames":27,"markers":13,
   "title":"3 Things to Do Before Quitting Your Job",
   "thumb":"WAIT BEFORE YOU QUIT","cta":"Career Decision Evidence Check",
   "next":"How to Change Jobs Without Starting Your Career Over",
   "forbidden":["Keep the Proof","Career Evidence Starter","fieldkit"],
   "memory":["Preserve the evidence","Name what the work built",
             "Test the next move"]},
4:{"base":"/tmp/da/v4","prior":R+"/video-4-slides/hit-final/Video_4_HIT_FINAL/LONG_FORM/Video4TeleprompterScriptwithslidemarkers_HIT_v3.0_DIRECT_ADDRESS.txt",
   "prior_ver":"v3.0 direct address","prior_zip":"600bbf407b7d4c7da2d2339d16c9670896c310894588f836db4a8b627cd6a65f",
   "main":R+"/video-4-slides/out/Video_4_Main_Slides.pptx",
   "rev":R+"/video-4-slides/out/Video_4_Reveal_Builds.pptx",
   "pdf":R+"/video-4-slides/out/Video_4_Slide_Preview.pdf",
   "slides":11,"frames":26,"markers":11,
   "title":"How to Explain Your Career Change",
   "thumb":"YOUR CAREER MAKES SENSE","cta":"Free Career Evidence Starter",
   "next":"Should I Make an Internal Move? 3 Questions to Decide",
   "forbidden":["Keep the Proof","fieldkit","career-decisions"],
   "memory":["My career has moved across","I kept being asked to",
             "That is why I am now focused on"]},
5:{"base":"/tmp/v4/v5","prior":R+"/video-5-slides/hit-final/Video_5_HIT_FINAL/LONG_FORM/Video5TeleprompterScriptwithslidemarkers_HIT_v3.1.txt",
   "prior_ver":"v3.1","prior_zip":"0067d5530c26d2625eec0bcce131ab61cfadba6f438f5febff30d557c04d2b23",
   "main":R+"/video-5-slides/out/Video_5_Main_Slides.pptx",
   "rev":R+"/video-5-slides/out/Video_5_Reveal_Builds.pptx",
   "pdf":R+"/video-5-slides/out/Video_5_Slide_Preview.pdf",
   "slides":12,"frames":25,"markers":12,
   "title":"Should I Make an Internal Move? 3 Questions to Decide",
   "thumb":"YOU MAY NOT NEED TO LEAVE","cta":"Career Decision Evidence Check",
   "next":"Are You Growing—or Just Being Given More Work?",
   "forbidden":["Keep the Proof","fieldkit","Career Evidence Starter","3 Cs"],
   "memory":["Will the work change","Will your judgment expand",
             "Will the evidence travel"]},
}
def text(p): return "\n".join(x.text for x in Document(p).paragraphs)
def paras(p): return [x.text.strip() for x in Document(p).paragraphs if x.text.strip()]
def block(t):
    a=t.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
    b=t.index("— END OF")
    return [x.strip() for x in t[a:b].split("\n") if x.strip()][1:]

OUT={}
for n,c in CFG.items():
    B=c["base"]; ROOT=os.path.join(B,"Video_%d_HIT_FINAL"%n)
    LF=os.path.join(ROOT,"LONG_FORM"); SH=os.path.join(ROOT,"SHORTS")
    res={}; F=[]
    def chk(k,l,cond,d=""):
        res[k]={"check":l,"pass":bool(cond),"detail":str(d)}
        if not cond: F.append((k,l,str(d)))
    canon=os.path.join(B,"canonical_v4.0.txt")
    nb=new_blocks(canon,"BEGIN APPROVED VIDEO %d v4.0 SCRIPT"%n,
                        "END APPROVED VIDEO %d v4.0 SCRIPT"%n)
    new=[x for x in nb if not x.startswith("[SLIDE:")]
    mk=[x for x in nb if x.startswith("[SLIDE:")]
    old=old_spoken(c["prior"]); rep=report(old,new)
    spoken="\n\n".join(new); words=sum(len(x.split()) for x in new)
    tel=glob.glob(LF+"/Video%dTeleprompter*_v4.0"%n+".txt")[0][:-4]
    rdg=glob.glob(LF+"/Video%dReading*_v4.0"%n+".txt")[0][:-4]
    edb=glob.glob(LF+"/*EDITOR_ONLY*v4.0.docx")[0]
    pub=glob.glob(LF+"/*Publishing*v4.0.docx")[0]
    seb=glob.glob(SH+"/*Shorts_EDITOR_ONLY*.docx")[0]
    desc=os.path.join(B,"Video_%d_YouTube_Description_HIT.docx"%n)
    zp=os.path.join(B,"Video_%d_HIT_FINAL_Recording_and_Shorts_Package.zip"%n)
    shorts={f:text(os.path.join(SH,f)) for f in sorted(os.listdir(SH))
            if f.startswith("Video_%d_Short_"%n)}
    telt=open(tel+".txt",encoding="utf-8").read(); rdgt=open(rdg+".txt",encoding="utf-8").read()
    pubt=text(pub); desct=text(desc); edbt=text(edb); sebt=text(seb)
    allpub=pubt+desct
    telb=[x.strip() for x in telt.split("\n\n") if x.strip()][1:]
    telnorm=["[SLIDE: %s]"%x[len("SLIDE  —  "):] if x.startswith("SLIDE  —  ") else x for x in telb]
    chk("C1","Teleprompter TXT == v4.0 canonical", telnorm==nb)
    chk("C2","Reading TXT == canonical spoken",
        [x.strip() for x in rdgt.split("\n\n") if x.strip()]==new)
    chk("C3","Teleprompter minus markers == reading script",
        [x for x in telnorm if not x.startswith("[SLIDE:")]==new)
    teld=paras(tel+".docx")[4:]
    chk("C4","Teleprompter DOCX == TXT",
        ["[SLIDE: %s]"%x[len("SLIDE  —  "):] if x.startswith("SLIDE  —  ") else x
         for x in teld]==nb)
    chk("C5","Reading DOCX == TXT", paras(rdg+".docx")[4:]==new)
    chk("C6","Marker count preserved", len(mk)==c["markers"], len(mk))
    chk("C7","No paragraph removed", rep["removed"]==0, rep["removed"])
    chk("C8","No unsupported factual claim added",
        not re.findall(r"\b\d+\s?%|\$\s?\d", spoken) and
        not re.search(r"\b30\s?%|\$2M|2 ?million", spoken, re.I))
    chk("D1","No detached phrasing",
        not scan_detached(spoken) and not any(scan_detached(t) for t in shorts.values()),
        scan_detached(spoken))
    vr=viewer_ratio(new)
    chk("D2","Viewer-facing second person", vr["ratio"]>=0.70,
        "%d/%d = %.0f%%"%(vr["second_person"],vr["viewer_facing"],100*vr["ratio"]))
    low=spoken.lower()
    br={b:low.count(b) for b in ["let me show you","i want to help you",
        "i want you to","here is what i would","i want to walk you through"]}
    chk("D3","First-person bridges present", sum(br.values())>=3,
        {k:v for k,v in br.items() if v})
    ep=spoken.count("experienced professionals")
    chk("D4","'experienced professionals' is positioning", ep<=1, ep)
    allsec=len([p for p in new if re.search(r"\b(you|your)\b",p,re.I)])/len(new)
    chk("D5","Not mechanically repetitive", allsec<=0.90, "%.0f%%"%(100*allsec))
    chk("D6","No keynote tone",
        not any(k in low for k in ["ladies and gentlemen","in this presentation",
        "welcome back to my channel","hey everyone","hi guys","let's dive in"]))
    hy=hype_scan(spoken)+[h for t in shorts.values() for h in hype_scan(t)]
    chk("H1","No hype or creator-bait phrasing", not hy, hy)
    warn=[p for p in new if "warning" in p.lower()]
    chk("H2","'warning' occurrences justified",
        all("leave you with one warning" in p for p in warn),
        [p[:70] for p in warn] or "none")
    chk("H3","Editor brief names the first-30 audit table",
        "First-30-second audit table" in edbt and "Viewer payoff" in edbt)
    chk("H4","Editor brief names all five hook layers",
        all(x in edbt for x in ["Verbal hook","On-screen hook","Visual / editor",
                                "Trust beat","Viewer payoff"]))
    chk("H5","Shorts brief gives every Short a hook and a route",
        sebt.count("On-screen hook")>=4 and sebt.count("Related Video")>=4,
        "%d hooks, %d routes"%(sebt.count("On-screen hook"),sebt.count("Related Video")))
    chk("M1","One memory structure, stated in plain language",
        all(m.lower() in low for m in c["memory"]),
        [m for m in c["memory"] if m.lower() not in low])
    chk("M2","No second framework",
        not re.search(r"\bCAR\b",spoken) and "3 Cs" not in spoken)
    chk("T1","Title exact", c["title"] in allpub)
    chk("T2","Thumbnail exact", c["thumb"] in allpub)
    chk("T3","One primary CTA", c["cta"].lower() in allpub.lower() and
        not any(f.lower() in allpub.lower() for f in c["forbidden"]),
        [f for f in c["forbidden"] if f.lower() in allpub.lower()])
    chk("T4","Watch Next correct", c["next"] in allpub and c["next"] in spoken)
    chk("E1","Emoji system present",
        all(e in pubt for e in ["✨","🧭","⏱️","▶️","🔗"]) and
        all(e in desct for e in ["✨","🧭","⏱️","▶️","🔗"]))
    chk("E2","Tag field present in both", "YouTube tag field" in pubt
        and "YouTube tag field" in desct)
    chk("E3","No [INSERT] placeholder", "[INSERT" not in (pubt+desct+edbt+sebt))
    chk("E4","Public block identical in both docs", block(pubt)==block(desct))
    chk("E5","Working estimates flagged outside public copy",
        "WORKING ESTIMATES — EDITOR MUST REPLACE FROM FINAL CUT" in pubt and
        pubt.index("— END OF")<pubt.index("WORKING ESTIMATES"))
    ED=["Visual:","On-screen","B-roll","EDITOR","Reveal:","End on:","FACTUAL BOUNDARY"]
    chk("E6","Shorts recording docs are spoken copy only",
        not [(f,w) for f,t in shorts.items() for w in ED if w in t],
        [(f,w) for f,t in shorts.items() for w in ED if w in t])
    mp=Presentation(c["main"]); rp=Presentation(c["rev"])
    chk("S1","Main slide count", len(mp.slides)==c["slides"], len(mp.slides))
    chk("S2","Reveal frame count", len(rp.slides)==c["frames"], len(rp.slides))
    prov=json.load(open(os.path.join(B,"_partdiff.json")))
    chk("S3","No slide XML, media or rels change",
        prov["main_nonnotes"]==[] and prov["reveal_nonnotes"]==[],
        prov["main_nonnotes"]+prov["reveal_nonnotes"])
    notes=[s.notes_slide.notes_text_frame.text for s in mp.slides]+\
          [s.notes_slide.notes_text_frame.text for s in rp.slides]
    chk("S4","All notes parts carry v4.0 notes",
        all("Timing:" in x for x in notes) and
        not any(re.search(r"v[23]\.[01]",x) for x in notes),
        len(notes))
    names=zipfile.ZipFile(zp).namelist()
    chk("Z1","Exactly 13 files in ZIP", len(names)==13, len(names))
    chk("Z2","No _source, Python, QA or render files in ZIP",
        not any(x.endswith((".py",".pyc",".html",".png")) or "_source" in x for x in names))
    sums=[l for l in open(os.path.join(ROOT,"SHA256SUMS.txt"),encoding="utf-8")
          .read().strip().split("\n") if l.strip() and not l.startswith("#")]
    chk("Z3","SHA256SUMS exactly 12 entries", len(sums)==12, len(sums))
    pr=subprocess.run("sha256sum -c SHA256SUMS.txt",shell=True,cwd=ROOT,
                      capture_output=True,text=True)
    chk("Z4","sha256sum -c successful", pr.returncode==0, pr.stdout[-160:])
    zs=sha256(zp)
    chk("Z5","Sibling checksum matches",
        open(zp+".sha256",encoding="utf-8").read().split()[0]==zs)
    OUT[n]={"checks":res,"failures":F,
      "change":{k:v for k,v in rep.items() if k!="blocks"},
      "voice":{k:v for k,v in vr.items() if k!="not_second_person"},
      "words":words,"markers":len(mk),"slides":len(mp.slides),"frames":len(rp.slides),
      "prior_ver":c["prior_ver"],"prior_zip":c["prior_zip"],
      "prior_source":os.path.basename(c["prior"]),
      "notes_rewritten":{"main":len(prov["main_changed"]),
                         "reveal":len(prov["reveal_changed"])},
      "hashes":{"zip":zs,"desc":sha256(desc),"main":sha256(c["main"]),
                "reveal":sha256(c["rev"]),"pdf":sha256(c["pdf"])}}
    print("VIDEO %d: %d checks, %d failed"%(n,len(res),len(F)))
    for f in F: print("   FAIL",f)
json.dump(OUT,open("/tmp/v4/QA_REPORT_v4.json","w"),indent=1)
print("\nTOTAL %d checks, %d failed"%(sum(len(v["checks"]) for v in OUT.values()),
                                      sum(len(v["failures"]) for v in OUT.values())))
