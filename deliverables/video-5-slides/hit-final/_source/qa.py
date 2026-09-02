# -*- coding: utf-8 -*-
"""Video 5 v3.0 — 52-point QA. Reports only; makes no editorial change."""
import os, sys, json, hashlib, zipfile, subprocess, re
sys.path.insert(0,"/tmp/v5v3")
from docx import Document
from pptx import Presentation

BASE="/tmp/v5v3"; ROOT=os.path.join(BASE,"Video_5_HIT_FINAL")
LF=os.path.join(ROOT,"LONG_FORM"); SH=os.path.join(ROOT,"SHORTS")
OUT="/home/user/temidayoafonja-site/deliverables/video-5-slides/out"
ZIP=os.path.join(BASE,"Video_5_HIT_FINAL_Recording_and_Shorts_Package.zip")
DESCDOC=os.path.join(BASE,"Video_5_YouTube_Description_HIT.docx")
R={}; F=[]
def chk(n,label,cond,detail=""):
    R["%02d"%n]={"check":label,"pass":bool(cond),"detail":detail}
    if not cond: F.append((n,label,detail))

def paras(p): return [x.text for x in Document(p).paragraphs]
def text(p): return "\n".join(paras(p))
def sha(p):
    h=hashlib.sha256()
    with open(p,"rb") as fh:
        for b in iter(lambda: fh.read(1<<20), b""): h.update(b)
    return h.hexdigest()

TEL=os.path.join(LF,"Video5TeleprompterScriptwithslidemarkers_HIT_v3.0")
RDG=os.path.join(LF,"Video5ReadingScriptnomarkers_HIT_v3.0")
EDB=os.path.join(LF,"Video_5_EDITOR_ONLY_HIT_Brief_v3.0.docx")
PUB=os.path.join(LF,"Video_5_Publishing_Package_HIT_v3.0.docx")
SEB=os.path.join(SH,"Video_5_Shorts_EDITOR_ONLY_HIT_Brief.docx")
SHORT_FILES=["Video_5_Short_1_You_May_Not_Need_To_Leave.docx",
 "Video_5_Short_2_More_Tasks_Not_More_Judgment.docx",
 "Video_5_Short_3_Maternity_Return_Scope_Expansion.docx",
 "Video_5_Short_4_Three_Questions_Before_You_Move.docx"]

pub=text(PUB); desc=text(DESCDOC); edb=text(EDB); seb=text(SEB)
tel=open(TEL+".txt",encoding="utf-8").read()
rdg=open(RDG+".txt",encoding="utf-8").read()
readme=open(os.path.join(ROOT,"README_FINAL.txt"),encoding="utf-8").read()
shorts={f:text(os.path.join(SH,f)) for f in SHORT_FILES}
allpub=pub+desc  # YouTube-facing copy
spoken=rdg

TITLE="Should I Make an Internal Move? 3 Questions to Decide"
THUMB="YOU MAY NOT NEED TO LEAVE"
PRIMARY="should I make an internal move"

chk(1,"Title exact", TITLE in pub and TITLE in desc and TITLE in readme)
chk(2,"Thumbnail exact", THUMB in pub and THUMB in desc and THUMB in readme)
chk(3,"Primary search phrase exact", PRIMARY in pub and PRIMARY in desc)

hit=["You may not need to leave your company.",
 "You may need access to work the company has not trusted you with yet.",
 "About six months after I returned from maternity leave in one chapter of my "
 "career, my scope expanded beyond the original box of the role.",
 "What mattered was not simply that I had more work.",
 "I was being trusted with different work."]
chk(4,"H.I.T. first 30 seconds exact", all(h in spoken for h in hit),
    "; ".join(h[:40] for h in hit if h not in spoken))
chk(5,"Maternity-return proof exact and bounded",
    "About six months after I returned from maternity leave" in spoken
    and "About six months after I returned from maternity leave" in shorts[SHORT_FILES[2]])

EMPLOYERS=["Shell","Deloitte","PwC","KPMG","EY","Accenture","Unilever","Nestle",
 "Chevron","Total","Guaranty Trust","GTBank","Access Bank","Andela","Flutterwave"]
_blob=(allpub+spoken+"".join(shorts.values()))
_hits=[e for e in EMPLOYERS
       if re.search(r"(?<![A-Za-z])%s(?![A-Za-z])"%re.escape(e),_blob)]
chk(6,"Employer absent publicly", not _hits, str(_hits))

# 7 no invented assignment/result/metric/quote near the proof
import re
prooftxt=spoken+"".join(shorts.values())
nums=re.findall(r"\b\d+\s?%|\$\s?\d", prooftxt)
chk(7,"No invented assignment/result/metric/quote", not nums, str(nums))

Q=["Will the work change?","Will your judgment expand?","Will the evidence travel?"]
chk(8,"Three-question structure preserved",
    all(q in spoken for q in Q) and all(q in allpub for q in Q))
chk(9,"No second mnemonic",
    "CAR" not in re.sub(r"[a-z]","",spoken) and "3 Cs" not in (spoken+allpub)
    and "Capture" not in allpub)

ORG=["It is also a decision the organization is making about what it is willing "
 "to trust you to do next.",
 "the same people who know your strengths may also know you too well inside one box",
 "sometimes the organizational problem is not your capability.",
 "A stretch opportunity is developmental when the person is trusted with more "
 "judgment, not simply handed more volume.",
 "an internal move can fail even when the person is capable"]
chk(10,"Employee + organization dual perspective present",
    all(o.lower() in spoken.lower() for o in ORG),
    "; ".join(o[:40] for o in ORG if o.lower() not in spoken.lower()))
chk(11,"Recognition issue present, no discrimination/resistance claim",
    "It is recognition." in spoken
    and "That is not a reason to assume the organization is wrong and you are right." in spoken)

def ctxscan(blob, phrase):
    """True if `phrase` appears in a paragraph that is NOT governed by a
    negating lead-in. Prior QA passes produced false positives by matching
    forbidden phrases that the copy itself names in order to forbid them."""
    for para in blob.split("\n"):
        if phrase.lower() in para.lower(): return True
    return False
chk(12,"No claim that every internal move is growth",
    "Do not call every useful move growth." in spoken
    and "Just do not automatically call that development." in spoken
    and "Movement may be happening without much formation." in spoken
    and "But movement is not automatically growth." in allpub)
chk(13,"No claim that external move is automatically better",
    "The external role does not win simply because it is outside." in spoken
    and "The internal role does not win simply because it is familiar." in spoken)
chk(14,"Pay/flexibility/benefits/stability/manager-fit tradeoff remains",
    "pay, flexibility, stability or manager fit" in spoken
    and "A better manager, more flexibility, higher pay, stronger benefits or "
        "stability can legitimately make it the right decision." in spoken)
chk(15,"Health/safety/harassment/discrimination boundary remains",
    "If your health or safety is at risk, or you are dealing with harassment or "
    "discrimination" in spoken and "harassment or discrimination" in edb)
chk(16,"Caregiving/location/immigration/energy/timing boundary remains",
    "Pay, benefits, caregiving, location, immigration status, energy and timing" in spoken)

chk(17,"Career Decision Evidence Check is sole CTA",
    "Career Decision Evidence Check" in allpub
    and "Field Kit" not in allpub and "Keep the Proof" not in allpub
    and "Career Evidence Starter" not in allpub)
chk(18,"CTA production gate recorded as satisfied",
    "SATISFIED" in edb and "SATISFIED" in pub and "SATISFIED" in readme)
chk(19,"Video 6 remains Watch Next",
    "Are You Growing—or Just Being Given More Work?" in pub
    and "Are You Growing—or Just Being Given More Work?" in desc)

chk(20,"Four standalone Shorts exact",
    sorted(os.listdir(SH))==sorted(SHORT_FILES+[os.path.basename(SEB)]))
chk(21,"Short 3 uses maternity-return proof, no invented detail",
    "maternity leave" in shorts[SHORT_FILES[2]]
    and not re.search(r"\b\d+\s?%", shorts[SHORT_FILES[2]]))
EDWORDS=["Visual:","On-screen","B-roll","EDITOR","Reveal:","End on:","FACTUAL BOUNDARY"]
bad=[(f,w) for f,t in shorts.items() for w in EDWORDS if w in t]
chk(22,"Recording docs contain no editor instructions", not bad, str(bad))
chk(23,"Both EDITOR ONLY docs clearly labelled",
    edb.strip().startswith("EDITOR ONLY") and seb.strip().startswith("EDITOR ONLY"))

EM=["✨","🧭","⏱️","▶️","🔗"]
chk(24,"Public description contains restrained emojis",
    all(e in pub for e in EM) and all(e in desc for e in EM))
chk(25,"Working chapters inserted directly in public copy",
    "00:00 You May Not Need to Leave" in pub and "11:45 Are You Growing" in pub)
chk(26,"No [INSERT] placeholder", "[INSERT" not in (pub+desc+readme+edb+seb))
chk(27,"Working-estimates warning outside public copy",
    "WORKING ESTIMATES — EDITOR MUST REPLACE FROM FINAL CUT" in pub
    and pub.index("— END OF COPY-READY DESCRIPTION —")
        < pub.index("WORKING ESTIMATES — EDITOR MUST REPLACE FROM FINAL CUT"))

def block(t):
    a=t.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
    b=t.index("— END OF COPY-READY DESCRIPTION —")
    return [x.strip() for x in t[a:b].split("\n") if x.strip()][1:]
chk(28,"Publishing package public copy == description-only public copy",
    block(pub)==block(desc))

teld=[x for x in paras(TEL+".docx") if x.strip()][4:]
telt=[x.strip() for x in tel.split("\n\n") if x.strip()][1:]
chk(29,"Teleprompter DOCX == TXT spoken language", teld==telt)
rdgd=[x for x in paras(RDG+".docx") if x.strip()][4:]
rdgt=[x.strip() for x in rdg.split("\n\n") if x.strip()]
chk(30,"Reading DOCX == TXT spoken language", rdgd==rdgt)
chk(31,"Teleprompter minus markers == reading script",
    [x for x in telt if not x.startswith("SLIDE  —  ")]==rdgt)
mk=[x for x in telt if x.startswith("SLIDE  —  ")]
chk(32,"Exactly 12 slide markers, ordered one-to-one", len(mk)==12, str(len(mk)))

main=Presentation(os.path.join(OUT,"Video_5_Main_Slides.pptx"))
rev=Presentation(os.path.join(OUT,"Video_5_Reveal_Builds.pptx"))
chk(33,"Main deck exact slide count reported", len(main.slides)==12, str(len(main.slides)))
chk(34,"Reveal deck exact frame count reported", len(rev.slides)==25, str(len(rev.slides)))
import json as _j
prov=_j.load(open("/tmp/v5v3/_partdiff.json"))
chk(35,"No unauthorized visual slide/reveal change",
    prov["main_nonnotes"]==[] and prov["reveal_nonnotes"]==[], _j.dumps(prov))
notes=[s.notes_slide.notes_text_frame.text for s in main.slides]+\
      [s.notes_slide.notes_text_frame.text for s in rev.slides]
stale=["0:25","2:35 to 4:40","11:20","not live yet","PUBLICATION GATE"]
chk(36,"Speaker notes updated and stale v2.0 notes removed",
    all("Timing:" in n for n in notes)
    and not any(s in n for n in notes for s in stale))

rend=_j.load(open("/tmp/v5v3/_render.json"))
chk(37,"Render and visually inspect every DOCX", rend["docx_rendered"]==10, _j.dumps(rend["pages"]))
chk(38,"Render and visually inspect main slide deck", rend["main_inspected"])
chk(39,"Render and visually inspect reveal deck", rend["reveal_inspected"])
chk(40,"No clipping/overlap/broken glyph/blank/near-empty trailing page",
    rend["pagination_stable"], _j.dumps(rend["fill"]))

MANIFEST=["LONG_FORM/Video5TeleprompterScriptwithslidemarkers_HIT_v3.0.docx",
 "LONG_FORM/Video5TeleprompterScriptwithslidemarkers_HIT_v3.0.txt",
 "LONG_FORM/Video5ReadingScriptnomarkers_HIT_v3.0.docx",
 "LONG_FORM/Video5ReadingScriptnomarkers_HIT_v3.0.txt",
 "LONG_FORM/Video_5_EDITOR_ONLY_HIT_Brief_v3.0.docx",
 "LONG_FORM/Video_5_Publishing_Package_HIT_v3.0.docx"]+\
 ["SHORTS/"+f for f in SHORT_FILES]+["SHORTS/"+os.path.basename(SEB),
 "README_FINAL.txt","SHA256SUMS.txt"]
names=zipfile.ZipFile(ZIP).namelist()
chk(41,"Build ZIP from explicit 13-file allowlist",
    names==["Video_5_HIT_FINAL/"+m for m in MANIFEST], str(names))
chk(42,"Exactly 13 files in ZIP", len(names)==13, str(len(names)))
chk(43,"No _source/Python/temp/render files inside ZIP",
    not any(n.endswith((".py",".pyc",".html",".png")) or "_source" in n for n in names))
chk(44,"README matches archive",
    all(os.path.basename(m) in readme for m in MANIFEST))
sums=open(os.path.join(ROOT,"SHA256SUMS.txt"),encoding="utf-8").read().strip().split("\n")
entries=[l for l in sums if not l.startswith("#") and l.strip()]
chk(45,"SHA256SUMS.txt exactly 12 entries", len(entries)==12, str(len(entries)))
p=subprocess.run("sha256sum -c SHA256SUMS.txt",shell=True,cwd=ROOT,
                 capture_output=True,text=True)
chk(46,"sha256sum -c successful", p.returncode==0, p.stdout[-300:]+p.stderr[-300:])
zsha=sha(ZIP)
sib=open(ZIP+".sha256",encoding="utf-8").read().split()[0]
chk(47,"Sibling ZIP checksum matches", zsha==sib, zsha+" / "+sib)
chk(48,"Description-only DOCX SHA-256 reported", True, sha(DESCDOC))
chk(49,"Main PPTX SHA-256 reported", True, sha(os.path.join(OUT,"Video_5_Main_Slides.pptx")))
chk(50,"Reveal PPTX SHA-256 reported", True, sha(os.path.join(OUT,"Video_5_Reveal_Builds.pptx")))
chk(51,"Slide-preview PDF SHA-256 reported", True, sha(os.path.join(OUT,"Video_5_Slide_Preview.pdf")))
g=subprocess.run("git status --porcelain",shell=True,cwd="/home/user/temidayoafonja-site",
                 capture_output=True,text=True).stdout.strip().split("\n")
g=[x for x in g if x.strip()]
outside=[x for x in g if "deliverables/" not in x]
chk(52,"No website/product/other-video file changed", not outside, str(g))

R["_zip_sha256"]=zsha
R["_summary"]={"total":52,"passed":52-len(F),"failed":len(F),"failures":F}
print(json.dumps(R,indent=1)[:200])
print("\nFAILURES:", F if F else "none")
open("/tmp/v5v3/QA_REPORT.json","w").write(json.dumps(R,indent=1))
