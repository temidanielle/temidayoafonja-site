"""
Video 1 slide deck - "How I Changed Jobs Without Starting My Career Over"

Layout is described once, as a list of absolutely-positioned elements on a
1920 x 1080 canvas, and rendered by two backends:

  * pptx_backend  -> a native, fully editable PowerPoint deck
  * html_backend  -> HTML that Chromium rasterises to exact 1920x1080 PNGs
                     and prints to PDF

Design rules (from the Video 1 un-script working sheet, section 08):
  - 16:9; consistent upper-left presenter safe area (~25% of frame width)
    kept free of essential content on every instructional slide
  - one idea per slide, large type, simple diagrams, progressive reveals
  - palette: medium blue #244B78, warm cream #F5F1E8, muted gold #C9A84C,
    deep navy #0F2346 for type and occasional bands, rust #C1440E as a
    small accent only
"""
import os
from PIL import Image
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(ROOT, "assets")
REPO = os.path.dirname(os.path.dirname(ROOT))
FONT_DIR = os.path.join(REPO, "fonts")

# ---------------------------------------------------------------- geometry
W, H = 1920, 1080
PX = 6350                       # EMU per design pixel (12192000 / 1920)
def E(px): return Emu(int(round(px * PX)))
def P(px): return Pt(px / 2.0)  # 1920px wide slide == 13.333in == 144 px/inch

# Presenter picture-in-picture safe area (upper left, 25% of frame width).
SAFE = dict(x=72, y=72, w=480, h=270)      # where the camera box sits
SAFE_CLEAR = dict(x=0, y=0, w=620, h=440)  # region kept free of essential content

TOP_X, TOP_Y, TOP_W = 660, 132, 1160       # headline zone, right of the safe area
LOW_X, LOW_Y, LOW_W = 120, 500, 1680       # body zone, below the safe area
ENDCARD_CLEAR = dict(x=1130, y=190, w=730, h=700)  # YouTube end-screen element

# ------------------------------------------------------------------ colour
NAVY  = RGBColor(0x0F, 0x23, 0x46)
BLUE  = RGBColor(0x24, 0x4B, 0x78)
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
RUST  = RGBColor(0xC1, 0x44, 0x0E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_DIM = RGBColor(0xB9, 0xC3, 0xD2)
NAVY_DIM  = RGBColor(0x5A, 0x6B, 0x82)
RULE_CREAM = RGBColor(0xE2, 0xDC, 0xCC)

DISPLAY = "Montserrat"
BODY    = "DM Sans"
SERIF   = "Cormorant Garamond"

def hexc(c): return "#%s" % str(c)


# ------------------------------------------------------------------ canvas
class Canvas(object):
    """One slide: an ordered list of drawing elements plus speaker notes."""

    def __init__(self, number, step=1):
        self.number, self.step = number, step
        self.els = []
        self.notes = ""

    def add(self, el):
        self.els.append(el)
        return el


DEFAULT_PARA = dict(size=32, font=DISPLAY, color=NAVY, bold=False, italic=False,
                    align="l", spacing=1.12, tracking=0, space_before=0,
                    space_after=0)


# --------------------------------------------------------------- primitives
def rect(sl, x, y, w, h, fill=None, line=None, lw=2, shape="rect", rot=0):
    return sl.add(dict(t="rect", x=x, y=y, w=w, h=h, fill=fill, line=line,
                       lw=lw, shape=shape, dash=None, rot=rot))


def block(sl, x, y, w, lines, anchor="t", h=None):
    """lines = list of (text, style-overrides)."""
    paras = []
    for txt, st in lines:
        p = dict(DEFAULT_PARA)
        p.update(st)
        p["text"] = txt
        paras.append(p)
    return sl.add(dict(t="text", x=x, y=y, w=w, h=h or 0, anchor=anchor,
                       paras=paras))


def eyebrow(sl, x, y, text, color=GOLD, size=26, w=1200, align="l"):
    return block(sl, x, y, w, [(text.upper(),
                 dict(size=size, font=DISPLAY, color=color, bold=True,
                      tracking=4.5, align=align, spacing=1.0))])


def hairline(sl, x, y, w, color=GOLD, h=3):
    return rect(sl, x, y, w, h, fill=color)


def img(sl, path, x, y, w, h, mode="cover", rotation=0):
    return sl.add(dict(t="img", path=path, x=x, y=y, w=w, h=h, mode=mode,
                       rotation=rotation))


def img_cover(sl, path, x, y, w, h, rotation=0):
    return img(sl, path, x, y, w, h, "cover", rotation)


def img_contain(sl, path, x, y, w, h, rotation=0):
    return img(sl, path, x, y, w, h, "contain", rotation)


def bg(sl, color):
    return rect(sl, 0, 0, W, H, fill=color)


def logomark(sl, x, y, unit=40, gap=11, gold=GOLD, rust=RUST, lw=3):
    """The Field Kit mark: four squares, the top-right one filled rust."""
    for (cx, cy), filled in (((0, 0), False), ((1, 0), True),
                             ((0, 1), False), ((1, 1), False)):
        rect(sl, x + cx * (unit + gap), y + cy * (unit + gap), unit, unit,
             fill=rust if filled else None,
             line=None if filled else gold, lw=lw)


def missing(sl, x, y, w, h, label):
    """A clearly-labelled placeholder for an asset that was not supplied."""
    s = rect(sl, x, y, w, h, fill=None, line=RUST, lw=3)
    s["dash"] = "dash"
    block(sl, x + 24, y, w - 48, [
        ("PLACEHOLDER", dict(size=22, color=RUST, bold=True, tracking=4, align="c")),
        (label, dict(size=20, font=BODY, color=RUST, align="c", spacing=1.25,
                     space_before=6)),
    ], anchor="m", h=h)
    return s


# ------------------------------------------------------------ pptx backend
_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
_ANCHOR = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


_PPTX_SHAPE = {"oval": MSO_SHAPE.OVAL,
               "tri_down": MSO_SHAPE.ISOSCELES_TRIANGLE,
               "tri_right": MSO_SHAPE.ISOSCELES_TRIANGLE}
_TRI_ROT = {"tri_down": 180, "tri_right": 90}


def _pptx_rect(sl, el):
    shape = _PPTX_SHAPE.get(el["shape"], MSO_SHAPE.RECTANGLE)
    s = sl.shapes.add_shape(shape, E(el["x"]), E(el["y"]), E(el["w"]), E(el["h"]))
    if el["shape"] in _TRI_ROT:
        s.rotation = _TRI_ROT[el["shape"]]
    if el.get("rot"):
        s.rotation = el["rot"]
    s.shadow.inherit = False
    if el["fill"] is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = el["fill"]
    if el["line"] is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = el["line"]
        s.line.width = E(el["lw"])
        if el.get("dash"):
            ln = s.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": el["dash"]}))
    return s


def _text_height(el):
    """Height a text box actually needs, so no box runs past the canvas."""
    total = 0
    for sp in el["paras"]:
        lines = sp["text"].count("\n") + 1
        total += lines * sp["size"] * sp["spacing"]
        total += sp["space_before"] + sp["space_after"]
    return max(int(round(total)) + 8, 24)


def _pptx_text(sl, el):
    h = el["h"] or _text_height(el)
    h = min(h, max(24, H - el["y"]))
    tb = sl.shapes.add_textbox(E(el["x"]), E(el["y"]), E(el["w"]), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = _ANCHOR[el["anchor"]]
    for i, sp in enumerate(el["paras"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = _ALIGN[sp["align"]]
        p.line_spacing = sp["spacing"]
        if sp["space_before"]:
            p.space_before = P(sp["space_before"])
        if sp["space_after"]:
            p.space_after = P(sp["space_after"])
        parts = sp["text"].split("\n")
        for j, line in enumerate(parts):
            r = p.add_run()
            r.text = line
            f = r.font
            f.name, f.size = sp["font"], P(sp["size"])
            f.bold, f.italic = sp["bold"], sp["italic"]
            f.color.rgb = sp["color"]
            rPr = r.font._rPr
            if sp["tracking"]:
                rPr.set("spc", str(int(sp["tracking"] * 50)))
            for tag in ("a:ea", "a:cs"):
                rPr.append(rPr.makeelement(qn(tag), {"typeface": sp["font"]}))
            if j < len(parts) - 1:
                p._p.append(p._p.makeelement(qn("a:br"), {}))
    return tb


def _pptx_img(sl, el):
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    iw, ih = Image.open(el["path"]).size
    if el["mode"] == "contain":
        s = min(w / iw, h / ih)
        dw, dh = iw * s, ih * s
        pic = sl.shapes.add_picture(el["path"], E(x + (w - dw) / 2),
                                    E(y + (h - dh) / 2), E(dw), E(dh))
    else:
        pic = sl.shapes.add_picture(el["path"], E(x), E(y), E(w), E(h))
        box_ar, img_ar = w / h, iw / ih
        if img_ar > box_ar:
            keep = box_ar / img_ar
            pic.crop_left = pic.crop_right = (1 - keep) / 2
        elif img_ar < box_ar:
            keep = img_ar / box_ar
            pic.crop_top = pic.crop_bottom = (1 - keep) / 2
    if el["rotation"]:
        pic.rotation = el["rotation"]
    return pic


def _is_backdrop(el):
    return (el["t"] == "rect" and el["fill"] is not None and el["line"] is None
            and el["x"] <= 0 and el["y"] <= 0 and el["w"] >= W and el["h"] >= H)


def render_pptx(canvases, path, titles=None):
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width, prs.slide_height = E(W), E(H)
    for cv in canvases:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        for el in cv.els:
            if _is_backdrop(el):
                fill = sl.background.fill
                fill.solid()
                fill.fore_color.rgb = el["fill"]
                continue
            {"rect": _pptx_rect, "text": _pptx_text, "img": _pptx_img}[el["t"]](sl, el)
        if cv.notes:
            sl.notes_slide.notes_text_frame.text = cv.notes
    prs.save(path)
    return path


# ------------------------------------------------------------ html backend
_FONT_FACES = [
    ("Montserrat", 400, "normal", "Montserrat-400-normal-latin-49e242.woff2"),
    ("Montserrat", 500, "normal", "Montserrat-500-normal-latin-49e242.woff2"),
    ("Montserrat", 600, "normal", "Montserrat-600-normal-latin-49e242.woff2"),
    ("Montserrat", 700, "normal", "Montserrat-700-normal-latin-49e242.woff2"),
    ("DM Sans", 400, "normal", "DMSans-400-normal-latin-1c49a6.woff2"),
    ("DM Sans", 500, "normal", "DMSans-500-normal-latin-1c49a6.woff2"),
    ("DM Sans", 700, "normal", "DMSans-600-normal-latin-1c49a6.woff2"),
    ("Cormorant Garamond", 400, "normal", "CormorantGaramond-400-normal-latin-abcaa8.woff2"),
    ("Cormorant Garamond", 700, "normal", "CormorantGaramond-600-normal-latin-abcaa8.woff2"),
    ("Cormorant Garamond", 400, "italic", "CormorantGaramond-400-italic-latin-4db21d.woff2"),
]

_CSS = """
* { margin:0; padding:0; box-sizing:border-box; }
html,body { background:#ffffff; }
.slide { position:relative; width:1920px; height:1080px; overflow:hidden;
         background:#F5F1E8; page-break-after:always; break-after:page; }
.slide:last-child { page-break-after:auto; break-after:auto; }
.el { position:absolute; }
.tx { display:flex; flex-direction:column; }
.tx.t { justify-content:flex-start; }
.tx.m { justify-content:center; }
.tx.b { justify-content:flex-end; }
p { white-space:pre-wrap; }
"""


def _esc(s):
    return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
             .replace("\n", "<br>"))


def _html_el(el, rel):
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    if el["t"] == "rect":
        st = ["left:%gpx" % x, "top:%gpx" % y, "width:%gpx" % w, "height:%gpx" % h]
        if el["fill"] is not None:
            st.append("background:%s" % hexc(el["fill"]))
        if el["line"] is not None:
            st.append("border:%gpx %s %s" % (el["lw"], el.get("dash") and "dashed"
                                             or "solid", hexc(el["line"])))
        if el["shape"] == "oval":
            st.append("border-radius:50%")
        elif el["shape"] == "tri_down":
            st.append("clip-path:polygon(0% 0%,100% 0%,50% 100%)")
        elif el["shape"] == "tri_right":
            st.append("clip-path:polygon(0% 0%,100% 50%,0% 100%)")
        if el.get("rot"):
            st.append("transform:rotate(%gdeg)" % el["rot"])
        return '<div class="el" style="%s"></div>' % ";".join(st)

    if el["t"] == "img":
        st = ["left:%gpx" % x, "top:%gpx" % y, "width:%gpx" % w, "height:%gpx" % h,
              "background-image:url('%s')" % rel(el["path"]),
              "background-size:%s" % el["mode"],
              "background-position:center", "background-repeat:no-repeat"]
        if el["rotation"]:
            st.append("transform:rotate(%gdeg)" % el["rotation"])
        return '<div class="el" style="%s"></div>' % ";".join(st)

    st = ["left:%gpx" % x, "top:%gpx" % y, "width:%gpx" % w]
    if h:
        st.append("height:%gpx" % h)
    out = ['<div class="el tx %s" style="%s">' % (el["anchor"], ";".join(st))]
    for sp in el["paras"]:
        ps = ["font-family:'%s',sans-serif" % sp["font"],
              "font-size:%gpx" % sp["size"],
              "font-weight:%d" % (700 if sp["bold"] else 400),
              "line-height:%g" % sp["spacing"],
              "color:%s" % hexc(sp["color"]),
              "text-align:%s" % {"l": "left", "c": "center", "r": "right"}[sp["align"]]]
        if sp["italic"]:
            ps.append("font-style:italic")
        if sp["tracking"]:
            ps.append("letter-spacing:%gpx" % sp["tracking"])
        if sp["space_before"]:
            ps.append("margin-top:%gpx" % sp["space_before"])
        if sp["space_after"]:
            ps.append("margin-bottom:%gpx" % sp["space_after"])
        out.append('<p style="%s">%s</p>' % (";".join(ps), _esc(sp["text"])))
    out.append("</div>")
    return "".join(out)


def render_html(canvases, path, title="Deck"):
    outdir = os.path.dirname(os.path.abspath(path))
    rel = lambda p: os.path.relpath(p, outdir).replace(os.sep, "/")
    faces = "\n".join(
        "@font-face{font-family:'%s';font-weight:%d;font-style:%s;"
        "src:url('%s') format('woff2');font-display:block;}"
        % (fam, wt, sty, rel(os.path.join(FONT_DIR, fn)))
        for fam, wt, sty, fn in _FONT_FACES)
    body = "\n".join(
        '<div class="slide" id="s%d">%s</div>'
        % (i + 1, "".join(_html_el(el, rel) for el in cv.els))
        for i, cv in enumerate(canvases))
    html = ("<!doctype html><html><head><meta charset='utf-8'><title>%s</title>"
            "<style>%s\n%s</style></head><body>%s</body></html>"
            % (title, faces, _CSS, body))
    with open(path, "w") as f:
        f.write(html)
    return path
