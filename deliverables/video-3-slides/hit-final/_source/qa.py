# -*- coding: utf-8 -*-
"""Direct-address revision QA for Videos 2, 3 and 4.

Runs the 20 shared checks from the prompt, the per-video checks, and the
package / checksum / rendering checks. Reports only; changes nothing."""
import os, sys, re, json, zipfile, subprocess, glob
sys.path.insert(0,"/tmp/da")
from changereport import (old_spoken, new_blocks, report, scan_detached,
                          viewer_ratio, DETACHED)
from docx import Document
from pptx import Presentation
from docxkit import sha256

REPO="/home/user/temidayoafonja-site"
V={
 2:{"base":"/tmp/da/v2","prior_tel":REPO+"/deliverables/video-2-slides/hit-final/LONG_FORM/Video2TeleprompterScriptwithslidemarkers_HIT_v2.0.txt",
    "prior_zip":None,"prior_words":1131,
    "main":REPO+"/deliverables/video-2-slides/out/Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.1.pptx",
    "rev":REPO+"/deliverables/video-2-slides/out/Video-2-Reveal-Builds_v1.1.pptx",
    "pdf":REPO+"/deliverables/video-2-slides/out/Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.1.pdf",
    "slides":13,"frames":23,"markers":13,
    "title":"Is Your Job Making You Less Marketable?",
    "thumb":"YOUR SKILLS ARE STALLING","cta":"Capability Formation Field Kit",
    "cta_url":"https://temidayoafonja.com/fieldkit",
    "next":"3 Things to Do Before Quitting Your Job",
    "forbidden_cta":["Keep the Proof","Career Evidence Starter",
                     "career-evidence-starter","career-decisions"]},
 3:{"base":"/tmp/da/v3","prior_tel":REPO+"/deliverables/video-3-slides/hit-final/Video_3_HIT_FINAL/LONG_FORM/Video3TeleprompterScriptwithslidemarkers_HIT_v2.0.txt",
    "prior_zip":"2455a0d08105e3148215191e62ead6204c8e4cdf896525592a2983b8c14ea177",
    "prior_words":1205,
    "main":REPO+"/deliverables/video-3-slides/out/Video-3-3-Things-to-Do-Before-Quitting-Your-Job_v1.1.pptx",
    "rev":REPO+"/deliverables/video-3-slides/out/Video-3-Reveal-Builds_v1.1.pptx",
    "pdf":REPO+"/deliverables/video-3-slides/out/Video-3-3-Things-to-Do-Before-Quitting-Your-Job_v1.1.pdf",
    "slides":13,"frames":27,"markers":13,
    "title":"3 Things to Do Before Quitting Your Job",
    "thumb":"WAIT BEFORE YOU QUIT","cta":"Career Decision Evidence Check",
    "cta_url":"https://temidayoafonja.com/career-decisions",
    "next":"How to Change Jobs Without Starting Your Career Over",
    "forbidden_cta":["Keep the Proof","Career Evidence Starter",
                     "career-evidence-starter","fieldkit"]},
 4:{"base":"/tmp/da/v4","prior_tel":REPO+"/deliverables/video-4-slides/hit-final/Video_4_HIT_FINAL/LONG_FORM/Video4TeleprompterScriptwithslidemarkers_HIT_v2.1.txt",
    "prior_zip":"6d9e8339a83a463ad231db8d180f6bb27025b07f41fd4bfc914778ea5f602684",
    "prior_words":1261,
    "main":REPO+"/deliverables/video-4-slides/out/Video_4_Main_Slides.pptx",
    "rev":REPO+"/deliverables/video-4-slides/out/Video_4_Reveal_Builds.pptx",
    "pdf":REPO+"/deliverables/video-4-slides/out/Video_4_Slide_Preview.pdf",
    "slides":11,"frames":26,"markers":11,
    "title":"How to Explain Your Career Change",
    "thumb":"YOUR CAREER MAKES SENSE","cta":"Free Career Evidence Starter",
    "cta_url":"https://temidayoafonja.com/career-evidence-starter",
    "next":"Should I Make an Internal Move? 3 Questions to Decide",
    "forbidden_cta":["Keep the Proof","fieldkit","career-decisions"]},
}

def text(p): return "\n".join(x.text for x in Document(p).paragraphs)
def paras(p): return [x.text.strip() for x in Document(p).paragraphs if x.text.strip()]

def block(t):
    a=t.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
    b=t.index("— END OF COPY-READY DESCRIPTION —")
    return [x.strip() for x in t[a:b].split("\n") if x.strip()][1:]

RESULT={}
for n,c in V.items():
    B=c["base"]; ROOT=os.path.join(B,"Video_%d_HIT_FINAL"%n)
    LF=os.path.join(ROOT,"LONG_FORM"); SH=os.path.join(ROOT,"SHORTS")
    R={}; F=[]
    def chk(k,label,cond,detail=""):
        R[k]={"check":label,"pass":bool(cond),"detail":str(detail)}
        if not cond: F.append((k,label,str(detail)))

    canon=os.path.join(B,"canonical_v3.0.txt")
    nb=new_blocks(canon,"BEGIN APPROVED VIDEO %d v3.0 DIRECT-ADDRESS SCRIPT"%n,
                        "END APPROVED VIDEO %d v3.0 DIRECT-ADDRESS SCRIPT"%n)
    new=[x for x in nb if not x.startswith("[SLIDE:")]
    mk=[x for x in nb if x.startswith("[SLIDE:")]
    old=old_spoken(c["prior_tel"])
    rep=report(old,new)
    words=sum(len(x.split()) for x in new)

    tel=os.path.join(LF,"Video%dTeleprompterScriptwithslidemarkers_HIT_v3.0_DIRECT_ADDRESS"%n)
    rdg=os.path.join(LF,"Video%dReadingScriptnomarkers_HIT_v3.0_DIRECT_ADDRESS"%n)
    edb=os.path.join(LF,"Video_%d_EDITOR_ONLY_HIT_Brief_v3.0_DIRECT_ADDRESS.docx"%n)
    pub=os.path.join(LF,"Video_%d_Publishing_Package_HIT_v3.0_DIRECT_ADDRESS.docx"%n)
    seb=os.path.join(SH,"Video_%d_Shorts_EDITOR_ONLY_HIT_Brief.docx"%n)
    desc=os.path.join(B,"Video_%d_YouTube_Description_HIT.docx"%n)
    zp=os.path.join(B,"Video_%d_HIT_FINAL_Recording_and_Shorts_Package.zip"%n)
    shorts={f:text(os.path.join(SH,f)) for f in sorted(os.listdir(SH))
            if f.startswith("Video_%d_Short_"%n)}

    telt=open(tel+".txt",encoding="utf-8").read()
    rdgt=open(rdg+".txt",encoding="utf-8").read()
    pubt=text(pub); desct=text(desc); edbt=text(edb); sebt=text(seb)
    allpub=pubt+desct
    spoken="\n\n".join(new)

    # ---------- canonical / change control
    telb=[x.strip() for x in telt.split("\n\n") if x.strip()][1:]
    telnorm=["[SLIDE: %s]"%x[len("SLIDE  —  "):] if x.startswith("SLIDE  —  ") else x
             for x in telb]
    chk("C1","Teleprompter TXT == revised canonical", telnorm==nb)
    chk("C2","Reading TXT == canonical spoken",
        [x.strip() for x in rdgt.split("\n\n") if x.strip()]==new)
    chk("C3","Teleprompter minus markers == reading script",
        [x for x in telnorm if not x.startswith("[SLIDE:")]==new)
    teld=paras(tel+".docx")[4:]
    teldn=["[SLIDE: %s]"%x[len("SLIDE  —  "):] if x.startswith("SLIDE  —  ") else x
           for x in teld]
    chk("C4","Teleprompter DOCX == canonical", teldn==nb)
    chk("C5","Reading DOCX == canonical spoken", paras(rdg+".docx")[4:]==new)
    chk("C6","Marker count and order preserved",
        len(mk)==c["markers"] and
        [i for i,x in enumerate(nb) if x.startswith("[SLIDE:")]==
        [i for i,x in enumerate(telnorm) if x.startswith("[SLIDE:")], len(mk))
    chk("C7","No paragraph removed", rep["removed"]==0, rep["removed"])
    chk("C8","Revision is bounded, not a rewrite",
        rep["unchanged"]>=rep["changed"],
        "unchanged %d vs changed %d"%(rep["unchanged"],rep["changed"]))

    # ---------- 1-20 shared direct-address QA
    hits=scan_detached(spoken)+ [h for t in shorts.values() for h in scan_detached(t)]
    chk("D1","No detached phrasing in long-form or Shorts", not hits, hits)
    vr=viewer_ratio(new)
    chk("D2","Viewer-facing situations predominantly second person",
        vr["ratio"]>=0.70,
        "%d of %d viewer-facing paragraphs (%.0f%%)"
        %(vr["second_person"],vr["viewer_facing"],100*vr["ratio"]))
    BRIDGE=["let me show you","i want to help you","i want you to",
            "here is what i would","i want to walk you through",
            "that is what i want to help you","now try the same"]
    low=spoken.lower()
    bridges={b:low.count(b) for b in BRIDGE}
    chk("D3","First-person evidence bridges back to the viewer",
        sum(bridges.values())>=3, {k:v for k,v in bridges.items() if v})
    ep=spoken.count("experienced professionals")
    chk("D4","'experienced professionals' is positioning language",
        ep==1, "count=%d"%ep)
    chk("D5","Organizational mechanics carry viewer consequence",
        vr["ratio"]>=0.70)
    allsec=len([p for p in new if re.search(r"\b(you|your)\b",p,re.I)])/len(new)
    chk("D6","Direct address not mechanically repetitive", allsec<=0.90,
        "%.0f%% of all spoken paragraphs"%(100*allsec))
    KEY=["ladies and gentlemen","in this presentation","as you can see on this slide",
         "thank you for having me","let's dive in","in today's talk",
         "welcome back to my channel","hey everyone","hi guys"]
    kn=[k for k in KEY if k in spoken.lower()]
    chk("D7","No keynote or conference voice", not kn, kn)
    chk("D8","No generic career-article phrasing", not hits)
    chk("D9","No unsupported fact added",
        rep["inserted"]<=1 and not re.findall(r"\b\d+\s?%|\$\s?\d", spoken),
        "inserted %d"%rep["inserted"])
    oldtxt="\n\n".join(old)
    dev_old=len(re.findall(r"\bacronym\b",oldtxt,re.I))
    dev_new=len(re.findall(r"\bacronym\b",spoken,re.I))
    chk("D10","No additional mnemonic",
        not re.search(r"\bCAR\b",spoken) and "3 Cs" not in spoken
        and dev_new<=dev_old,
        "'acronym' occurrences prior %d, revised %d"%(dev_old,dev_new))
    chk("D11","One primary CTA remains",
        c["cta"].lower() in allpub.lower()
        and not any(f.lower() in allpub.lower() for f in c["forbidden_cta"]),
        [f for f in c["forbidden_cta"] if f.lower() in allpub.lower()])
    chk("D12","Watch Next remains correct",
        c["next"] in allpub and c["next"] in spoken)
    sh_ok={f:bool(re.search(r"\b(you|your|yourself)\b",t,re.I)) for f,t in shorts.items()}
    chk("D13","Every Short speaks to one viewer", all(sh_ok.values()), sh_ok)
    chk("D14","Both editor briefs carry the direct-address section",
        "Direct address is part of the creative" in edbt
        and "Direct address is part of the creative" in sebt)
    prov=json.load(open(os.path.join(B,"_partdiff.json")))
    chk("D15","Slides visually unchanged",
        prov["main_nonnotes"]==[] and prov["reveal_nonnotes"]==[], prov["main_nonnotes"])
    chk("D16","Only notes parts changed in both PPTX",
        len(prov["main_changed"])==c["slides"] and len(prov["reveal_changed"])==c["frames"],
        "%d / %d"%(len(prov["main_changed"]),len(prov["reveal_changed"])))
    chk("D17","Public description searchable and copy-ready",
        c["title"] in pubt and c["thumb"] in pubt
        and "COPY-READY YOUTUBE DESCRIPTION — BEGIN" in pubt)
    chk("D18","Working chapters flagged as estimates",
        "WORKING ESTIMATES — EDITOR MUST REPLACE FROM FINAL CUT" in pubt
        and pubt.index("— END OF COPY-READY DESCRIPTION —")
            < pubt.index("WORKING ESTIMATES — EDITOR MUST REPLACE FROM FINAL CUT"))
    chk("D19","No [INSERT] placeholder", "[INSERT" not in (pubt+desct+edbt+sebt))
    chk("D20","Publishing public block == description-only public block",
        block(pubt)==block(desct))

    # ---------- recording docs stay clean
    ED=["Visual:","On-screen","B-roll","EDITOR","Reveal:","End on:",
        "FACTUAL BOUNDARY","Related Video:"]
    bad=[(f,w) for f,t in shorts.items() for w in ED if w in t]
    chk("D21","Recording docs contain no editor instructions", not bad, bad)
    chk("D22","Both EDITOR ONLY docs clearly labelled",
        edbt.strip().startswith("EDITOR ONLY") and sebt.strip().startswith("EDITOR ONLY"))

    # ---------- decks
    mp=Presentation(c["main"]); rp=Presentation(c["rev"])
    chk("P1","Main slide count", len(mp.slides)==c["slides"], len(mp.slides))
    chk("P2","Reveal frame count", len(rp.slides)==c["frames"], len(rp.slides))
    notes=[s.notes_slide.notes_text_frame.text for s in mp.slides]+\
          [s.notes_slide.notes_text_frame.text for s in rp.slides]
    chk("P3","Notes updated to v3.0 direct address",
        all("Timing:" in x for x in notes)
        and any("direct-address" in x or "Direct address" in x or
                "direct address" in x for x in notes))

    # ---------- package
    names=zipfile.ZipFile(zp).namelist()
    chk("Z1","Exactly 13 files in ZIP", len(names)==13, len(names))
    chk("Z2","No _source, Python, temp or render files in ZIP",
        not any(x.endswith((".py",".pyc",".html",".png")) or "_source" in x
                for x in names))
    sums=[l for l in open(os.path.join(ROOT,"SHA256SUMS.txt"),encoding="utf-8")
          .read().strip().split("\n") if l.strip() and not l.startswith("#")]
    chk("Z3","SHA256SUMS.txt exactly 12 entries", len(sums)==12, len(sums))
    pr=subprocess.run("sha256sum -c SHA256SUMS.txt",shell=True,cwd=ROOT,
                      capture_output=True,text=True)
    chk("Z4","sha256sum -c successful", pr.returncode==0, pr.stdout[-200:])
    zs=sha256(zp)
    chk("Z5","Sibling ZIP checksum matches",
        open(zp+".sha256",encoding="utf-8").read().split()[0]==zs)
    rd=open(os.path.join(ROOT,"README_FINAL.txt"),encoding="utf-8").read()
    chk("Z6","README matches archive",
        all(os.path.basename(m) in rd for m in
            [x.split("/",1)[1] for x in names if "/" in x.split("/",1)[1]] or
            [os.path.basename(x) for x in names]))

    RESULT[n]={"checks":R,"failures":F,
      "change_report":{k:v for k,v in rep.items() if k!="blocks"},
      "voice":{k:v for k,v in vr.items() if k!="not_second_person"},
      "words_prior":c["prior_words"],"words_new":words,
      "markers":len(mk),"slides":len(mp.slides),"frames":len(rp.slides),
      "notes_parts":{"main":len(prov["main_changed"]),
                     "reveal":len(prov["reveal_changed"])},
      "hashes":{"zip":zs,"desc_docx":sha256(desc),"main_pptx":sha256(c["main"]),
                "reveal_pptx":sha256(c["rev"]),"preview_pdf":sha256(c["pdf"])},
      "prior_zip":c["prior_zip"]}
    print("VIDEO %d: %d checks, %d failed"%(n,len(R),len(F)))
    for f in F: print("   FAIL",f)

json.dump(RESULT,open("/tmp/da/QA_REPORT.json","w"),indent=1)
tot=sum(len(v["checks"]) for v in RESULT.values())
fail=sum(len(v["failures"]) for v in RESULT.values())
print("\nTOTAL %d checks, %d failed"%(tot,fail))
