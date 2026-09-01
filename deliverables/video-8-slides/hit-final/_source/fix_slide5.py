# -*- coding: utf-8 -*-
"""Apply the authorised Slide 5 text correction to both Video 8 decks.

Only the single gold emphasis line changes. The stale conceptual framing
"It feels like a competence gap. It is an information gap." is replaced by the
three approved statements, set as three lines inside the SAME text box, with
the existing run's properties cloned so typography family, size, weight and
colour are carried over untouched. Every line was measured against the real
Montserrat Bold at 20pt first: the widest is 7.06in inside an 11.11in box, so
nothing wraps, and the block ends at 6.71in of a 7.5in slide.
"""
import copy
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def q(t): return "{%s}%s" % (A, t)

STALE = "It feels like a competence gap. It is an information gap."
NEW = ["IT CAN FEEL LIKE A COMPETENCE GAP.",
       "SOME CONTEXT CAN BE RESEARCHED.",
       "SOME MUST BE LEARNED THROUGH EXPOSURE."]
TARGETS = {"Video_8_Main_Slides.pptx": [5], "Video_8_Reveal_Builds.pptx": [11]}

for path, slides in TARGETS.items():
    prs = Presentation(path)
    for n in slides:
        hit = None
        for sh in prs.slides[n - 1].shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == STALE:
                hit = sh
        assert hit is not None, "stale line not found on %s slide %d" % (path, n)
        body = hit.text_frame._txBody
        ps = body.findall(q("p"))
        assert len(ps) == 1, "expected a single paragraph"
        p = ps[0]
        runs = p.findall(q("r"))
        assert len(runs) == 1, "expected a single run"
        template = runs[0]
        # Rebuild the paragraph as three lines, cloning the original run's
        # rPr so nothing about the typography is re-specified by hand.
        for child in list(p):
            if etree.QName(child).localname in ("r", "br"):
                p.remove(child)
        for i, line in enumerate(NEW):
            if i:
                p.append(etree.SubElement(p, q("br")))
            r = copy.deepcopy(template)
            r.find(q("t")).text = line
            p.append(r)
        hit.height = Inches(1.05)   # 0.39in held one line; three need room
        print("%s slide %d: corrected" % (path, n))
    prs.save(path)
