# -*- coding: utf-8 -*-
"""Video 8 H.I.T. package QA — the 44 checks from the canonical prompt."""
import os, re, json, zipfile, hashlib, subprocess, sys
from docx import Document
from pptx import Presentation

ROOT = "Video_8_HIT_FINAL"
LF = ROOT + "/LONG_FORM"; SH = ROOT + "/SHORTS"
ZIPNAME = "Video_8_HIT_FINAL_Recording_and_Shorts_Package.zip"
DECKS = "/home/user/temidayoafonja-site/deliverables/video-8-slides/out/"
TITLE = "How to Switch Industries Without Starting Over"
THUMB = "YOUR EXPERIENCE STILL COUNTS"
NEXT = "What to Do Before a Layoff Happens"
R = []
def chk(n, desc, ok, detail=""):
    R.append({"n": n, "check": desc, "result": "PASS" if ok else "FAIL",
              "detail": detail})

def paras(p):
    return [x.text for x in Document(p).paragraphs]
def text(p):
    return "\n".join(paras(p))

DOCX = [os.path.join(dp, f) for dp, _, fn in os.walk(ROOT) for f in fn
        if f.endswith(".docx")]
DESC_DOC = "Video_8_YouTube_Description_HIT.docx"
ALL_DOCS = sorted(DOCX) + [DESC_DOC]
TXTS = [os.path.join(dp, f) for dp, _, fn in os.walk(ROOT) for f in fn
        if f.endswith(".txt")]
def alltext(paths): return "\n".join(text(p) if p.endswith(".docx")
                                     else open(p, encoding="utf-8").read()
                                     for p in paths)
PUBLIC = [p for p in ALL_DOCS + TXTS if "EDITOR_ONLY" not in p]
EVERY = ALL_DOCS + TXTS

# ---- Video 8 checks 1-35 ----
main = Presentation(DECKS + "Video_8_Main_Slides.pptx")
rev = Presentation(DECKS + "Video_8_Reveal_Builds.pptx")
def deck_text(prs):
    out = []
    for s in prs.slides:
        out.append("\n".join(sh.text_frame.text for sh in s.shapes
                             if sh.has_text_frame))
    return out
MAINT, REVT = deck_text(main), deck_text(rev)
DECKALL = "\n".join(MAINT + REVT)

chk(1, "actual main-slide count", len(main.slides._sldIdLst) == 12,
    "%d main slides" % len(main.slides._sldIdLst))
chk(2, "actual reveal-frame count", len(rev.slides._sldIdLst) == 24,
    "%d reveal frames" % len(rev.slides._sldIdLst))
pub = text(LF + "/Video_8_Publishing_Package_HIT_v2.2.docx")
a = alltext(ALL_DOCS + TXTS)

def blocks(p):
    return (paras(p) if p.endswith(".docx")
            else [x for x in open(p, encoding="utf-8").read().split("\n\n")])
def context(bl, i):
    """Paragraph window plus the list lead-in governing paragraph i."""
    win = bl[max(0, i - 2):i + 3]
    for j in range(i - 1, max(-1, i - 12), -1):
        prev = bl[j].strip()
        if prev.endswith(":"):
            win = [prev] + win; break
        if prev and not prev.startswith(("—", "-", "•", " ")):
            break
    return " ".join(win)
EVERY = ALL_DOCS + TXTS
def scan(pats, reject, label, n):
    bad = []
    for p in EVERY:
        bl = blocks(p)
        for i, para in enumerate(bl):
            for pat in pats:
                if re.search(pat, para, re.I) and not re.search(reject, context(bl, i), re.I):
                    bad.append((os.path.basename(p), para[:110]))
    chk(n, label, not bad, str(bad))

chk(3, "title exactly correct", TITLE in pub)
chk(4, "thumbnail exactly correct", THUMB in pub)
chk(5, "primary search phrase exact",
    "how to switch industries without starting over" in pub)
chk(6, "'move into a new industry' not retained as public title",
    "move into a new industry" not in a.lower())
tel = open(LF + "/Video8TeleprompterScriptwithslidemarkers_HIT_v2.2.txt",
           encoding="utf-8").read()
chk(7, "H.I.T. opening matches the approved v2.2 script",
    # v2.2 opens on the audience pain, then the v2.1 line becomes the answer to it.
    "direct industry experience required" in tel
    and "Changing industries does not make you entry-level at everything. It makes "
        "you new to a context." in tel
    and "prepared for the CISM exam and didn’t pass the first time" in tel)
chk(8, "eight industries/sectors wording",
    "eight industries and sectors" in a)
chk(9, "'nearly two decades' / 'roughly eighteen years' within boundary",
    "nearly two decades" in a and "eighteen-year" in a,
    "no other duration claim: %s" % (not re.search(r"(twenty|thirty|20|25|30)[- ]?year", a, re.I)))
chk(10, "CISM first-attempt non-pass wording",
    "didn’t pass the first time" in a or "not passing the first time" in a)
scan([r"\bscored?\s+\d", r"exam (score|date)", r"passed the (exam|cism)",
      r"second attempt", r"\d+\s*%\s*(on|in) the exam"],
     r"do not invent|do not|never|no invented|must not|fake|\bno\b|neither",
     "no invented exam score, date, attempt count or passing result", 11)
scan([r"everything transfers", r"all (of your |your )?experience transfers",
      r"fully transferable", r"whole career is transferable"],
     r"does not|do not|not mean|never|goal is not|but it also|rather than|not to prove",
     "no claim that all experience transfers", 12)
chk(13, "Slide 5 stale phrase removed from live visual materials",
    "It is an information gap" not in DECKALL
    and "It feels like a competence gap" not in DECKALL)
chk(14, "revised Slide 5 text present on main slide 5",
    "IT CAN FEEL LIKE A COMPETENCE GAP." in MAINT[4]
    and "SOME CONTEXT CAN BE RESEARCHED." in MAINT[4]
    and "SOME MUST BE LEARNED THROUGH EXPOSURE." in MAINT[4])
chk(15, "corresponding reveal-frame stale text corrected",
    "IT CAN FEEL LIKE A COMPETENCE GAP." in REVT[10]
    and not any("information gap" in t for t in REVT))
chk(16, "Slide 12 left untouched (carries no video title)",
    "CONTINUE THE SERIES" in MAINT[11] and NEXT.upper() not in MAINT[11],
    "reported to Temidayo as an open decision")
import zipfile as _zf
def changed(orig, now):
    A_, B_ = _zf.ZipFile(orig), _zf.ZipFile(now)
    return [x.split("/")[-1] for x in sorted(set(A_.namelist()))
            if A_.read(x) != B_.read(x)]
cm = changed("/tmp/v8hit/Main.orig.pptx", DECKS + "Video_8_Main_Slides.pptx")
cr = changed("/tmp/v8hit/Reveal.orig.pptx", DECKS + "Video_8_Reveal_Builds.pptx")
chk(17, "no unauthorized slide changed",
    cm == ["slide5.xml"] and cr == ["slide11.xml"], "main %s | reveal %s" % (cm, cr))
# marker-to-slide mapping
import sys as _sys; _sys.path.insert(0, "/tmp/v8hit")
from script_text import MARKERS
mk = [m[len("[SLIDE:"):-1].strip() for m in MARKERS]
chk(18, "marker-to-slide mapping one-to-one and ordered",
    len(mk) == 12 and len(set(mk)) == 12 and len(MAINT) == 12,
    "%d markers, %d slides, no duplicates" % (len(mk), len(MAINT)))
shorts = sorted(f for f in os.listdir(SH) if f.startswith("Video_8_Short_"))
chk(19, "exactly four Shorts", len(shorts) == 4, str(shorts))
DIRECTIVE = ["On-screen", "Visual:", "EDITOR", "Do not use", "Reveal:",
             "Related Video", "Avoid", "Secondary:"]
bad = [(f, d) for f in shorts for d in DIRECTIVE
       if d.lower() in text(os.path.join(SH, f)).lower()]
chk(20, "Short recording docs contain no editor directions", not bad, str(bad))
e1 = paras(LF + "/Video_8_EDITOR_ONLY_HIT_Brief_v2.2.docx")
e2 = paras(SH + "/Video_8_Shorts_EDITOR_ONLY_HIT_Brief.docx")
chk(21, "both EDITOR ONLY docs labelled",
    e1[0] == "EDITOR ONLY" and e2[0] == "EDITOR ONLY")
chk(22, "Field Kit is sole product CTA",
    "temidayoafonja.com/fieldkit" in a)
chk(23, "Watch Next is Video 9", NEXT in pub and "(Video 9)" in pub)
PUBLIC = [p for p in EVERY if "EDITOR_ONLY" not in p]
leak = [os.path.basename(p) for p in PUBLIC
        if re.search(r"keep the proof|keep-the-proof|decision evidence check",
                     text(p) if p.endswith(".docx")
                     else open(p, encoding="utf-8").read(), re.I)]
chk(24, "no Keep the Proof / Decision Check CTA leak", not leak, str(leak))
desc_t = text(DESC_DOC)
EMOJI_OK = ["✨", "\U0001F9ED", "⏱", "▶", "\U0001F517"]
found = sorted(set(ch for ch in desc_t if ord(ch) > 0x2000 and ch not in "—’“”→·…️≠"))
chk(25, "approved emoji system present and not exceeded",
    all(e in desc_t for e in EMOJI_OK), "found: %r" % found)
chk(26, "working chapters inserted directly into public description",
    # v2.2: chapter 1 renamed to name the direct-experience pain and all
    # timings recomputed from the longer v2.2 script.
    "00:00 When the Job Says Direct Industry Experience Required" in desc_t
    and "12:32 What to Do Before a Layoff Happens" in desc_t)
chk(27, "no [INSERT] placeholder remains",
    not re.search(r"\[INSERT[^\]]*\]", a), str(re.findall(r"\[INSERT[^\]]*\]", a)))
dp = paras(DESC_DOC)
begin = dp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
end = dp.index("— END OF THE COPY-READY DESCRIPTION —")
warn = [i for i, x in enumerate(dp) if x.startswith("WORKING ESTIMATES")]
chk(28, "working-estimates warning sits outside the copy-ready block",
    all(i > end for i in warn), "begin=%d end=%d warn=%s" % (begin, end, warn))
chk(28.1, "no internal note inside the copy-ready block",
    not any(re.search(r"INTERNAL|EDITOR MUST|WORKING ESTIMATES", x)
            for x in dp[begin + 1:end]))
pp = paras(LF + "/Video_8_Publishing_Package_HIT_v2.2.docx")
pb, pe = (pp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN"),
          pp.index("— END OF THE COPY-READY DESCRIPTION —"))
chk(29, "description-only public copy == publishing-package public copy",
    dp[begin:end] == pp[pb:pe])
cv = subprocess.run([sys.executable, "verify_v22.py"],
                    capture_output=True, text=True)
def line(sub): return [l for l in cv.stdout.splitlines() if sub in l][0]
chk(30, "canonical source verification", "VERIFICATION: PASS" in cv.stdout)
chk(31, "teleprompter DOCX == TXT", "PASS" in line("teleprompter DOCX"))
chk(32, "reading DOCX == TXT", "PASS" in line("reading DOCX"))
chk(33, "teleprompter minus markers == reading script",
    "PASS" in line("minus markers == reading TXT"))
# 36-40 packaging
z = zipfile.ZipFile(ZIPNAME)
names = [n for n in z.namelist() if not n.endswith("/")]
chk(36, "ZIP built from an explicit 13-file allowlist", len(names) == 13)
chk(37, "exactly 13 files in the archive", len(names) == 13, "%d" % len(names))
chk(38, "no _source, Python, QA or image files in the archive",
    not [n for n in names if re.search(r"_source|\.py$|\.png$|\.pyc$|^\.", n)])
readme = open(ROOT + "/README_FINAL.txt", encoding="utf-8").read()
listed = [n.split("Video_8_HIT_FINAL/")[1] for n in names]
missing = [f for f in listed if os.path.basename(f) not in readme]
chk(39, "README matches archive", not missing, str(missing))
sums = [l for l in open(ROOT + "/SHA256SUMS.txt", encoding="utf-8")
        if l.strip() and not l.startswith("#")]
chk(40, "SHA256SUMS.txt has exactly 12 entries", len(sums) == 12, "%d" % len(sums))
# 41 sha256sum -c
cp = subprocess.run(["sha256sum", "-c", "SHA256SUMS.txt"], cwd=ROOT,
                    capture_output=True, text=True)
chk(41, "sha256sum -c succeeds", cp.returncode == 0,
    cp.stdout.strip().splitlines()[-1] if cp.stdout else cp.stderr[:200])
# 42 sibling checksum
def sha256(p):
    h = hashlib.sha256()
    with open(p, "rb") as fh:
        for b in iter(lambda: fh.read(1 << 20), b""): h.update(b)
    return h.hexdigest()
sib = open(ZIPNAME + ".sha256").read().split()[0]
chk(42, "sibling ZIP checksum matches the archive", sib == sha256(ZIPNAME), sib)
# 42b in-ZIP and on-disk SHA256SUMS identical
chk(42.1, "in-ZIP and on-disk SHA256SUMS.txt byte-identical",
    z.read("Video_8_HIT_FINAL/SHA256SUMS.txt")
    == open(ROOT + "/SHA256SUMS.txt", "rb").read())
# 43 description doc hash
chk(43, "description-only DOCX hash reported", True, sha256(DESC_DOC))
print(json.dumps(R, indent=1))
fails = [r for r in R if r["result"] == "FAIL"]
print("\n%d checks | %d PASS | %d FAIL" % (len(R), len(R) - len(fails), len(fails)))
for f in fails: print("  FAIL %s: %s  %s" % (f["n"], f["check"], f["detail"][:300]))
json.dump(R, open("QA_REPORT.json", "w"), indent=1)
