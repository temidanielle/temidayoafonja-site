"""Build Video 8 and run its QA.

  PPTX  Video_8_Main_Slides.pptx      - 12 editable slides, final revealed state
  PPTX  Video_8_Reveal_Builds.pptx    - every reveal state as a duplicate slide
  PDF   Video_8_Slide_Preview.pdf     - printed from Chromium at true 16:9
  PNG   exact 1920 x 1080 frames, contact sheet, reveal order, phone check

    python3 build/build.py
"""
import os, sys, json, shutil

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck import W, H, ENDCARD_CLEAR, Canvas, render_pptx, render_html
import slides as S
from PIL import Image, ImageDraw, ImageFont

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "out")
PNG = os.path.join(OUT, "png")
REV = os.path.join(OUT, "reveals")
RENDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_render")
CHROME = "/opt/pw-browsers/chromium-1194/chrome-linux/chrome"
FD = os.path.expanduser("~/.fonts/")

MAIN = "Video_8_Main_Slides"
RECDECK = "Video_8_Reveal_Builds"
PREVIEW = "Video_8_Slide_Preview"
N = 12
PLAN = [(n, s) for n in range(1, N + 1) for s in range(1, S.STEPS[n] + 1)]


def canvases(plan):
    out = []
    for n, step in plan:
        cv = Canvas(n, step)
        S.BUILDERS[n](cv, step)
        out.append(cv)
    return out


def shoot(html_path, png_dir, names, pdf_path=None):
    from playwright.sync_api import sync_playwright
    os.makedirs(png_dir, exist_ok=True)
    made = []
    with sync_playwright() as pw:
        b = pw.chromium.launch(executable_path=CHROME, args=["--no-sandbox"])
        pg = b.new_page(viewport={"width": W, "height": H}, device_scale_factor=1)
        pg.goto("file://" + os.path.abspath(html_path), wait_until="load")
        pg.wait_for_timeout(600)
        for i, name in enumerate(names, start=1):
            p = os.path.join(png_dir, name)
            pg.query_selector("#s%d" % i).screenshot(path=p)
            made.append(p)
        if pdf_path:
            pg.pdf(path=pdf_path, width="13.333in", height="7.5in",
                   print_background=True, scale=0.66665,
                   margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
        b.close()
    return made


def f(sz, w="Bold"):
    return ImageFont.truetype(FD + "MontserratTB-%s.ttf" % w, sz)


def sheet(pngs, path, tile_w, title, sub, labels=None, cols=4):
    tw = tile_w; th = round(tw * 9 / 16)
    rows = (len(pngs) + cols - 1) // cols
    lab_h = 46 if labels else 16
    pad, top = 46, 190
    Wd = pad * 2 + cols * tw + (cols - 1) * pad
    Hd = top + rows * (th + lab_h + pad) + pad
    im = Image.new("RGB", (Wd, Hd), (250, 247, 241)); d = ImageDraw.Draw(im)
    d.text((pad, 52), title, font=f(46), fill=(15, 35, 70))
    d.text((pad, 116), sub, font=f(26, "Regular"), fill=(120, 132, 150))
    d.rectangle([pad, 166, Wd - pad, 170], fill=(201, 168, 76))
    for i, p in enumerate(pngs):
        r, c = divmod(i, cols)
        x = pad + c * (tw + pad); y = top + r * (th + lab_h + pad)
        im.paste(Image.open(p).convert("RGB").resize((tw, th), Image.LANCZOS), (x, y))
        d.rectangle([x, y, x + tw, y + th], outline=(214, 204, 186), width=2)
        if labels:
            d.text((x, y + th + 12), labels[i], font=f(20, "Regular"),
                   fill=(90, 100, 118))
    im.save(path); return path


# Every slide's copy, verbatim from Section 5 of the production package.
SEC5 = {
    1:  ["CHANGING INDUSTRIES", "DOES NOT MAKE YOU",
         "ENTRY-LEVEL AT EVERYTHING.", "IT MAKES YOU NEW", "TO A CONTEXT."],
    2:  ["WHAT ACTUALLY CHANGES", "Language", "Stakeholders", "Incentives",
         "Regulation", "Operating rhythm", "Risks",
         "NONE OF THAT IS YOUR COMPETENCE.", "ALL OF IT IS CONTEXT.",
         "And context is learnable."],
    3:  ["MOVE ONE", "1", "CAPABILITY",
         "Judgment that stays useful when the setting changes.",
         "2", "CONTEXT", "What the new field knows and you do not — yet.",
         "3", "CREDENTIAL",
         "Formal evidence or permission the destination requires."],
    4:  ["CAPABILITY", "WHAT TRAVELS",
         "Finding the real decision in a room of stated positions.",
         "Structuring a problem nobody has framed yet.",
         "Knowing when a number is too clean to be true.",
         "Holding a position when the room would rather you did not.",
         "None of these belong to an industry."],
    5:  ["CONTEXT", "WHAT MUST BE RELEARNED",
         "The names the same idea goes by here.",
         "Who must be convinced, and who can quietly stop it.",
         "What the pressure is, and where it comes from.",
         "What everyone here is quietly afraid of.",
         "It feels like a competence gap. It is an information gap."],
    6:  ["CREDENTIAL", "WHAT MUST BE EARNED", "Genuinely required",
         "There is no way around it. Plan for it.", "A signal",
         "Not required, but it shortens the conversation.", "Neither",
         "Chased because it feels like progress."],
    7:  ["MOVE TWO", "START FROM THE DESTINATION,", "NOT FROM YOUR CV.",
         "What does that industry keep failing to solve?",
         "Which of those decisions have you made before?",
         "Which patterns do you already recognise early?",
         "Which constraints have you worked inside?"],
    8:  ["TRANSLATE, DO NOT RECITE", "“ADAPTABLE.”", "“STRATEGIC.”",
         "“A FAST LEARNER.”",
         "These describe how you would like to be seen.",
         "THE UNIT THAT TRANSFERS IS JUDGMENT,", "WITH EVIDENCE ATTACHED.",
         "Not a job title. Not a task list."],
    9:  ["MOVE THREE", "BRIDGE EVIDENCE", "A RELEVANT PROJECT",
         "CROSS-FUNCTIONAL WORK", "RESEARCH TURNED INTO A POINT OF VIEW",
         "A WORK SAMPLE", "PRACTITIONER CONVERSATIONS",
         "A FIRST-90-DAYS PLAN"],
    10: ["ONE PAGE, THREE COLUMNS", "WHAT", "TRAVELS",
         "Judgment, decisions,", "patterns, ways of working.",
         "CHANGES", "Language, stakeholders,", "incentives, rhythm, risk.",
         "WHAT I", "MUST EARN", "The credential or exposure",
         "nothing substitutes for."],
    11: ["CAPABILITY FORMATION", "FIELD KIT",
         "What has your work built in you?",
         "How portable is it? What is still missing?",
         "temidayoafonja.com/fieldkit"],
    12: ["CONTINUE THE SERIES", "CAREER", "PORTABILITY",
         "CAREER PIVOTS  ·  INTERNAL MOVES  ·  GROWTH"],
}


def main():
    for dpath in (OUT, PNG, REV, RENDER):
        os.makedirs(dpath, exist_ok=True)
    for dpath in (PNG, REV):
        shutil.rmtree(dpath); os.makedirs(dpath)

    # ---------------------------------------------------------- main deck
    cvs = canvases([(n, S.STEPS[n]) for n in range(1, N + 1)])
    main_pptx = os.path.join(OUT, MAIN + ".pptx")
    preview_pdf = os.path.join(OUT, PREVIEW + ".pdf")
    render_pptx(cvs, main_pptx)
    html = render_html(cvs, os.path.join(RENDER, "deck.html"),
                       "How to Move Into a New Industry Without Starting Over")
    names = ["slide-%02d-%s.png" % (n, S.SLUGS[n]) for n in range(1, N + 1)]
    pngs = shoot(html, PNG, names, pdf_path=preview_pdf)

    # -------------------------------------------------------- reveal builds
    rcvs = canvases(PLAN)
    rec_pptx = os.path.join(OUT, RECDECK + ".pptx")
    render_pptx(rcvs, rec_pptx)
    rhtml = render_html(rcvs, os.path.join(RENDER, "reveals.html"),
                        "Video 8 reveal builds")
    rnames = ["frame-%02d-slide-%02d-build-%d.png" % (i + 1, n, s)
              for i, (n, s) in enumerate(PLAN)]
    rpngs = shoot(rhtml, REV, rnames)

    sheet(pngs, os.path.join(OUT, "contact-sheet.png"), 560,
          "HOW TO MOVE INTO A NEW INDUSTRY WITHOUT STARTING OVER",
          "Video 8   ·   12 slides   ·   1920 x 1080   ·   16:9",
          labels=["%02d  %s" % (n, S.TITLES[n]) for n in range(1, N + 1)])
    rlabels = []
    for i, (n, s) in enumerate(PLAN):
        t = "%02d  slide %d" % (i + 1, n)
        if S.STEPS[n] > 1:
            t += ", state %d of %d" % (s, S.STEPS[n])
        rlabels.append(t)
    sheet(rpngs, os.path.join(OUT, "reveal-order.png"), 420,
          "REVEAL BUILDS, ADVANCE ORDER",
          "Video_8_Reveal_Builds   ·   %d frames   ·   duplicate sequential "
          "slides, no animations" % len(rpngs), labels=rlabels)
    sheet(pngs, os.path.join(OUT, "phone-check.png"), 320,
          "PHONE LEGIBILITY CHECK", "Each slide at 320 x 180.")

    # ------------------------------------------------------------------ QA
    from pptx import Presentation
    import pymupdf
    rep = {}
    prs = Presentation(main_pptx)
    flat = lambda t: " ".join(t.split())
    vis = [flat("\n".join(sh.text_frame.text for sh in sl.shapes
                          if sh.has_text_frame and sh.text_frame.text.strip()))
           for sl in prs.slides]

    rep["main_slides"] = len(prs.slides)
    rep["reveal_slides"] = len(Presentation(rec_pptx).slides)
    rep["reveal_plan"] = PLAN
    doc = pymupdf.open(preview_pdf)
    rep["pdf_pages"] = doc.page_count
    rep["pdf_size_in"] = [round(doc[0].rect.width / 72, 3),
                          round(doc[0].rect.height / 72, 3)]

    missing = {}
    for n, phrases in SEC5.items():
        gone = [p for p in phrases if flat(p) not in vis[n - 1]]
        if gone:
            missing[n] = gone
    rep["section5_missing_phrases"] = missing
    rep["section5_all_present"] = not missing

    rep["no_title_card"] = "How to Move Into a New Industry" not in " ".join(vis)
    rep["reveal_spec_matches"] = ([S.STEPS[n] for n in range(1, 13)]
                                  == [2, 2, 3, 2, 2, 1, 2, 2, 3, 3, 1, 1])

    rep["text_shapes"] = sum(1 for sl in prs.slides for sh in sl.shapes
                             if sh.has_text_frame and sh.text_frame.text.strip())
    rep["picture_shapes"] = sum(1 for sl in prs.slides for sh in sl.shapes
                                if sh.shape_type == 13)

    # A shape ending exactly on the canvas edge is a flush-edge full-bleed
    # panel, not an overflow. A checker using >= rather than > flags those, so
    # the two cases are separated and both are reported.
    over, flush = [], []
    for deck_name, pth in (("main", main_pptx), ("reveals", rec_pptx)):
        pp = Presentation(pth)
        SWp, SHp = pp.slide_width, pp.slide_height
        for i, sl in enumerate(pp.slides, start=1):
            for sh in sl.shapes:
                r, b = sh.left + sh.width, sh.top + sh.height
                if sh.left < 0 or sh.top < 0 or r > SWp or b > SHp:
                    over.append((deck_name, i, str(sh.shape_type),
                                 sh.left, sh.top, sh.width, sh.height))
                elif sh.left == 0 or sh.top == 0 or r == SWp or b == SHp:
                    flush.append({"deck": deck_name, "slide": i,
                                  "left": sh.left, "top": sh.top,
                                  "right": r, "bottom": b,
                                  "slide_width": SWp, "slide_height": SHp,
                                  "overhang_emu": max(0, r - SWp, b - SHp)})
    rep["shapes_outside_canvas"] = over
    rep["no_shape_overflow"] = not over
    rep["flush_to_canvas_edge"] = flush
    rep["flush_edge_slides_main"] = sorted({f["slide"] for f in flush
                                            if f["deck"] == "main"})
    rep["flush_edge_frames_reveals"] = sorted({f["slide"] for f in flush
                                               if f["deck"] == "reveals"})
    rep["flush_edge_max_overhang_emu"] = max([f["overhang_emu"]
                                              for f in flush] or [0])
    rep["flush_edges_are_benign"] = rep["flush_edge_max_overhang_emu"] == 0

    rep["cta_name_on_slide_11"] = ("CAPABILITY FORMATION" in vis[10].upper()
                                   and "FIELD KIT" in vis[10].upper())
    rep["cta_url_slides"] = [i + 1 for i, t in enumerate(vis)
                             if "temidayoafonja.com/fieldkit" in t]
    joined = " ".join(vis).lower()
    rep["no_competing_offer"] = not any(x in joined for x in
        ("keep the proof", "keep-the-proof", "career decision evidence check",
         "career-decisions", "maven", "book", "workshop"))
    rep["slide12_names_playlist_only"] = ("CAREER PORTABILITY" in vis[11].upper()
                                          and "CONTINUE THE SERIES" in vis[11].upper())
    rep["slide12_names_no_unpublished_video"] = not any(
        x in vis[11].upper() for x in ("LAYOFF", "VIDEO 9", "WHAT TO DO BEFORE"))
    rev_vis = [flat("\n".join(sh.text_frame.text for sh in sl.shapes
                              if sh.has_text_frame and sh.text_frame.text.strip()))
               for sl in Presentation(rec_pptx).slides]
    rep["watch_next_in_last_frame"] = "CAREER PORTABILITY" in rev_vis[-1].upper()
    rep["final_states_match_main"] = all(
        vis[n - 1] == rev_vis[sum(S.STEPS[k] for k in range(1, n + 1)) - 1]
        for n in range(1, 13))

    def bounds_ok(cv):
        for el in cv.els:
            if el["x"] < 0 or el["y"] < 0 or el["x"] + el["w"] > W:
                return False
        return True
    rep["all_elements_on_canvas"] = all(bounds_ok(c) for c in cvs)

    end = ENDCARD_CLEAR
    def full_bleed(el):
        return el["x"] <= 0 and el["y"] <= 0 and el["w"] >= W and el["h"] >= H
    intrude = [el for el in cvs[11].els
               if not full_bleed(el) and el["x"] + el["w"] > end["x"]]
    rep["slide12_endcard_clear"] = not intrude

    seams = []
    for n in [k for k in S.STEPS if S.STEPS[k] > 1]:
        states = []
        for s in range(1, S.STEPS[n] + 1):
            cv = Canvas(n, s); S.BUILDERS[n](cv, s); states.append(cv)
        for a, b in zip(states, states[1:]):
            ka = [(e["t"], e["x"], e["y"]) for e in a.els]
            kb = [(e["t"], e["x"], e["y"]) for e in b.els]
            if not all(k in kb for k in ka):
                seams.append(n)
    rep["reveal_positions_stable"] = not seams
    rep["reveal_seam_failures"] = sorted(set(seams))

    json.dump(rep, open(os.path.join(OUT, "qa.json"), "w"), indent=2)
    for k, v in rep.items():
        if k != "reveal_plan":
            print("%-34s %s" % (k, v))
    print("frames: %d" % len(PLAN))


if __name__ == "__main__":
    main()
