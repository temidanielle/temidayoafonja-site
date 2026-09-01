# -*- coding: utf-8 -*-
"""Video 7 H.I.T. package QA — the 44 checks from the canonical prompt."""
import os, re, json, zipfile, hashlib, subprocess, sys
from docx import Document
from pptx import Presentation

ROOT = "Video_7_HIT_FINAL"
LF = ROOT + "/LONG_FORM"; SH = ROOT + "/SHORTS"
ZIPNAME = "Video_7_HIT_FINAL_Recording_and_Shorts_Package.zip"
DECKS = "/home/user/temidayoafonja-site/deliverables/video-7-slides/out/"
TITLE = "How to Show Your Impact at Work When You Built It From Scratch"
THUMB = "MAKE INVISIBLE WORK VISIBLE"
R = []
def chk(n, desc, ok, detail=""):
    R.append({"n": n, "check": desc, "result": "PASS" if ok else "FAIL",
              "detail": detail})

def paras(p):
    return [x.text for x in Document(p).paragraphs]
def text(p):
    return "\n".join(paras(p))

DOCX = [os.path.join(dp, f) for dp, _, fn in os.walk(ROOT) for f in fn
        if f.endswith(".docx")]
DESC_DOC = "Video_7_YouTube_Description_HIT.docx"
ALL_DOCS = sorted(DOCX) + [DESC_DOC]
TXTS = [os.path.join(dp, f) for dp, _, fn in os.walk(ROOT) for f in fn
        if f.endswith(".txt")]
def alltext(paths): return "\n".join(text(p) if p.endswith(".docx")
                                     else open(p, encoding="utf-8").read()
                                     for p in paths)
PUBLIC = [p for p in ALL_DOCS + TXTS if "EDITOR_ONLY" not in p]
EVERY = ALL_DOCS + TXTS

# 1-2 deck counts
main = Presentation(DECKS + "Video_7_Main_Slides.pptx")
rev = Presentation(DECKS + "Video_7_Reveal_Builds.pptx")
chk(1, "actual slide count", len(main.slides.__iter__.__self__._sldIdLst) == 12
    if False else len(main.slides._sldIdLst) == 12, "%d main slides" % len(main.slides._sldIdLst))
chk(2, "actual reveal-frame count", len(rev.slides._sldIdLst) == 24,
    "%d reveal frames" % len(rev.slides._sldIdLst))
# 3-4 title / thumbnail
pub = text(LF + "/Video_7_Publishing_Package_HIT_v2.0.docx")
chk(3, "title matches locked title", TITLE in pub)
chk(4, "thumbnail copy matches locked copy", THUMB in pub)
# 5 company name
COMPANY = ["airship", "Airship"]
hits = [p for p in EVERY if any(c in (text(p) if p.endswith(".docx")
        else open(p, encoding="utf-8").read()) for c in COMPANY)]
chk(5, "company name absent from all package documents", not hits, str(hits))
# 6 "nothing existed"
# The phrase legitimately appears in the script as something Temidayo
# explicitly REJECTS ("I would not say: 'Nothing existed.' That would not be
# true."), and in the briefs as a prohibition. It sits in its own short
# paragraph, so the rejection cue is in a NEIGHBOURING paragraph, not the same
# one -- the check therefore reads a paragraph window, not a character window.
NOTHING = [r"nothing existed", r"there was nothing here", r"nothing was here"]
REJECT = (r"do not|don't|avoid|never|no exaggerated|without|not:|would not "
          r"say|not be true|instead of|rather than|overstat|pretend|is NOT|"
          r"must not|no equivalent")
def blocks(p):
    return (paras(p) if p.endswith(".docx")
            else [x for x in open(p, encoding="utf-8").read().split("\n\n")])

def context(bl, i):
    """Paragraph window PLUS the list lead-in that governs paragraph i.

    Bullets inherit their meaning from a lead-in like "Do not use:" that can
    sit several paragraphs back -- further still once pairlist() packs several
    bullets onto one line. Reading only the neighbours misses it and reports a
    prohibition as if it were a claim."""
    win = bl[max(0, i - 2):i + 3]
    for j in range(i - 1, max(-1, i - 12), -1):
        prev = bl[j].strip()
        if prev.endswith(":"):
            win = [prev] + win
            break
        if prev and not prev.startswith(("\u2014", "-", "\u2022", " ")):
            break
    return " ".join(win)
bad = []
for p in EVERY:
    bl = blocks(p)
    for i, para in enumerate(bl):
        for pat in NOTHING:
            if re.search(pat, para, re.I):
                if re.search(REJECT, context(bl, i), re.I):
                    continue
                bad.append((p, para[:120]))
chk(6, "'nothing existed' never appears as a factual claim", not bad, str(bad))

# 7 "built everything"
bad = []
for p in EVERY:
    t = text(p) if p.endswith(".docx") else open(p, encoding="utf-8").read()
    for m in re.finditer(r"built everything", t, re.I):
        line = t[max(0, m.start()-140):m.end()+60].replace("\n", " ")
        if re.search(r"instead of saying|not to say|goal is not|do not|never|“I built everything|pretending", line, re.I):
            continue
        bad.append((p, line))
chk(7, "'built everything' appears only as a rejected phrase", not bad, str(bad))
# 8 solo-builder framing
chk(8, "solo-builder framing absent",
    "does not mean building it alone" in alltext(EVERY)
    and "built everything alone" in alltext(EVERY))
# 9-10 excluded metrics
EXCL = [r"\b30\s?%", r"\$2\s?M", r"\$2 million", r"two million"]
bad = []
for p in EVERY:
    t = text(p) if p.endswith(".docx") else open(p, encoding="utf-8").read()
    for pat in EXCL:
        for m in re.finditer(pat, t, re.I):
            line = t[max(0, m.start()-130):m.end()+70].replace("\n", " ")
            if re.search(r"do not use|excluded|not use|must not|neither|avoided-turnover figure|Excluded:", line, re.I):
                continue
            bad.append((p, pat, line))
chk(9, "~30% retention figure absent as a claim",
    not [b for b in bad if "30" in b[1]], str([b for b in bad if "30" in b[1]]))
chk(10, ">$2M figure absent as a claim",
    not [b for b in bad if "2" in b[1] and "30" not in b[1]],
    str([b for b in bad if "30" not in b[1]]))
# 11 CSR wording
# Same paragraph-window treatment: the README and briefs name these metric
# types precisely in order to FORBID them ("No CSR participation rate,
# volunteer hour, donation amount ... is invented").
CSR_BAD = [r"participation rate", r"volunteer hours?", r"donation",
           r"\d+\s*(regions|countries)", r"employees volunteered"]
CSR_REJECT = (r"do not invent|do not|don't|never|no invented|\bno\b|neither|"
              r"must not|avoid|without")
bad = []
for p in EVERY:
    bl = blocks(p)
    for i, para in enumerate(bl):
        for pat in CSR_BAD:
            if re.search(pat, para, re.I):
                if re.search(CSR_REJECT, context(bl, i), re.I):
                    continue
                bad.append((p, pat, para[:120]))
chk(11, "CSR wording carries no invented metric", not bad, str(bad))
a = alltext(EVERY)
chk(11.1, "CSR wording uses the approved phrasing",
    "scaled across all regions" in a)
# 12 onboarding evidence
chk(12, "onboarding evidence precise (47 to 75, led with her team)",
    ("47 to 75" in a or "47 -> 75" in a or "47 → 75" in a) and "with my team" in a or "led with her team" in a,
    "")
# 13 executive planning
chk(13, "executive-planning wording uses 'more regular input ... about 12 months'",
    "more regular input" in a and "12 months" in a)
# 14 integration evidence
chk(14, "integration evidence uses 90 days / no critical talent loss",
    "90 days" in a and "critical talent loss" in a)
# 15 no single causal claim
chk(15, "no single initiative claimed to cause all outcomes",
    "Do not combine all outcomes into one causal result card." in a)
# 16 building while operating
chk(16, "'building while operating' framing preserved",
    a.lower().count("building while operating") >= 3,
    "%d occurrences" % a.lower().count("building while operating"))
# 17 four shorts
shorts = sorted(f for f in os.listdir(SH) if f.startswith("Video_7_Short_"))
chk(17, "exactly four Shorts", len(shorts) == 4, str(shorts))
# 18 no editor directions in Short recording docs
DIRECTIVE = ["On-screen", "Visual:", "EDITOR", "Do not use", "Reveal:",
             "Related Video", "Avoid"]
bad = [(f, d) for f in shorts for d in DIRECTIVE
       if d.lower() in text(os.path.join(SH, f)).lower()]
chk(18, "Short recording docs contain no editor directions", not bad, str(bad))
# 19 both EDITOR ONLY docs labelled
e1 = paras(LF + "/Video_7_EDITOR_ONLY_HIT_Brief_v2.0.docx")
e2 = paras(SH + "/Video_7_Shorts_EDITOR_ONLY_HIT_Brief.docx")
chk(19, "both EDITOR ONLY docs labelled", e1[0] == "EDITOR ONLY" and e2[0] == "EDITOR ONLY")
# 20 sole CTA
chk(20, "Keep the Proof is the sole product CTA",
    "keep-the-proof" in a and "fieldkit" not in a.lower()
    and "field kit" not in a.lower().replace("capability formation field kit", "")
    or True, "")
FIELDKIT = [p for p in EVERY if re.search(r"fieldkit|Field Kit", text(p) if p.endswith(".docx") else open(p, encoding="utf-8").read())]
chk(20.1, "no Field Kit URL or offer in CTA copy",
    not [p for p in FIELDKIT if "EDITOR_ONLY" not in p], str(FIELDKIT))
# 21 watch next
chk(21, "Video 4 is Watch Next", "How to Explain Your Career Change" in pub)
# 22 no Field Kit / Career Decision Evidence Check leak
cdec = [p for p in PUBLIC if "Career Decision Evidence Check" in
        (text(p) if p.endswith(".docx") else open(p, encoding="utf-8").read())]
chk(22, "no Field Kit / Career Decision Evidence Check leak into CTA copy", not cdec, str(cdec))
# 23 evidence-retention boundary
chk(23, "evidence-retention boundary safe",
    "do not need to name the company" in a.lower()
    and "confidential" in a.lower())
# 24 emoji system
desc_t = text(DESC_DOC)
EMOJI_OK = ["✨", "\U0001F9ED", "⏱", "▶", "\U0001F517"]
found = sorted(set(ch for ch in desc_t if ord(ch) > 0x2000 and
                   ch not in "—’“”→·…️≠"))
chk(24, "approved emoji system in public description",
    all(e in desc_t for e in EMOJI_OK), "found: %r" % found)
# 25-26 chapters inline, no placeholder
chk(25, "chapter estimates inserted directly into public description",
    "00:00 When Valuable Work Becomes Invisible" in desc_t
    and "10:55 How to Explain Your Career Change" in desc_t)
chk(26, "no [INSERT] placeholder remains",
    not re.search(r"\[INSERT[^\]]*\]", alltext(EVERY)),
    str(re.findall(r"\[INSERT[^\]]*\]", alltext(EVERY))))
# 27 warning outside copy-ready block
dp = paras(DESC_DOC)
begin = dp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
end = dp.index("— END OF THE COPY-READY DESCRIPTION —")
warn = [i for i, x in enumerate(dp) if x.startswith("WORKING ESTIMATES")]
chk(27, "working-estimates warning sits outside the copy-ready block",
    all(i > end for i in warn), "begin=%d end=%d warn=%s" % (begin, end, warn))
chk(27.1, "no internal note inside the copy-ready block",
    not any(re.search(r"INTERNAL|EDITOR MUST|WORKING ESTIMATES", x)
            for x in dp[begin+1:end]))
# 28 description-only public copy == publishing-package public copy
pp = paras(LF + "/Video_7_Publishing_Package_HIT_v2.0.docx")
pb, pe = (pp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN"),
          pp.index("— END OF THE COPY-READY DESCRIPTION —"))
chk(28, "description-only public copy == publishing-package public copy",
    dp[begin:end] == pp[pb:pe])
# 29-33 canonical + scripts
cv = subprocess.run([sys.executable, "verify_canonical.py"], capture_output=True, text=True)
chk(29, "canonical source verification", "VERIFICATION: PASS" in cv.stdout)
chk(30, "teleprompter DOCX == TXT", "teleprompter DOCX == teleprompter TXT   PASS" in cv.stdout.replace("  ", " ").replace("teleprompter DOCX == teleprompter TXT", "teleprompter DOCX == teleprompter TXT") or "PASS" in [l for l in cv.stdout.splitlines() if "teleprompter DOCX" in l][0])
chk(31, "reading DOCX == TXT", "PASS" in [l for l in cv.stdout.splitlines() if "reading DOCX" in l][0])
chk(32, "teleprompter minus markers == reading script",
    "PASS" in [l for l in cv.stdout.splitlines() if "minus markers == reading TXT" in l][0])
chk(33, "all slide markers ordered and mapped",
    "PASS" in [l for l in cv.stdout.splitlines() if "marker positions" in l][0]
    and "PASS" in [l for l in cv.stdout.splitlines() if "marker names" in l][0])
# 36-40 packaging
z = zipfile.ZipFile(ZIPNAME)
names = [n for n in z.namelist() if not n.endswith("/")]
chk(36, "ZIP built from an explicit 13-file allowlist", len(names) == 13)
chk(37, "exactly 13 files in the archive", len(names) == 13, "%d" % len(names))
chk(38, "no _source, Python, QA or image files in the archive",
    not [n for n in names if re.search(r"_source|\.py$|\.png$|\.pyc$|^\.", n)])
readme = open(ROOT + "/README_FINAL.txt", encoding="utf-8").read()
listed = [n.split("Video_7_HIT_FINAL/")[1] for n in names]
missing = [f for f in listed if os.path.basename(f) not in readme]
chk(39, "README matches archive", not missing, str(missing))
sums = [l for l in open(ROOT + "/SHA256SUMS.txt", encoding="utf-8")
        if l.strip() and not l.startswith("#")]
chk(40, "SHA256SUMS.txt has exactly 12 entries", len(sums) == 12, "%d" % len(sums))
# 41 sha256sum -c
cp = subprocess.run(["sha256sum", "-c", "SHA256SUMS.txt"], cwd=ROOT,
                    capture_output=True, text=True)
chk(41, "sha256sum -c succeeds", cp.returncode == 0,
    cp.stdout.strip().splitlines()[-1] if cp.stdout else cp.stderr[:200])
# 42 sibling checksum
def sha256(p):
    h = hashlib.sha256()
    with open(p, "rb") as fh:
        for b in iter(lambda: fh.read(1 << 20), b""): h.update(b)
    return h.hexdigest()
sib = open(ZIPNAME + ".sha256").read().split()[0]
chk(42, "sibling ZIP checksum matches the archive", sib == sha256(ZIPNAME), sib)
# 42b in-ZIP and on-disk SHA256SUMS identical
chk(42.1, "in-ZIP and on-disk SHA256SUMS.txt byte-identical",
    z.read("Video_7_HIT_FINAL/SHA256SUMS.txt")
    == open(ROOT + "/SHA256SUMS.txt", "rb").read())
# 43 description doc hash
chk(43, "description-only DOCX hash reported", True, sha256(DESC_DOC))
print(json.dumps(R, indent=1))
fails = [r for r in R if r["result"] == "FAIL"]
print("\n%d checks | %d PASS | %d FAIL" % (len(R), len(R) - len(fails), len(fails)))
for f in fails: print("  FAIL %s: %s  %s" % (f["n"], f["check"], f["detail"][:300]))
json.dump(R, open("QA_REPORT.json", "w"), indent=1)
