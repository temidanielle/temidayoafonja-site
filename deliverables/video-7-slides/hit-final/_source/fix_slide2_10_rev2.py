# -*- coding: utf-8 -*-
"""Apply the revised Slide 2 and Slide 10 wording to both Video 7 decks.

This supersedes the earlier Slide 2 / Slide 10 wording. Titles are untouched
("WHY THE WORK GOES QUIET", "THREE PARAGRAPHS") and every replacement clones
the existing run's rPr, so typography family, size, weight, colour and
text-box position are carried over unchanged.

Fit, measured against the real embedded fonts first:
  slide 2 headline  9.11in in an 11.11in box  -> one line
  slide 2 sub-copy  7.18in / 2.32in           -> two lines, mirroring the
                    original block's own split point
  slide 10 rows     the BEFORE and BUILD sentences are 8.29in and 8.10in in a
                    7.36in box, so they cannot hold one line. The box already
                    ends at 12.22in, the design's right margin, so it cannot be
                    widened. They are therefore broken at a comma into two
                    lines each -- the narrow fit adjustment the brief allows.
                    Two lines end 0.12in above the divider rule, so nothing
                    collides. RETURN is 6.10in and stays on one line.
"""
import copy
from pptx import Presentation
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def q(t): return "{%s}%s" % (A, t)

EDITS = [
    # slide 2
    ("FOUNDATIONAL WORK CAN LOSE ITS “BEFORE.”",
     ["FOUNDATIONAL WORK CAN BE HARD TO SEE."]),
    ("Once the mechanism works,\x0bthe starting condition becomes harder to see.",
     ["Sometimes the mechanism that would have recorded it",
      "was still maturing."]),
    # slide 10
    ("What was incomplete, inconsistent or difficult?",
     ["What was incomplete, inconsistent,", "or difficult before the work?"]),
    ("What you created underneath the output.",
     ["What did you help put in place,", "improve, or make more usable?"]),
    ("What is different now—and how someone else could tell.",
     ["What changed afterward because of the work?"]),
]
TARGETS = {"Video_7_Main_Slides.pptx": [2, 10],
           "Video_7_Reveal_Builds.pptx": [4, 20, 21, 22]}

def set_lines(shape, lines):
    body = shape.text_frame._txBody
    ps = body.findall(q("p"))
    assert len(ps) == 1, "expected a single paragraph"
    p = ps[0]
    runs = p.findall(q("r"))
    assert runs, "expected a run to clone"
    template = runs[0]
    for child in list(p):
        if etree.QName(child).localname in ("r", "br"):
            p.remove(child)
    for i, line in enumerate(lines):
        if i:
            etree.SubElement(p, q("br"))
        r = copy.deepcopy(template)
        r.find(q("t")).text = line
        p.append(r)

total = 0
for path, slides in TARGETS.items():
    prs = Presentation(path)
    for n in slides:
        for sh in prs.slides[n - 1].shapes:
            if not sh.has_text_frame:
                continue
            for frm, to in EDITS:
                if sh.text_frame.text == frm:
                    set_lines(sh, to)
                    total += 1
                    print("%-26s slide %-3d -> %r" % (path, n, to[0][:42]))
    prs.save(path)
print("edits applied:", total)
