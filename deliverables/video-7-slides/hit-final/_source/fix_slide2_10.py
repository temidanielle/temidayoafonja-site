# -*- coding: utf-8 -*-
"""Apply the authorised Slide 2 and Slide 10 text corrections to both decks.

Both slides carried absolute-absence framing the approved Video 7 factual
boundary rules out. Only run text changes: every replacement clones the
existing run's rPr, so typography family, size, weight, colour and the
text-box position are untouched. Widths were measured against the real
embedded fonts first (Montserrat Bold 27pt, DM Sans 20pt):

  slide 2 headline   9.61in in an 11.11in box  (replaces 9.07in) - one line
  slide 2 sub-copy   3.64in / 5.96in           (replaces 5.63in / 2.13in)
                     kept as two lines, broken at the clause boundary, so the
                     two-line composition of the block is preserved
  slide 10 column    6.09in in a 7.36in box    (replaces 2.42in) - one line,
                     and shorter than the RETURN row already in that column
"""
import copy
from pptx import Presentation
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def q(t): return "{%s}%s" % (A, t)

HEAD_FROM = "FOUNDATIONAL WORK HAS NO PRIOR STATE."
HEAD_TO = ["FOUNDATIONAL WORK CAN LOSE ITS “BEFORE.”"]
SUB_FROM = "The instrument that would have recorded it\x0bdid not exist yet."
SUB_TO = ["Once the mechanism works,",
          "the starting condition becomes harder to see."]
TMPL_FROM = "What did not exist."
TMPL_TO = ["What was incomplete, inconsistent or difficult?"]

EDITS = [(HEAD_FROM, HEAD_TO), (SUB_FROM, SUB_TO), (TMPL_FROM, TMPL_TO)]
TARGETS = {"Video_7_Main_Slides.pptx": [2, 10],
           "Video_7_Reveal_Builds.pptx": [4, 20, 21, 22]}

def set_lines(shape, lines):
    """Rewrite the frame as `lines`, separated by <a:br/>, cloning run props."""
    body = shape.text_frame._txBody
    ps = body.findall(q("p"))
    assert len(ps) == 1, "expected a single paragraph"
    p = ps[0]
    runs = p.findall(q("r"))
    assert runs, "expected at least one run to clone"
    template = runs[0]
    for child in list(p):
        if etree.QName(child).localname in ("r", "br"):
            p.remove(child)
    for i, line in enumerate(lines):
        if i:
            etree.SubElement(p, q("br"))
        r = copy.deepcopy(template)
        r.find(q("t")).text = line
        p.append(r)

total = 0
for path, slides in TARGETS.items():
    prs = Presentation(path)
    for n in slides:
        for sh in prs.slides[n - 1].shapes:
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text
            for frm, to in EDITS:
                if txt == frm:
                    set_lines(sh, to)
                    total += 1
                    print("%-28s slide %-3d %r -> %r"
                          % (path, n, frm[:38], to[0][:38]))
    prs.save(path)
print("edits applied:", total)
