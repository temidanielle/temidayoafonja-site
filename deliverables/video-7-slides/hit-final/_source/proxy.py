"""Geometric proxy renderer for PPTX slides.

LibreOffice cannot render anything in this environment, so to *see* a slide we
reconstruct it in HTML at true EMU positions with the real embedded fonts and
screenshot it with Chromium. Colour resolution is the part that has bitten us:
a run whose colour comes from the theme has no <a:srgbClr>, and if we default
such runs to inherit they paint in the page background and silently vanish.
So we resolve explicit -> scheme-via-theme -> NAVY, and never leave it unset.
"""
import sys, os, re, subprocess
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def q(t): return "{%s}%s" % (A, t)

NAVY = "0F2346"
EMU_IN = 914400

def theme_map(prs):
    """srgb hex for each scheme colour name, read from the presentation theme."""
    out = {}
    try:
        part = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        root = etree.fromstring(part.blob)
    except Exception:
        return out
    scheme = root.find(".//" + q("clrScheme"))
    if scheme is None:
        return out
    for child in scheme:
        name = etree.QName(child).localname
        srgb = child.find(q("srgbClr"))
        sysc = child.find(q("sysClr"))
        if srgb is not None:
            out[name] = srgb.get("val")
        elif sysc is not None:
            out[name] = sysc.get("lastClr") or "000000"
    # dk1/lt1 are addressed as tx1/bg1 (and dk2/lt2 as tx2/bg2) from runs
    for a, b in (("dk1", "tx1"), ("lt1", "bg1"), ("dk2", "tx2"), ("lt2", "bg2")):
        if a in out:
            out[b] = out[a]
    return out

def run_colour(r, tmap):
    rPr = r.find(q("rPr"))
    if rPr is not None:
        fill = rPr.find(q("solidFill"))
        if fill is not None:
            srgb = fill.find(q("srgbClr"))
            if srgb is not None:
                return "#" + srgb.get("val")
            sch = fill.find(q("schemeClr"))
            if sch is not None:
                return "#" + tmap.get(sch.get("val"), NAVY)
    return "#" + NAVY

def run_props(r):
    rPr = r.find(q("rPr"))
    sz, bold, face = 1800, False, "DM Sans"
    if rPr is not None:
        if rPr.get("sz"): sz = int(rPr.get("sz"))
        bold = rPr.get("b") == "1"
        lat = rPr.find(q("latin"))
        if lat is not None and lat.get("typeface"):
            face = lat.get("typeface")
    return sz / 100.0, bold, face

def walk(shapes, out, dx=0, dy=0):
    for sh in shapes:
        if sh.shape_type == 6:  # group
            walk(sh.shapes, out, dx + (sh.left or 0), dy + (sh.top or 0))
        else:
            out.append((sh, dx, dy))
    return out

def render(pptx_path, indices, png_path, scale=110):
    prs = Presentation(pptx_path)
    tmap = theme_map(prs)
    W = prs.slide_width / EMU_IN
    H = prs.slide_height / EMU_IN
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    parts = []
    for idx in indices:
        slide = prs.slides[idx - 1]
        # Slides carry their own <p:bg>; painting a fixed page colour made
        # light text on the navy slides render invisibly against cream.
        bg = "#F3F0E8"
        bgel = slide._element.find(".//{%s}bg" % P)
        if bgel is not None:
            c = bgel.find(".//" + q("srgbClr"))
            if c is not None:
                bg = "#" + c.get("val")
        blocks = []
        for sh, dx, dy in walk(slide.shapes, []):
            L = ((sh.left or 0) + dx) / EMU_IN
            T = ((sh.top or 0) + dy) / EMU_IN
            Wd = (sh.width or 0) / EMU_IN
            Ht = (sh.height or 0) / EMU_IN
            # solid-filled shapes with no text are rules / bands: draw them
            if not sh.has_text_frame:
                fill_hex = None
                try:
                    sp = sh._element.find(".//" + q("solidFill"))
                    if sp is not None:
                        s = sp.find(q("srgbClr"))
                        if s is not None: fill_hex = "#" + s.get("val")
                except Exception:
                    pass
                if fill_hex:
                    blocks.append(
                        "<div style='position:absolute;left:%.4fin;top:%.4fin;"
                        "width:%.4fin;height:%.4fin;background:%s'></div>"
                        % (L, T, Wd, max(Ht, 0.01), fill_hex))
                continue
            paras = []
            for p in sh.text_frame._txBody.findall(q("p")):
                pieces = []
                for node in p:
                    tag = etree.QName(node).localname
                    if tag == "br":
                        pieces.append("<br>")
                    elif tag == "r":
                        t = node.find(q("t"))
                        if t is None or not t.text:
                            continue
                        sz, bold, face = run_props(node)
                        pieces.append(
                            "<span style='font-family:\"%s\";font-size:%.1fpt;"
                            "font-weight:%s;color:%s'>%s</span>"
                            % (face, sz, "700" if bold else "400",
                               run_colour(node, tmap),
                               t.text.replace("&", "&amp;").replace("<", "&lt;")))
                if pieces:
                    paras.append("<div style='margin:0'>%s</div>" % "".join(pieces))
            if paras:
                blocks.append(
                    "<div style='position:absolute;left:%.4fin;top:%.4fin;"
                    "width:%.4fin;line-height:1.18'>%s</div>"
                    % (L, T, Wd, "".join(paras)))
        parts.append(
            "<div class='slide' style='background:%s'>"
            "<div class='n'>slide %d</div>%s</div>"
            % (bg, idx, "".join(blocks)))
    html = """<meta charset=utf-8><style>
@font-face{font-family:'Montserrat';src:url('file:///root/.fonts/MontserratTB-Bold.ttf');font-weight:700}
@font-face{font-family:'DM Sans';src:url('file:///root/.fonts/DMSans-Regular.ttf')}
body{margin:0;background:#888}
.slide{position:relative;width:%.4fin;height:%.4fin;
       margin:0 0 14px 0;overflow:hidden}
.n{position:absolute;right:6px;bottom:4px;font:10px monospace;color:#999;z-index:9}
</style>%s""" % (W, H, "".join(parts))
    open(png_path.replace(".png", ".html"), "w").write(html)
    from playwright.sync_api import sync_playwright
    with sync_playwright() as pw:
        b = pw.chromium.launch(
            executable_path="/opt/pw-browsers/chromium-1194/chrome-linux/chrome",
            args=["--no-sandbox"])
        pg = b.new_page(viewport={"width": int(W * scale), "height": 800})
        pg.goto("file://" + png_path.replace(".png", ".html"))
        pg.wait_for_timeout(600)
        pg.screenshot(path=png_path, full_page=True)
        b.close()
    print("rendered", png_path, indices)

if __name__ == "__main__":
    path = sys.argv[1]
    idxs = [int(x) for x in sys.argv[2].split(",")]
    render(path, idxs, sys.argv[3])
