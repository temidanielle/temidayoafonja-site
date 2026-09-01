"""Remove the two orphaned divider rules left behind by the Slide 5 correction.

The authorized Slide 5 change cut the item list from 8 to 7. The 8th item's
text box was removed but its divider rule was not, leaving a gold rule under
empty space in the right-hand column. These are the only two shapes removed,
addressed by explicit (slide, shape index, position) so nothing else can move.
"""
from pptx import Presentation
EMU = 914400
TARGETS = {
    "Video_7_Main_Slides.pptx":   [(5, 16, 7.08, 5.62)],
    "Video_7_Reveal_Builds.pptx": [(11, 16, 7.08, 5.62)],
}
for f, targets in TARGETS.items():
    prs = Presentation(f)
    for slide_no, shape_idx, exp_l, exp_t in targets:
        s = prs.slides[slide_no - 1]
        sh = list(s.shapes)[shape_idx]
        assert not sh.text_frame.text.strip(), "refusing: shape has text"
        assert abs(sh.left / EMU - exp_l) < 0.01 and abs(sh.top / EMU - exp_t) < 0.01, \
            "refusing: shape is not at the expected position"
        assert sh.height / EMU <= 0.02, "refusing: shape is not a thin rule"
        sh._element.getparent().remove(sh._element)
        print("%s slide %d: removed rule at %.2f,%.2f" % (f, slide_no, exp_l, exp_t))
    prs.save(f)
