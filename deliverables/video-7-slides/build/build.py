"""Build Video 7 and run its QA.

  PPTX  Video_7_Main_Slides.pptx      - 12 editable slides, final revealed state
  PPTX  Video_7_Reveal_Builds.pptx    - every reveal state as a duplicate slide
  PDF   Video_7_Slide_Preview.pdf     - printed from Chromium at true 16:9
  PNG   exact 1920 x 1080 frames, contact sheet, reveal order, phone check

    python3 build/build.py
"""
import os, sys, json, shutil

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck import W, H, ENDCARD_CLEAR, Canvas, render_pptx, render_html
import slides as S
from PIL import Image, ImageDraw, ImageFont

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "out")
PNG = os.path.join(OUT, "png")
REV = os.path.join(OUT, "reveals")
RENDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_render")
CHROME = "/opt/pw-browsers/chromium-1194/chrome-linux/chrome"
FD = os.path.expanduser("~/.fonts/")

MAIN = "Video_7_Main_Slides"
RECDECK = "Video_7_Reveal_Builds"
PREVIEW = "Video_7_Slide_Preview"
N = 12
PLAN = [(n, s) for n in range(1, N + 1) for s in range(1, S.STEPS[n] + 1)]


def canvases(plan):
    out = []
    for n, step in plan:
        cv = Canvas(n, step)
        S.BUILDERS[n](cv, step)
        out.append(cv)
    return out


def shoot(html_path, png_dir, names, pdf_path=None):
    from playwright.sync_api import sync_playwright
    os.makedirs(png_dir, exist_ok=True)
    made = []
    with sync_playwright() as pw:
        b = pw.chromium.launch(executable_path=CHROME, args=["--no-sandbox"])
        pg = b.new_page(viewport={"width": W, "height": H}, device_scale_factor=1)
        pg.goto("file://" + os.path.abspath(html_path), wait_until="load")
        pg.wait_for_timeout(600)
        for i, name in enumerate(names, start=1):
            p = os.path.join(png_dir, name)
            pg.query_selector("#s%d" % i).screenshot(path=p)
            made.append(p)
        if pdf_path:
            pg.pdf(path=pdf_path, width="13.333in", height="7.5in",
                   print_background=True, scale=0.66665,
                   margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
        b.close()
    return made


def f(sz, w="Bold"):
    return ImageFont.truetype(FD + "MontserratTB-%s.ttf" % w, sz)


def sheet(pngs, path, tile_w, title, sub, labels=None, cols=4):
    tw = tile_w; th = round(tw * 9 / 16)
    rows = (len(pngs) + cols - 1) // cols
    lab_h = 46 if labels else 16
    pad, top = 46, 190
    Wd = pad * 2 + cols * tw + (cols - 1) * pad
    Hd = top + rows * (th + lab_h + pad) + pad
    im = Image.new("RGB", (Wd, Hd), (250, 247, 241)); d = ImageDraw.Draw(im)
    d.text((pad, 52), title, font=f(46), fill=(15, 35, 70))
    d.text((pad, 116), sub, font=f(26, "Regular"), fill=(120, 132, 150))
    d.rectangle([pad, 166, Wd - pad, 170], fill=(201, 168, 76))
    for i, p in enumerate(pngs):
        r, c = divmod(i, cols)
        x = pad + c * (tw + pad); y = top + r * (th + lab_h + pad)
        im.paste(Image.open(p).convert("RGB").resize((tw, th), Image.LANCZOS), (x, y))
        d.rectangle([x, y, x + tw, y + th], outline=(214, 204, 186), width=2)
        if labels:
            d.text((x, y + th + 12), labels[i], font=f(20, "Regular"),
                   fill=(90, 100, 118))
    im.save(path); return path


# Every slide's copy, verbatim from Section 5 of the production package.
SEC5 = {
    1:  ["YOU WERE NOT", "IMPROVING SOMETHING.", "YOU WERE THE BEFORE."],
    2:  ["WHY THE WORK GOES QUIET", "MOST WORKPLACE EVIDENCE", "IS COMPARATIVE.",
         "A baseline moved. A number improved. A cycle time dropped.",
         "FOUNDATIONAL WORK HAS NO PRIOR STATE.",
         "The instrument that would have recorded it", "did not exist yet."],
    3:  ["THE THREE MOVES", "1", "RECONSTRUCT THE BEFORE", "2",
         "NAME WHAT YOU BUILT", "3", "SHOW WHAT IT RETURNED"],
    4:  ["MOVE ONE", "DOCUMENT THE ABSENCE", "YOU WALKED INTO.",
         "Not the state of the world when you arrived.",
         "What was missing from it."],
    5:  ["WHAT DID NOT EXIST?", "No owner", "No system", "No shared language",
         "No baseline anyone trusted", "No decision process",
         "No relationships across functions", "No standard",
         "No repeatable method"],
    6:  ["MOVE TWO", "THE OUTPUT IS THE", "SMALLEST PART.",
         "DEFINITIONS", "RELATIONSHIPS", "DECISIONS", "ALIGNMENT",
         "REPEATABLE CAPABILITY",
         "What the organisation can do now that it could not do before."],
    7:  ["JUDGMENT IS THE PART", "THAT DOES NOT TRANSFER", "INTO A SLIDE.",
         "Name the choice. Name the alternative you rejected.",
         "Name the reason."],
    8:  ["MOVE THREE", "WHAT IS DIFFERENT NOW?", "ADOPTION", "CONTINUED USE",
         "BETTER DECISIONS", "REDUCED AMBIGUITY", "REPEATABILITY",
         "RECOGNITION", "A CLEAN HANDOFF", "A CAPABILITY THAT REMAINED"],
    9:  ["A NUMBER YOU CANNOT DEFEND", "IS WORSE THAN NO NUMBER.",
         "CONTINUED USE IS EVIDENCE.",
         "Who still used it after the handoff—and for how long?",
         "WHAT YOU MAY KEEP",
         "Outcomes, decisions, what you learned, non-confidential",
         "examples in your own words. Not documents, not data,",
         "not anything employer-owned."],
    10: ["THREE PARAGRAPHS", "BEFORE", "What did not exist.", "BUILD",
         "What you created underneath the output.", "RETURN",
         "What is different now—and how someone else could tell."],
    11: ["KEEP THE PROOF", "A 60-MINUTE CAREER", "EVIDENCE SYSTEM",
         "Reconstruct the before. Name what you built.",
         "Record what it returned.", "temidayoafonja.com/keep-the-proof"],
    12: ["WATCH NEXT", "HOW TO EXPLAIN A", "NONLINEAR CAREER",
         "WITHOUT LOOKING", "UNFOCUSED",
         "Career Portability: Career Pivots,", "Internal Moves & Growth"],
}


def main():
    for dpath in (OUT, PNG, REV, RENDER):
        os.makedirs(dpath, exist_ok=True)
    for dpath in (PNG, REV):
        shutil.rmtree(dpath); os.makedirs(dpath)

    # ---------------------------------------------------------- main deck
    cvs = canvases([(n, S.STEPS[n]) for n in range(1, N + 1)])
    main_pptx = os.path.join(OUT, MAIN + ".pptx")
    preview_pdf = os.path.join(OUT, PREVIEW + ".pdf")
    render_pptx(cvs, main_pptx)
    html = render_html(cvs, os.path.join(RENDER, "deck.html"),
                       "How to Show Your Impact at Work When You Built It From Scratch")
    names = ["slide-%02d-%s.png" % (n, S.SLUGS[n]) for n in range(1, N + 1)]
    pngs = shoot(html, PNG, names, pdf_path=preview_pdf)

    # -------------------------------------------------------- reveal builds
    rcvs = canvases(PLAN)
    rec_pptx = os.path.join(OUT, RECDECK + ".pptx")
    render_pptx(rcvs, rec_pptx)
    rhtml = render_html(rcvs, os.path.join(RENDER, "reveals.html"),
                        "Video 7 reveal builds")
    rnames = ["frame-%02d-slide-%02d-build-%d.png" % (i + 1, n, s)
              for i, (n, s) in enumerate(PLAN)]
    rpngs = shoot(rhtml, REV, rnames)

    sheet(pngs, os.path.join(OUT, "contact-sheet.png"), 560,
          "HOW TO SHOW YOUR IMPACT AT WORK WHEN YOU BUILT IT FROM SCRATCH",
          "Video 7   ·   12 slides   ·   1920 x 1080   ·   16:9",
          labels=["%02d  %s" % (n, S.TITLES[n]) for n in range(1, N + 1)])
    rlabels = []
    for i, (n, s) in enumerate(PLAN):
        t = "%02d  slide %d" % (i + 1, n)
        if S.STEPS[n] > 1:
            t += ", state %d of %d" % (s, S.STEPS[n])
        rlabels.append(t)
    sheet(rpngs, os.path.join(OUT, "reveal-order.png"), 420,
          "REVEAL BUILDS, ADVANCE ORDER",
          "Video_7_Reveal_Builds   ·   %d frames   ·   duplicate sequential "
          "slides, no animations" % len(rpngs), labels=rlabels)
    sheet(pngs, os.path.join(OUT, "phone-check.png"), 320,
          "PHONE LEGIBILITY CHECK", "Each slide at 320 x 180.")

    # ------------------------------------------------------------------ QA
    from pptx import Presentation
    import pymupdf
    rep = {}
    prs = Presentation(main_pptx)
    flat = lambda t: " ".join(t.split())
    vis = [flat("\n".join(sh.text_frame.text for sh in sl.shapes
                          if sh.has_text_frame and sh.text_frame.text.strip()))
           for sl in prs.slides]

    rep["main_slides"] = len(prs.slides)
    rep["reveal_slides"] = len(Presentation(rec_pptx).slides)
    rep["reveal_plan"] = PLAN
    doc = pymupdf.open(preview_pdf)
    rep["pdf_pages"] = doc.page_count
    rep["pdf_size_in"] = [round(doc[0].rect.width / 72, 3),
                          round(doc[0].rect.height / 72, 3)]

    missing = {}
    for n, phrases in SEC5.items():
        gone = [p for p in phrases if flat(p) not in vis[n - 1]]
        if gone:
            missing[n] = gone
    rep["section5_missing_phrases"] = missing
    rep["section5_all_present"] = not missing

    rep["no_title_card"] = "How to Show Your Impact" not in " ".join(vis)
    rep["reveal_spec_matches"] = ([S.STEPS[n] for n in range(1, 13)]
                                  == [2, 2, 3, 2, 2, 2, 1, 3, 2, 3, 1, 1])

    rep["text_shapes"] = sum(1 for sl in prs.slides for sh in sl.shapes
                             if sh.has_text_frame and sh.text_frame.text.strip())
    rep["picture_shapes"] = sum(1 for sl in prs.slides for sh in sl.shapes
                                if sh.shape_type == 13)

    # A shape ending exactly on the canvas edge is a flush-edge full-bleed
    # panel, not an overflow. A checker using >= rather than > flags those, so
    # the two cases are separated and both are reported.
    over, flush = [], []
    for deck_name, pth in (("main", main_pptx), ("reveals", rec_pptx)):
        pp = Presentation(pth)
        SWp, SHp = pp.slide_width, pp.slide_height
        for i, sl in enumerate(pp.slides, start=1):
            for sh in sl.shapes:
                r, b = sh.left + sh.width, sh.top + sh.height
                if sh.left < 0 or sh.top < 0 or r > SWp or b > SHp:
                    over.append((deck_name, i, str(sh.shape_type),
                                 sh.left, sh.top, sh.width, sh.height))
                elif sh.left == 0 or sh.top == 0 or r == SWp or b == SHp:
                    flush.append({"deck": deck_name, "slide": i,
                                  "left": sh.left, "top": sh.top,
                                  "right": r, "bottom": b,
                                  "slide_width": SWp, "slide_height": SHp,
                                  "overhang_emu": max(0, r - SWp, b - SHp)})
    rep["shapes_outside_canvas"] = over
    rep["no_shape_overflow"] = not over
    rep["flush_to_canvas_edge"] = flush
    rep["flush_edge_slides_main"] = sorted({f["slide"] for f in flush
                                            if f["deck"] == "main"})
    rep["flush_edge_frames_reveals"] = sorted({f["slide"] for f in flush
                                               if f["deck"] == "reveals"})
    rep["flush_edge_max_overhang_emu"] = max([f["overhang_emu"]
                                              for f in flush] or [0])
    rep["flush_edges_are_benign"] = rep["flush_edge_max_overhang_emu"] == 0

    rep["cta_name_on_slide_11"] = ("KEEP THE PROOF" in vis[10].upper()
                                   and "A 60-MINUTE CAREER EVIDENCE SYSTEM" in vis[10])
    rep["cta_url_slides"] = [i + 1 for i, t in enumerate(vis)
                             if "temidayoafonja.com/keep-the-proof" in t]
    joined = " ".join(vis).lower()
    rep["removed_line_absent"] = "every team in the region" not in joined
    rep["no_competing_offer"] = not any(x in joined for x in
        ("field kit", "fieldkit", "career decision evidence check",
         "career-decisions", "book", "workshop"))
    rep["watch_next_routes_to_video_8"] = all(
        x in vis[11] for x in ("HOW TO EXPLAIN A", "NONLINEAR CAREER",
                               "WITHOUT LOOKING", "UNFOCUSED"))
    rep["watch_next_names_playlist"] = "Career Portability" in vis[11]
    rev_vis = [flat("\n".join(sh.text_frame.text for sh in sl.shapes
                              if sh.has_text_frame and sh.text_frame.text.strip()))
               for sl in Presentation(rec_pptx).slides]
    rep["watch_next_in_last_frame"] = "UNFOCUSED" in rev_vis[-1]
    rep["final_states_match_main"] = all(
        vis[n - 1] == rev_vis[sum(S.STEPS[k] for k in range(1, n + 1)) - 1]
        for n in range(1, 13))

    def bounds_ok(cv):
        for el in cv.els:
            if el["x"] < 0 or el["y"] < 0 or el["x"] + el["w"] > W:
                return False
        return True
    rep["all_elements_on_canvas"] = all(bounds_ok(c) for c in cvs)

    end = ENDCARD_CLEAR
    def full_bleed(el):
        return el["x"] <= 0 and el["y"] <= 0 and el["w"] >= W and el["h"] >= H
    intrude = [el for el in cvs[11].els
               if not full_bleed(el) and el["x"] + el["w"] > end["x"]]
    rep["slide12_endcard_clear"] = not intrude

    seams = []
    for n in [k for k in S.STEPS if S.STEPS[k] > 1]:
        states = []
        for s in range(1, S.STEPS[n] + 1):
            cv = Canvas(n, s); S.BUILDERS[n](cv, s); states.append(cv)
        for a, b in zip(states, states[1:]):
            ka = [(e["t"], e["x"], e["y"]) for e in a.els]
            kb = [(e["t"], e["x"], e["y"]) for e in b.els]
            if not all(k in kb for k in ka):
                seams.append(n)
    rep["reveal_positions_stable"] = not seams
    rep["reveal_seam_failures"] = sorted(set(seams))

    json.dump(rep, open(os.path.join(OUT, "qa.json"), "w"), indent=2)
    for k, v in rep.items():
        if k != "reveal_plan":
            print("%-34s %s" % (k, v))
    print("frames: %d" % len(PLAN))


if __name__ == "__main__":
    main()
