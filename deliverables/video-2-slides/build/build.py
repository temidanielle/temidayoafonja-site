"""Build the Video 2 deck and run its QA.

  PPTX  - editable main deck, 13 slides
  PPTX  - recording deck, every reveal state as a sequential duplicate slide
  PDF   - printed from Chromium at true 16:9 slide size
  PNG   - exact 1920 x 1080 frames, plus sheets and safe-area guides

    python3 build/build.py
"""
import os, sys, json, shutil, hashlib

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck import (W, H, SAFE, SAFE_CLEAR, ENDCARD_CLEAR, Canvas,
                  render_pptx, render_html)
import slides as S
from PIL import Image, ImageDraw, ImageFont, ImageChops

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "out")
PNG = os.path.join(OUT, "png")
REV = os.path.join(OUT, "reveals")
GUIDES = os.path.join(OUT, "guides")
RENDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_render")
CHROME = "/opt/pw-browsers/chromium-1194/chrome-linux/chrome"

DECK = "Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.0"
RECDECK = "Video-2-Reveal-Builds_v1.0"
N = 13
RECORDING_PLAN = [(n, s) for n in range(1, N + 1) for s in range(1, S.STEPS[n] + 1)]


def canvases(plan):
    out = []
    for n, step in plan:
        cv = Canvas(n, step)
        S.BUILDERS[n](cv, step)
        out.append(cv)
    return out


def shoot(html_path, png_dir, names, pdf_path=None):
    from playwright.sync_api import sync_playwright
    os.makedirs(png_dir, exist_ok=True)
    made = []
    with sync_playwright() as pw:
        b = pw.chromium.launch(executable_path=CHROME, args=["--no-sandbox"])
        pg = b.new_page(viewport={"width": W, "height": H}, device_scale_factor=1)
        pg.goto("file://" + os.path.abspath(html_path), wait_until="load")
        pg.wait_for_timeout(600)
        for i, name in enumerate(names, start=1):
            p = os.path.join(png_dir, name)
            pg.query_selector("#s%d" % i).screenshot(path=p)
            made.append(p)
        if pdf_path:
            pg.pdf(path=pdf_path, width="13.333in", height="7.5in",
                   print_background=True, scale=0.66665,
                   margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
        b.close()
    return made


def font(size, bold=False):
    return ImageFont.truetype(
        os.path.expanduser("~/.fonts/Montserrat-%s.ttf" % ("Bold" if bold else "Regular")),
        size)


def sheet(paths, out_path, tw, title, subtitle, labels=None, cols=5):
    th = int(tw * H / W)
    rows = (len(paths) + cols - 1) // cols
    gx, gy, pad, top = 40, (86 if labels else 46), 56, 168
    cw = pad * 2 + cols * tw + (cols - 1) * gx
    ch = top + rows * (th + gy) + pad - 24
    im = Image.new("RGB", (cw, ch), (245, 241, 232))
    d = ImageDraw.Draw(im)
    d.text((pad, 56), title, font=font(34, True), fill=(15, 35, 70))
    d.text((pad, 104), subtitle, font=font(22), fill=(90, 107, 130))
    d.rectangle([pad, 146, cw - pad, 148], fill=(201, 168, 76))
    for i, p in enumerate(paths):
        r, c = divmod(i, cols)
        x, y = pad + c * (tw + gx), top + r * (th + gy)
        im.paste(Image.open(p).resize((tw, th), Image.LANCZOS), (x, y))
        d.rectangle([x, y, x + tw - 1, y + th - 1], outline=(224, 217, 200))
        if labels:
            d.text((x, y + th + 14), labels[i], font=font(18, True), fill=(15, 35, 70))
    im.save(out_path)
    return out_path


def guide_overlay(paths, outdir):
    os.makedirs(outdir, exist_ok=True)
    for i, p in enumerate(paths, start=1):
        im = Image.open(p).convert("RGB")
        d = ImageDraw.Draw(im, "RGBA")
        d.rectangle([0, 0, SAFE_CLEAR["w"], SAFE_CLEAR["h"]],
                    fill=(193, 68, 14, 38), outline=(193, 68, 14, 200), width=3)
        d.rectangle([SAFE["x"], SAFE["y"], SAFE["x"] + SAFE["w"],
                     SAFE["y"] + SAFE["h"]], outline=(193, 68, 14, 255), width=4)
        d.text((SAFE["x"] + 18, SAFE["y"] + 18),
               "CAMERA\n480 x 270 px\n25% of frame width", font=font(24, True),
               fill=(193, 68, 14))
        if i == N:
            z = ENDCARD_CLEAR
            d.rectangle([z["x"], z["y"], z["x"] + z["w"], z["y"] + z["h"]],
                        fill=(201, 168, 76, 45), outline=(201, 168, 76, 230), width=4)
            d.text((z["x"] + 20, z["y"] + 20), "YOUTUBE END-SCREEN ELEMENT",
                   font=font(24, True), fill=(160, 128, 40))
        im.save(os.path.join(outdir, "guide-%02d.png" % i))


MEASURE_JS = """
() => {
  const out = [];
  document.querySelectorAll('.slide').forEach((slide, si) => {
    const sb = slide.getBoundingClientRect();
    const items = [];
    slide.querySelectorAll('p').forEach(p => {
      const r = p.getBoundingClientRect();
      const range = document.createRange();
      range.selectNodeContents(p);
      const tops = new Set();
      for (const cr of range.getClientRects()) tops.add(Math.round(cr.top));
      const wanted = (p.innerHTML.match(/<br>/g) || []).length + 1;
      items.push({ text: p.textContent.slice(0, 60),
        x: r.left - sb.left, y: r.top - sb.top, w: r.width, h: r.height,
        linesWanted: wanted, linesActual: tops.size });
    });
    out.push({ slide: si + 1, items });
  });
  return out;
}
"""


def measure(html_path):
    from playwright.sync_api import sync_playwright
    with sync_playwright() as pw:
        b = pw.chromium.launch(executable_path=CHROME, args=["--no-sandbox"])
        pg = b.new_page(viewport={"width": W, "height": H})
        pg.goto("file://" + os.path.abspath(html_path), wait_until="load")
        pg.wait_for_timeout(600)
        data = pg.evaluate(MEASURE_JS)
        b.close()
    return data


def overlaps(a, b, pad=0):
    return not (a["x"] + a["w"] <= b["x"] + pad or b["x"] + b["w"] <= a["x"] + pad
                or a["y"] + a["h"] <= b["y"] + pad or b["y"] + b["h"] <= a["y"] + pad)


def geometry(measured, plan=None):
    rep = {"clipped": [], "wrapped": [], "crowded": [], "safe_area": [],
           "endscreen": [], "bottom_edge": []}
    for s in measured:
        n = plan[s["slide"] - 1][0] if plan else s["slide"]
        items = s["items"]
        for it in items:
            tag = (n, it["text"])
            if (it["x"] < -1 or it["y"] < -1 or it["x"] + it["w"] > W + 1
                    or it["y"] + it["h"] > H + 1):
                rep["clipped"].append(tag)
            if it["linesActual"] > it["linesWanted"]:
                rep["wrapped"].append((n, it["text"], it["linesWanted"],
                                       it["linesActual"]))
            if it["y"] + it["h"] > H - 60:
                rep["bottom_edge"].append((n, it["text"], round(it["y"] + it["h"])))
            z = SAFE_CLEAR
            if overlaps(it, {"x": z["x"], "y": z["y"], "w": z["w"], "h": z["h"]}):
                rep["safe_area"].append(tag)
            if n == N:
                e = ENDCARD_CLEAR
                if overlaps(it, {"x": e["x"], "y": e["y"], "w": e["w"], "h": e["h"]}):
                    rep["endscreen"].append(tag)
        for i in range(len(items)):
            for j in range(i + 1, len(items)):
                if overlaps(items[i], items[j], pad=2):
                    rep["crowded"].append((n, items[i]["text"][:34],
                                           items[j]["text"][:34]))
    return rep


def fingerprint(cv):
    out = []
    for el in cv.els:
        e = dict(el)
        for k in ("fill", "line"):
            if e.get(k) is not None:
                e[k] = str(e[k])
        if e["t"] == "text":
            e["paras"] = [dict(p, color=str(p["color"])) for p in e["paras"]]
        out.append(e)
    return json.dumps(out, sort_keys=True)


def texts(cv):
    return [p["text"] for el in cv.els if el["t"] == "text" for p in el["paras"]]


def pixel_delta(a, b):
    d = ImageChops.difference(Image.open(a).convert("RGB"),
                              Image.open(b).convert("RGB")).convert("L")
    if d.getbbox() is None:
        return {"max": 0, "pixels": 0}
    return {"max": d.getextrema()[1],
            "pixels": sum(1 for v in d.get_flattened_data() if v > 0)}


def main():
    for d in (OUT, PNG, REV, GUIDES, RENDER):
        os.makedirs(d, exist_ok=True)
    for d in (PNG, REV, GUIDES):
        shutil.rmtree(d); os.makedirs(d)

    cvs = canvases([(n, S.STEPS[n]) for n in range(1, N + 1)])
    pptx_path = os.path.join(OUT, DECK + ".pptx")
    pdf_path = os.path.join(OUT, DECK + ".pdf")
    render_pptx(cvs, pptx_path)
    html = render_html(cvs, os.path.join(RENDER, "deck.html"),
                       "Is Your Job Making You Less Marketable?")
    names = ["slide-%02d-%s.png" % (n, S.SLUGS[n]) for n in range(1, N + 1)]
    pngs = shoot(html, PNG, names, pdf_path=pdf_path)

    rcvs = canvases(RECORDING_PLAN)
    rec_pptx = os.path.join(OUT, RECDECK + ".pptx")
    render_pptx(rcvs, rec_pptx)
    rhtml = render_html(rcvs, os.path.join(RENDER, "recording.html"),
                        "Video 2 recording deck")
    rnames = ["frame-%02d-slide-%02d-build-%d.png" % (i + 1, n, s)
              for i, (n, s) in enumerate(RECORDING_PLAN)]
    rpngs = shoot(rhtml, REV, rnames,
                  pdf_path=os.path.join(OUT, RECDECK + ".pdf"))

    sheet(pngs, os.path.join(OUT, "contact-sheet-v1.0.png"), 560,
          "IS YOUR JOB MAKING YOU LESS MARKETABLE?",
          "Video 2, version 1.0   ·   13 slides   ·   1920 x 1080   ·   16:9",
          labels=["%02d  %s" % (n, S.TITLES[n]) for n in range(1, N + 1)])
    rlabels = []
    for i, (n, s) in enumerate(RECORDING_PLAN):
        t = "%02d  slide %d" % (i + 1, n)
        if S.STEPS[n] > 1:
            t += ", build %d of %d" % (s, S.STEPS[n])
        rlabels.append(t)
    sheet(rpngs, os.path.join(OUT, "recording-deck-order-v1.0.png"), 420,
          "RECORDING DECK, ADVANCE ORDER",
          "Video-2-Reveal-Builds_v1.0   ·   %d frames   ·   sequential duplicate "
          "slides, no animations" % len(rpngs), labels=rlabels)
    sheet(pngs, os.path.join(OUT, "phone-thumbnail-check-v1.0.png"), 320,
          "PHONE-THUMBNAIL LEGIBILITY CHECK",
          "Each slide at 320 x 180.")
    guide_overlay(pngs, GUIDES)

    # ------------------------------------------------------------------ QA
    from pptx import Presentation
    import pymupdf
    rep = {}
    prs = Presentation(pptx_path)
    flat = lambda t: " ".join(t.split())
    vis = [flat("\n".join(sh.text_frame.text for sh in sl.shapes
                          if sh.has_text_frame and sh.text_frame.text.strip()))
           for sl in prs.slides]
    nts = [flat(sl.notes_slide.notes_text_frame.text) for sl in prs.slides]

    rep["main_slides"] = len(prs.slides)
    rep["pdf_pages"] = pymupdf.open(pdf_path).page_count
    rep["pdf_size_in"] = [round(pymupdf.open(pdf_path)[0].rect.width / 72, 3),
                          round(pymupdf.open(pdf_path)[0].rect.height / 72, 3)]
    rep["recording_slides"] = len(Presentation(rec_pptx).slides)
    rep["recording_plan"] = RECORDING_PLAN

    tests = ["Remove the company nouns", "Find outside-context evidence",
             "Read the last 90 days"]
    def only(i, has, hasnt):
        return all(h in vis[i - 1] for h in has) and not any(h in vis[i - 1]
                                                             for h in hasnt)
    rep["section_01_alone"] = only(3, ["THREE MARKETABILITY TESTS", "01", tests[0]],
                                   [tests[1], tests[2]])
    rep["section_02_alone"] = only(5, ["THREE MARKETABILITY TESTS", "02", tests[1]],
                                   [tests[0], tests[2]])
    rep["section_03_alone"] = only(7, ["THREE MARKETABILITY TESTS", "03", tests[2]],
                                   [tests[0], tests[1]])
    rep["slides_with_all_three"] = [i + 1 for i, t in enumerate(vis)
                                    if all(x in t for x in tests)]
    rep["recap_only"] = rep["slides_with_all_three"] == [9]
    rep["section_slides_have_no_builds"] = all(S.STEPS[n] == 1 for n in (3, 5, 7, 9))
    rep["fieldkit_url_slides"] = [i + 1 for i, t in enumerate(vis)
                                  if "temidayoafonja.com/fieldkit" in t]
    rep["watch_next"] = "Before You Quit Your Job, Check These 3 Things" in vis[12]
    rep["evidence_kept_generic"] = ("I own the QBR process for this business unit."
                                    in vis[3])
    banned = ["—", "–", "résumé", "Maven", "Keep the Proof",
              "Career Portability Map", "Role Relevance", "$150"]
    rep["banned_found"] = {b: [i + 1 for i, t in enumerate(vis + nts) if b in t]
                           for b in banned}
    rep["banned_found"] = {k: v for k, v in rep["banned_found"].items() if v}
    import re
    rep["timings"] = re.findall(r"Timing: (\d+:\d\d)-(\d+:\d\d)", "\n".join(nts))

    rep["geometry_main"] = geometry(measure(html))
    rep["geometry_recording"] = geometry(measure(rhtml), RECORDING_PLAN)

    growth, matches, by_slide = [], [], {}
    for i, (n, s) in enumerate(RECORDING_PLAN):
        by_slide.setdefault(n, []).append((s, rcvs[i], rpngs[i]))
    for n, frames in by_slide.items():
        for k in range(1, len(frames)):
            removed = set(texts(frames[k - 1][1])) - set(texts(frames[k][1]))
            growth.append({"slide": n, "step": frames[k][0],
                           "removed": sorted(removed)})
        matches.append({"slide": n,
                        "content_identical_to_main":
                            fingerprint(frames[-1][1]) == fingerprint(cvs[n - 1]),
                        "pixel_delta": pixel_delta(frames[-1][2], pngs[n - 1])})
    rep["reveal_growth"] = growth
    rep["last_build_matches_main"] = matches

    with open(os.path.join(OUT, "qa-raw-v1.0.json"), "w") as f:
        json.dump(rep, f, indent=1)

    print("main slides %d | pdf pages %d %s | recording slides %d"
          % (rep["main_slides"], rep["pdf_pages"], rep["pdf_size_in"],
             rep["recording_slides"]))
    print("section breaks alone:", rep["section_01_alone"], rep["section_02_alone"],
          rep["section_03_alone"], "| all three on:", rep["slides_with_all_three"],
          "| recap only:", rep["recap_only"])
    print("section/recap unbuilt:", rep["section_slides_have_no_builds"],
          "| fieldkit url on:", rep["fieldkit_url_slides"],
          "| watch next:", rep["watch_next"])
    for k, v in rep["geometry_main"].items():
        print("  main %-12s %d %s" % (k, len(v), v[:2]))
    for k, v in rep["geometry_recording"].items():
        print("  rec  %-12s %d %s" % (k, len(v), v[:2]))
    print("frames removing content:", sum(1 for g in growth if g["removed"]))
    print("last build == main:", all(m["content_identical_to_main"] for m in matches))
    print("banned:", rep["banned_found"])
    print("timings:", ["%s-%s" % t for t in rep["timings"]])


if __name__ == "__main__":
    main()
