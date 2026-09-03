# -*- coding: utf-8 -*-
"""Run the v5.0 belonging/identity QA plus package, deck and rendering checks
across Videos 1-5."""
import sys, os, json, glob, zipfile, subprocess, re
sys.path.insert(0,"/tmp/v5"); sys.path.insert(0,"/tmp/da")
from qa5 import audit
from docxkit import sha256
from pptx import Presentation
from docx import Document

R="/home/user/temidayoafonja-site/deliverables"
DECK={1:(R+"/video-1-slides/out/Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pptx",
         R+"/video-1-slides/out/Video-1-Reveal-Builds_v2.4.pptx",
         R+"/video-1-slides/out/Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pdf"),
      2:(R+"/video-2-slides/out/Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.1.pptx",
         R+"/video-2-slides/out/Video-2-Reveal-Builds_v1.1.pptx",
         R+"/video-2-slides/out/Video-2-Is-Your-Job-Making-You-Less-Marketable_v1.1.pdf"),
      3:(R+"/video-3-slides/out/Video-3-3-Things-to-Do-Before-Quitting-Your-Job_v1.1.pptx",
         R+"/video-3-slides/out/Video-3-Reveal-Builds_v1.1.pptx",
         R+"/video-3-slides/out/Video-3-3-Things-to-Do-Before-Quitting-Your-Job_v1.1.pdf"),
      4:(R+"/video-4-slides/out/Video_4_Main_Slides.pptx",
         R+"/video-4-slides/out/Video_4_Reveal_Builds.pptx",
         R+"/video-4-slides/out/Video_4_Slide_Preview.pdf"),
      5:(R+"/video-5-slides/out/Video_5_Main_Slides.pptx",
         R+"/video-5-slides/out/Video_5_Reveal_Builds.pptx",
         R+"/video-5-slides/out/Video_5_Slide_Preview.pdf")}
CFG=json.load(open("/tmp/v5/qacfg.json"))
OUT={}
for n in range(1,6):
    B="/tmp/v5/v%d"%n if n!=1 else "/tmp/v5/v1"
    cfg=CFG[str(n)]
    Rq,F,words,mk,spoken=audit(n,os.path.join(B,"canonical_v5.0.txt"),cfg)
    def chk(k,l,c,d=""):
        Rq[k]={"check":l,"pass":bool(c),"detail":str(d)}
        if not c: F.append((k,l,str(d)))
    ROOT=os.path.join(B,"Video_%d_HIT_FINAL"%n)
    SH=os.path.join(ROOT,"SHORTS")
    zp=os.path.join(B,"Video_%d_HIT_FINAL_Recording_and_Shorts_Package.zip"%n)
    desc=os.path.join(B,"Video_%d_YouTube_Description_HIT.docx"%n)
    def text(p): return "\n".join(x.text for x in Document(p).paragraphs)
    pub=glob.glob(ROOT+"/LONG_FORM/*Publishing*v5.0.docx")[0]
    edb=glob.glob(ROOT+"/LONG_FORM/*EDITOR_ONLY*v5.0.docx")[0]
    seb=glob.glob(SH+"/*Shorts_EDITOR_ONLY*v5.0.docx")[0]
    pubt=text(pub); desct=text(desc); edbt=text(edb); sebt=text(seb)
    shorts={f:text(os.path.join(SH,f)) for f in sorted(os.listdir(SH))
            if f.startswith("Video_%d_Short_"%n)}
    # Shorts
    chk("SH1","Four Shorts, all rewritten for v5.0", len(shorts)==4, len(shorts))
    chk("SH2","Every Short speaks to one viewer",
        all(re.search(r"\b(you|your)\b",t,re.I) for t in shorts.values()))
    ED=["Visual:","On-screen","B-roll","EDITOR","End on:","FACTUAL BOUNDARY"]
    chk("SH3","Shorts recording docs are spoken copy only",
        not [(f,w) for f,t in shorts.items() for w in ED if w in t],
        [(f,w) for f,t in shorts.items() for w in ED if w in t])
    chk("SH4","Shorts brief gives each Short a hook and a route",
        sebt.count("On-screen hook")>=4 and sebt.count("Related Video")>=4)
    # editor brief 14 sections
    chk("EB1","Editor brief carries all fourteen sections",
        all(("%d.  "%i) in edbt for i in range(1,15)),
        [i for i in range(1,15) if ("%d.  "%i) not in edbt])
    chk("EB2","Editor brief carries the identity promise and exit",
        "Identity promise" in edbt and "Identity exit" in edbt)
    # description
    chk("D1","Emoji system present in both docs",
        all(e in pubt for e in ["✨","🧭","⏱️","▶️","🔗"]) and
        all(e in desct for e in ["✨","🧭","⏱️","▶️","🔗"]))
    def block(t):
        a=t.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
        b=t.index("— END OF")
        return [x.strip() for x in t[a:b].split("\n") if x.strip()][1:]
    chk("D2","Public block identical in both docs", block(pubt)==block(desct))
    chk("D3","No [INSERT] placeholder", "[INSERT" not in (pubt+desct+edbt+sebt))
    chk("D4","Working estimates flagged outside public copy",
        "WORKING ESTIMATES — EDITOR MUST REPLACE FROM FINAL CUT" in pubt and
        pubt.index("— END OF")<pubt.index("WORKING ESTIMATES"))
    chk("D5","Tag field in both docs",
        "YouTube tag field" in pubt and "YouTube tag field" in desct)
    chk("D6","Pinned comment invites self-recognition",
        "Which part of this is happening in your situation" in pubt)
    # decks
    mp=Presentation(DECK[n][0]); rp=Presentation(DECK[n][1])
    chk("P1","Main slide count", len(mp.slides)==cfg["markers"], len(mp.slides))
    chk("P2","Reveal frame count", len(rp.slides)==cfg["frames"], len(rp.slides))
    prov=json.load(open(os.path.join(B,"_partdiff.json")))
    chk("P3","No slide XML, media or rels change",
        prov["main_nonnotes"]==[] and prov["reveal_nonnotes"]==[])
    notes=[s.notes_slide.notes_text_frame.text for s in mp.slides]+\
          [s.notes_slide.notes_text_frame.text for s in rp.slides]
    chk("P4","All notes parts carry v5.0 narration",
        all("Timing:" in x for x in notes) and
        not any(re.search(r"v[234]\.[01]",x) for x in notes), len(notes))
    # package
    names=zipfile.ZipFile(zp).namelist()
    chk("Z1","Exactly 13 files in ZIP", len(names)==13, len(names))
    chk("Z2","No source, Python, QA or render files in ZIP",
        not any(x.endswith((".py",".pyc",".html",".png")) or "_source" in x
                for x in names))
    sums=[l for l in open(os.path.join(ROOT,"SHA256SUMS.txt"),encoding="utf-8")
          .read().strip().split("\n") if l.strip() and not l.startswith("#")]
    chk("Z3","SHA256SUMS exactly 12 entries", len(sums)==12, len(sums))
    pr=subprocess.run("sha256sum -c SHA256SUMS.txt",shell=True,cwd=ROOT,
                      capture_output=True,text=True)
    chk("Z4","sha256sum -c successful", pr.returncode==0, pr.stdout[-140:])
    zs=sha256(zp)
    chk("Z5","Sibling checksum matches",
        open(zp+".sha256",encoding="utf-8").read().split()[0]==zs)
    OUT[n]={"checks":Rq,"failures":F,"words":words,"markers":mk,
            "slides":len(mp.slides),"frames":len(rp.slides),
            "notes":{"main":len(prov["main_changed"]),
                     "reveal":len(prov["reveal_changed"])},
            "hashes":{"zip":zs,"desc":sha256(desc),"main":sha256(DECK[n][0]),
                      "reveal":sha256(DECK[n][1]),"pdf":sha256(DECK[n][2])}}
    print("VIDEO %d: %d checks, %d failed"%(n,len(Rq),len(F)))
    for f in F: print("   FAIL",f)
json.dump(OUT,open("/tmp/v5/QA_REPORT_v5.json","w"),indent=1)
print("\nTOTAL %d checks, %d failed"%(sum(len(v["checks"]) for v in OUT.values()),
                                      sum(len(v["failures"]) for v in OUT.values())))
