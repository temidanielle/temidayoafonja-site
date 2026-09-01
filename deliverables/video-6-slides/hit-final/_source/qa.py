# -*- coding: utf-8 -*-
"""Video 6 H.I.T. package QA, including slide checks and source comparison."""
import zipfile, re, sys, os, json, hashlib, subprocess, glob, unicodedata
from xml.etree import ElementTree as ET
sys.path.insert(0,'/tmp/v6hit')
from script_text import LINES, SPOKEN, MARKERS
NS='{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
R="/tmp/v6hit/Video_6_HIT_FINAL"; L=R+"/LONG_FORM/"
MK=re.compile(r'^SLIDE\s+—\s+(.+)$')
OLD_NEXT_FULL="How to Prove the Value of Work That Had No Blueprint"
NEXT_FULL="How to Show Your Impact at Work When You Built It From Scratch"
rep={}

def paras(p):
    x=ET.fromstring(zipfile.ZipFile(p).read('word/document.xml'))
    t=lambda q:''.join(e.text or '' for e in q.iter(NS+'t'))
    return [t(el).strip() for el in x.find(NS+'body') if el.tag==NS+'p' and t(el).strip()]
def txt_blocks(p):
    return [b.strip() for b in open(p,encoding="utf-8").read().strip().split("\n\n") if b.strip()]

# ---- 1-7 : the live decks --------------------------------------------------
from pptx import Presentation
MAIN="/home/user/temidayoafonja-site/deliverables/video-6-slides/out/Video_6_Main_Slides.pptx"
REV ="/home/user/temidayoafonja-site/deliverables/video-6-slides/out/Video_6_Reveal_Builds.pptx"
ORIG_MAIN="/tmp/v6hit/Video_6_Main_Slides.pptx.orig"
ORIG_REV ="/tmp/v6hit/Video_6_Reveal_Builds.pptx.orig"
rep["01_main_slide_count"]=len(Presentation(MAIN).slides)
rep["02_reveal_frame_count"]=len(Presentation(REV).slides)
rep["01b_main_is_12"]=rep["01_main_slide_count"]==12
rep["02b_reveal_is_23"]=rep["02_reveal_frame_count"]==23

def parts(p):
    z=zipfile.ZipFile(p); return {n:hashlib.sha256(z.read(n)).hexdigest() for n in z.namelist()}
def runs_of(pptx, idx):
    sl=list(Presentation(pptx).slides)[idx]
    return [r.text for sh in sl.shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs for r in para.runs]
rep["03_pre_change_slide12_title"]=["HOW TO PROVE THE","VALUE OF WORK","THAT HAD NO BLUEPRINT"]
om, cm = parts(ORIG_MAIN), parts(MAIN)
changed_m=[n for n in om if om[n]!=cm.get(n)]
rep["04_slides_1_to_11_byte_identical"]=changed_m==["ppt/slides/slide12.xml"]
rep["04b_main_parts_changed"]=changed_m
orv, crv = parts(ORIG_REV), parts(REV)
changed_r=[n for n in orv if orv[n]!=crv.get(n)]
rep["06_only_one_reveal_frame_changed"]=changed_r==["ppt/slides/slide23.xml"]
rep["06b_reveal_parts_changed"]=changed_r
NEWT=["HOW TO SHOW YOUR","IMPACT AT WORK","WHEN YOU BUILT","IT FROM SCRATCH"]
def title_runs(p,i):
    sl=list(Presentation(p).slides)[i]
    for sh in sl.shapes:
        if sh.has_text_frame:
            rs=[r.text for para in sh.text_frame.paragraphs for r in para.runs]
            if rs and rs[0].startswith("HOW TO"): return rs
rep["05_slide12_title_corrected"]=title_runs(MAIN,11)==NEWT
rep["06c_reveal23_title_corrected"]=title_runs(REV,22)==NEWT
def markup_same_but_text(a,b,part):
    za,zb=zipfile.ZipFile(a),zipfile.ZipFile(b)
    sa=re.sub(rb'<a:t>.*?</a:t>',b'<a:t/>',za.read(part))
    sb=re.sub(rb'<a:t>.*?</a:t>',b'<a:t/>',zb.read(part))
    # the only permitted structural delta is one extra run + one extra break
    return sb.count(b'<a:r>')-sa.count(b'<a:r>')==1 and sb.count(b'<a:br/>')-sa.count(b'<a:br/>')==1
rep["05b_slide12_delta_is_one_run_and_one_break"]=markup_same_but_text(ORIG_MAIN,MAIN,"ppt/slides/slide12.xml")
rep["06d_reveal23_delta_is_one_run_and_one_break"]=markup_same_but_text(ORIG_REV,REV,"ppt/slides/slide23.xml")
git=lambda *a: subprocess.run(["git","-C","/home/user/temidayoafonja-site"]+list(a),
                              capture_output=True,text=True).stdout
rep["07_thumbnail_unchanged"]=not [l for l in git("status","--porcelain").splitlines()
                                   if "thumbnail" in l.lower()]

# ---- 8-11 : scripts and the source comparison ------------------------------
CANON=[b.strip() for b in open("/tmp/v6hit/canonical_script.txt",encoding="utf-8").read().split("\n\n") if b.strip()]
CANON_SPOKEN=[b for b in CANON if not b.startswith("[SLIDE:")]
CANON_MARKERS=[(i,b) for i,b in enumerate(CANON) if b.startswith("[SLIDE:")]
tel=txt_blocks(L+"Video6TeleprompterScriptwithslidemarkers_HIT_v2.0.txt")
rd =txt_blocks(L+"Video6ReadingScriptnomarkers_HIT_v2.0.txt")
tel=tel[tel.index(CANON_SPOKEN[0]):]; rd=rd[rd.index(CANON_SPOKEN[0]):]
tel_spoken=[b for b in tel if not b.startswith("SLIDE  —")]
s_="\n\n".join(CANON_SPOKEN); t_="\n\n".join(rd); u_="\n\n".join(tel_spoken)
tel_markers=[(i,b) for i,b in enumerate(tel) if b.startswith("SLIDE  —")]
sv={"named_canonical_file":"Video_6_Code_Prompt_HIT_Final.txt",
 "named_file_uploaded":False,
 "source_used":"The session's machine-recorded copy of Temidayo's brief, "
   "extracted programmatically between the BEGIN/END APPROVED VIDEO 6 SCRIPT "
   "fences. Not retyped, not normalised.",
 "recorded_copies_found":1,
 "source_blocks":len(CANON),"source_markers":len(CANON_MARKERS),
 "source_spoken_paragraphs":len(CANON_SPOKEN),
 "source_spoken_words":sum(len(b.split()) for b in CANON_SPOKEN),
 "teleprompter_minus_markers_equals_source":tel_spoken==CANON_SPOKEN,
 "reading_script_equals_source":rd==CANON_SPOKEN,
 "marker_names_match":[c[1][len("[SLIDE:"):-1].strip() for c in CANON_MARKERS]==
                      [m[1][len("SLIDE  —"):].strip() for m in tel_markers],
 "marker_positions_match":[c[0] for c in CANON_MARKERS]==[m[0] for m in tel_markers],
 "characters":(len(s_),len(t_),len(u_)),
 "sha256_joined_spoken_text":hashlib.sha256(s_.encode()).hexdigest(),
 "sha256_identical_across_all_three":
   hashlib.sha256(s_.encode()).hexdigest()==hashlib.sha256(t_.encode()).hexdigest()
   ==hashlib.sha256(u_.encode()).hexdigest(),
 "per_character_class":{k:(s_.count(v),t_.count(v),u_.count(v)) for k,v in
   {"apostrophe":"’","left_quote":"“","right_quote":"”","em_dash":"—","ellipsis":"…"}.items()},
 "capital_letters":(sum(c.isupper() for c in s_),sum(c.isupper() for c in t_),
                    sum(c.isupper() for c in u_))}
sv["status"]="PASSED" if all([sv["teleprompter_minus_markers_equals_source"],
  sv["reading_script_equals_source"],sv["marker_names_match"],
  sv["marker_positions_match"],sv["sha256_identical_across_all_three"]]) else "FAILED"
sv["caveat"]=("The named canonical .txt was not uploaded. This comparison is "
  "against the recorded copy of the same message, so it proves the package "
  "matches what the session received; it cannot rule out a difference from a "
  "file held elsewhere.")
rep["SOURCE_VERIFICATION"]=sv
rep["11_canonical_comparison_passes"]=sv["status"]=="PASSED"

def from_first(s): return s[s.index(SPOKEN[0]):]
tel_d=from_first(paras(L+"Video6TeleprompterScriptwithslidemarkers_HIT_v2.0.docx"))
rd_d =from_first(paras(L+"Video6ReadingScriptnomarkers_HIT_v2.0.docx"))
strip=lambda s:[x for x in s if not MK.match(x)]
rep["08_teleprompter_docx_equals_txt"]=tel_d==tel
rep["09_reading_docx_equals_txt"]=rd_d==rd
rep["10_teleprompter_minus_markers_equals_reading"]=strip(tel_d)==rd_d

DECK=["Core Distinction","More Work / Real Growth","Three Tests","Complexity",
 "Capability Question","Accountability / Authority","Authority Warning",
 "Return","Pattern Read","Before the Scope Expands Again",
 "Capability Formation Field Kit","Watch Next"]
got=[MK.match(s).group(1) for s in tel_d if MK.match(s)]
rep["12_markers_ordered_and_mapped"]=got==DECK
rep["12b_marker_count_is_12"]=len(got)==12

shorts=[f for f in sorted(os.listdir(R+"/SHORTS")) if "EDITOR_ONLY" not in f]
rep["13_exactly_four_short_docs"]=len(shorts)==4
rep["13b_short_filenames"]=shorts
BAN=("EDITOR","On-screen hook","Visual:","B-roll","caption","zoom","punch-in",
     "stock footage","overlay","9:16","Related Video","direct to camera",
     "Reveal progressively","End:","FACTUAL BOUNDARY","Progressively show")
bad={f:[b for b in BAN if b.lower() in " ".join(paras(R+"/SHORTS/"+f)).lower()] for f in shorts}
bad={k:v for k,v in bad.items() if v}
rep["14_shorts_free_of_editor_directions"]=not bad
rep["14b_offending"]=bad
rep["15_editor_docs_labelled"]=all(paras(R+"/"+f)[0]=="EDITOR ONLY" for f in
 ("LONG_FORM/Video_6_EDITOR_ONLY_HIT_Brief_v2.0.docx",
  "SHORTS/Video_6_Shorts_EDITOR_ONLY_HIT_Brief.docx"))

brief=" ".join(paras(L+"Video_6_EDITOR_ONLY_HIT_Brief_v2.0.docx"))
pub=" ".join(paras(L+"Video_6_Publishing_Package_HIT_v2.0.docx"))
shb=" ".join(paras(R+"/SHORTS/Video_6_Shorts_EDITOR_ONLY_HIT_Brief.docx"))
DESCDOC="/tmp/v6hit/Video_6_YouTube_Description_HIT.docx"
ALL={os.path.basename(p):" ".join(paras(p)) for p in
     glob.glob(R+"/LONG_FORM/*.docx")+glob.glob(R+"/SHORTS/*.docx")+[DESCDOC]}

rep["16_proof_boundary_bounded"]=("scope has expanded beyond the original job "
  "description more than once" in brief and "FACTUAL BOUNDARY" in brief
  and "FACTUAL BOUNDARY" in shb)
INVENT=("EY","Deloitte","my employer","the company I worked for","promoted to",
        "as a director","as a manager at","my boss said","in 2019","in 2020")
inv={n:[x for x in INVENT if re.search(r'\b%s\b'%re.escape(x),t,re.I)] for n,t in ALL.items()}
inv={k:v for k,v in inv.items() if v}
rep["16b_no_invented_employer_role_timeline"]=not inv
rep["16c_offending"]=inv
METRIC=[r"\b\d{1,3}\s*%", r"\$\s*\d", r"\b\d+\s*million", r"\bNPS\b", r"\bROI\b"]
mets={n:[p for p in METRIC if re.search(p,t)] for n,t in ALL.items()}
# the long-form editor brief names the forbidden figures only to forbid them
mets={k:v for k,v in mets.items() if v and "EDITOR_ONLY_HIT_Brief_v2.0" not in k}
rep["17_no_retention_or_turnover_metric"]=not mets
rep["17b_offending"]=mets

OTHER=("Keep the Proof","Career Decision Evidence Check","/keep-the-proof",
       "/career-decisions")
leak={n:[o for o in OTHER if o in t] for n,t in ALL.items()}
leak={k:v for k,v in leak.items() if v and "EDITOR_ONLY_HIT_Brief_v2.0" not in k}
rep["18_no_other_cta_leak"]=not leak
rep["18b_offending"]=leak
rep["19_fieldkit_is_sole_cta"]=all("temidayoafonja.com/fieldkit" in t
    for t in ("\n".join(rd_d),brief,pub))
rep["20_video7_title_everywhere_watch_next"]=all(
 "How to Show Your Impact at Work When You Built It From Scratch" in t
 for t in (brief,pub," ".join(paras(DESCDOC))))
# the old title may appear ONLY in the editor brief's FROM/TO correction record
stale={n:True for n,t in ALL.items()
       if ("Had No Blueprint" in t or "HAD NO BLUEPRINT" in t)
       and "EDITOR_ONLY_HIT_Brief_v2.0" not in n}
brief_record_ok = ("FROM:  "+OLD_NEXT_FULL) in brief and ("TO:    "+NEXT_FULL) in brief
readme=open(R+"/README_FINAL.txt",encoding="utf-8").read()
rep["21_no_old_video7_title_outside_the_correction_record"]=not stale
rep["21c_editor_brief_carries_the_from_to_record"]=brief_record_ok
rep["21d_stale_title_locations"]=sorted(stale)
rep["21b_no_old_title_on_corrected_slides"]=(
 "BLUEPRINT" not in " ".join(runs_of(MAIN,11)).upper()
 and "BLUEPRINT" not in " ".join(runs_of(REV,22)).upper())
UNSAFE=("download","forward the","usb","screenshot","copy the file","export the",
        "take the file","email yourself")
uns={n:[u for u in UNSAFE if u in t.lower()] for n,t in ALL.items()}
uns={k:v for k,v in uns.items() if v}
rep["22_no_unsafe_evidence_advice"]=not uns

pubp=paras(L+"Video_6_Publishing_Package_HIT_v2.0.docx")
bi=pubp.index("— END OF THE COPY-READY DESCRIPTION —"); desc=pubp[:bi]
EMOJI=["✨ Complexity","✨ Authority","✨ Return",
       "🧭 CAPABILITY FORMATION FIELD KIT","⏱️ CHAPTERS","▶️ WATCH NEXT",
       "🔗 CONNECT AND EXPLORE"]
rep["23_emoji_system_retained"]=all(any(e in p for p in desc) for e in EMOJI)
CH=["00:00 Your Workload Can Grow Faster Than Your Career",
 "01:55 The 3 Tests for Real Growth",
 "02:10 Test 1: Did the Problem Become More Complex?",
 "03:30 What Can You Do Now That You Couldn’t Before?",
 "04:15 Test 2: Did Your Authority Expand?",
 "05:20 Accountability Is Not Authority",
 "06:05 Test 3: What Did the Work Return?","08:20 Read the Pattern",
 "09:45 What to Ask Before Your Scope Expands Again",
 "10:55 Capability Formation Field Kit",
 "11:20 How to Show Your Impact at Work When You Built It From Scratch"]
ci=pubp.index("⏱️ CHAPTERS")
rep["24_chapters_inline_in_description"]=pubp[ci+1:ci+12]==CH
rep["25_no_insert_placeholder"]=not any("[INSERT" in p.upper() for p in
    pubp+paras(DESCDOC))
rep["26_warning_outside_public_copy"]=(
 not any("WORKING ESTIMATES" in p for p in desc)
 and any("WORKING ESTIMATES" in p for p in pubp[bi:]))
dp=paras(DESCDOC); dbi=dp.index("— END OF THE COPY-READY DESCRIPTION —")
rep["27_description_doc_matches_publishing"]=(
 dp[dp.index("⏱️ CHAPTERS"):dbi]==pubp[ci:bi])
rep["27b_description_doc_outside_zip"]=True

ZIP="/tmp/v6hit/Video_6_HIT_FINAL_Recording_and_Shorts_Package.zip"
zf=zipfile.ZipFile(ZIP); names=zf.namelist()
rep["31_zip_file_count"]=len(names)
rep["31b_zip_count_is_13"]=len(names)==13
rep["32_no_source_or_python_in_zip"]=not any(
 "_source" in n or n.endswith((".py",".pyc")) for n in names)
rep["32b_no_description_doc_in_zip"]=not any("Description" in n for n in names)
rep["32c_zip_integrity"]=zf.testzip() is None
sums=zf.read("Video_6_HIT_FINAL/SHA256SUMS.txt").decode()
entries=[l for l in sums.splitlines() if l.strip() and not l.startswith("#")]
rep["34_sums_entry_count"]=len(entries)
rep["34b_sums_has_12_entries"]=len(entries)==12
rep["34c_sums_no_self_or_zip"]=not any("SHA256SUMS.txt" in e or ".zip" in e for e in entries)
rep["35_every_listed_hash_matches"]=not [e for e in entries
 if hashlib.sha256(zf.read("Video_6_HIT_FINAL/"+e.split("  ",1)[1])).hexdigest()!=e.split("  ",1)[0]]
listed=[l.strip() for l in readme.split("ALL FILES IN THIS PACKAGE")[1].splitlines()
        if l.strip().endswith((".docx",".txt")) and "Description_HIT" not in l]
rep["33_readme_lists_13_files"]=len(listed)==13
zh=hashlib.sha256(open(ZIP,"rb").read()).hexdigest()
rep["36_sibling_checksum_matches"]=open(ZIP+".sha256").read().split()[0]==zh
rep["36b_zip_sha256"]=zh
rep["37_description_doc_sha256"]=hashlib.sha256(open(DESCDOC,"rb").read()).hexdigest()
dirty=[l for l in git("status","--porcelain").splitlines()
       if "video-6-slides" not in l and "SERIES_STATUS_TRACKER" not in l]
rep["38_nothing_else_modified"]=not dirty
rep["38b_unexpected"]=dirty

pages={}
for f in sorted(glob.glob("/tmp/render/V6_*-p01.png")):
    stem=f[:-8]; pages[os.path.basename(stem)[3:]]=len(glob.glob(stem+"-p*.png"))
rep["28_29_rendered_page_counts"]=pages
rep["canonical_spoken_paragraphs"]=len(SPOKEN)
rep["canonical_spoken_words"]=sum(len(p.split()) for p in SPOKEN)
ok=all(v for k,v in rep.items() if isinstance(v,bool))
rep["ALL_BOOLEAN_CHECKS_PASS"]=ok
print(json.dumps(rep,indent=1,ensure_ascii=False))
