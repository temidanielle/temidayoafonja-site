"""Build version 2 of the deck and run the slide QA pass.

Writes only *_v2 outputs. The version 1 files are never touched.

    python3 build/build_v2.py
"""
import os, sys, json, shutil, re

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck import W, H, SAFE, SAFE_CLEAR, ENDCARD_CLEAR, Canvas, render_pptx, render_html
from build import shoot, font, RENDER
import slides_v2 as S
from PIL import Image, ImageDraw

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "out")
V2 = os.path.join(OUT, "v2")
PNG = os.path.join(V2, "png")
REV = os.path.join(V2, "reveals")
GUIDES = os.path.join(V2, "guides")
DECK = "Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2"
BUILDS = "Video-1-Reveal-Builds_v2"
N = 10


def canvases(plan):
    out = []
    for n, step in plan:
        cv = Canvas(n, step)
        S.BUILDERS[n](cv, step)
        out.append(cv)
    return out


def sheet(paths, out_path, tw, title, subtitle, cols=5, labels=True):
    th = int(tw * H / W)
    rows = (len(paths) + cols - 1) // cols
    gx, gy, pad, top = 40, (86 if labels else 46), 56, 168
    cw = pad * 2 + cols * tw + (cols - 1) * gx
    ch = top + rows * (th + gy) + pad - 24
    im = Image.new("RGB", (cw, ch), (245, 241, 232))
    d = ImageDraw.Draw(im)
    d.text((pad, 56), title, font=font(34, True), fill=(15, 35, 70))
    d.text((pad, 104), subtitle, font=font(23), fill=(90, 107, 130))
    d.rectangle([pad, 146, cw - pad, 148], fill=(201, 168, 76))
    for i, p in enumerate(paths):
        r, c = divmod(i, cols)
        x, y = pad + c * (tw + gx), top + r * (th + gy)
        im.paste(Image.open(p).resize((tw, th), Image.LANCZOS), (x, y))
        d.rectangle([x, y, x + tw - 1, y + th - 1], outline=(224, 217, 200))
        if labels:
            d.text((x, y + th + 14), "%02d  %s" % (i + 1, S.TITLES[i + 1]),
                   font=font(18, True), fill=(15, 35, 70))
    im.save(out_path)
    return out_path


def guide_overlay(paths, outdir):
    os.makedirs(outdir, exist_ok=True)
    for i, p in enumerate(paths, start=1):
        im = Image.open(p).convert("RGB")
        d = ImageDraw.Draw(im, "RGBA")
        d.rectangle([0, 0, SAFE_CLEAR["w"], SAFE_CLEAR["h"]],
                    fill=(193, 68, 14, 38), outline=(193, 68, 14, 200), width=3)
        d.rectangle([SAFE["x"], SAFE["y"], SAFE["x"] + SAFE["w"],
                     SAFE["y"] + SAFE["h"]], outline=(193, 68, 14, 255), width=4)
        d.text((SAFE["x"] + 18, SAFE["y"] + 18),
               "CAMERA\n480 x 270 px\n25% of frame width", font=font(24, True),
               fill=(193, 68, 14))
        if i == N:
            z = ENDCARD_CLEAR
            d.rectangle([z["x"], z["y"], z["x"] + z["w"], z["y"] + z["h"]],
                        fill=(201, 168, 76, 45), outline=(201, 168, 76, 230), width=4)
            d.text((z["x"] + 20, z["y"] + 20), "YOUTUBE END-SCREEN ELEMENT",
                   font=font(24, True), fill=(160, 128, 40))
        im.save(os.path.join(outdir, "guide-%02d.png" % i))


# ----------------------------------------------------------------------- QA
MEASURE_JS = """
() => {
  const out = [];
  document.querySelectorAll('.slide').forEach((slide, si) => {
    const sb = slide.getBoundingClientRect();
    const items = [];
    slide.querySelectorAll('p').forEach(p => {
      const r = p.getBoundingClientRect();
      const range = document.createRange();
      range.selectNodeContents(p);
      const tops = new Set();
      for (const cr of range.getClientRects()) tops.add(Math.round(cr.top));
      const wanted = (p.innerHTML.match(/<br>/g) || []).length + 1;
      items.push({
        text: p.textContent.slice(0, 60),
        x: r.left - sb.left, y: r.top - sb.top, w: r.width, h: r.height,
        linesWanted: wanted, linesActual: tops.size
      });
    });
    out.push({ slide: si + 1, items });
  });
  return out;
}
"""


def measure(html_path):
    from playwright.sync_api import sync_playwright
    with sync_playwright() as pw:
        b = pw.chromium.launch(
            executable_path="/opt/pw-browsers/chromium-1194/chrome-linux/chrome",
            args=["--no-sandbox"])
        pg = b.new_page(viewport={"width": W, "height": H})
        pg.goto("file://" + os.path.abspath(html_path), wait_until="load")
        pg.wait_for_timeout(600)
        data = pg.evaluate(MEASURE_JS)
        b.close()
    return data


def overlaps(a, b, pad=0):
    return not (a["x"] + a["w"] <= b["x"] + pad or b["x"] + b["w"] <= a["x"] + pad
                or a["y"] + a["h"] <= b["y"] + pad or b["y"] + b["h"] <= a["y"] + pad)


def qa(measured, pptx_path, pdf_path):
    from pptx import Presentation
    import pymupdf

    rep = {"clipped": [], "wrapped": [], "crowded": [], "safe_area": [],
           "endscreen": [], "bottom_edge": []}
    for s in measured:
        n, items = s["slide"], s["items"]
        for it in items:
            if (it["x"] < -1 or it["y"] < -1 or it["x"] + it["w"] > W + 1
                    or it["y"] + it["h"] > H + 1):
                rep["clipped"].append((n, it["text"], round(it["x"]), round(it["y"]),
                                       round(it["w"]), round(it["h"])))
            if it["linesActual"] > it["linesWanted"]:
                rep["wrapped"].append((n, it["text"], it["linesWanted"],
                                       it["linesActual"]))
            if it["y"] + it["h"] > H - 60:
                rep["bottom_edge"].append((n, it["text"], round(it["y"] + it["h"])))
            z = SAFE_CLEAR
            if overlaps(it, {"x": z["x"], "y": z["y"], "w": z["w"], "h": z["h"]}):
                rep["safe_area"].append((n, it["text"]))
            if n == N:
                e = ENDCARD_CLEAR
                if overlaps(it, {"x": e["x"], "y": e["y"], "w": e["w"], "h": e["h"]}):
                    rep["endscreen"].append((n, it["text"]))
        for i in range(len(items)):
            for j in range(i + 1, len(items)):
                if overlaps(items[i], items[j], pad=2):
                    rep["crowded"].append((n, items[i]["text"][:34],
                                           items[j]["text"][:34]))

    prs = Presentation(pptx_path)
    visible, notes_txt = [], []
    for sl in prs.slides:
        v = []
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                v.append(sh.text_frame.text)
        visible.append("\n".join(v))
        notes_txt.append(sl.notes_slide.notes_text_frame.text)
    alltext = "\n".join(visible + notes_txt)

    banned = ["Maven", "Keep the Proof", "Career Portability Map",
              "Role Relevance", "Portability Map", "$150", "Gem", "gem:",
              "résumé", "Résumé", "—", "–"]
    found = {b: [i + 1 for i, t in enumerate(visible) if b in t] +
                [("notes", i + 1) for i, t in enumerate(notes_txt) if b in t]
             for b in banned}
    found = {k: v for k, v in found.items() if v}

    url_slides = [i + 1 for i, t in enumerate(visible) if "temidayoafonja.com/book" in t]
    fieldkit_slides = [i + 1 for i, t in enumerate(visible)
                       if "Field Kit" in t or "Capability Formation Field Kit" in t]

    doc = pymupdf.open(pdf_path)
    times = re.findall(r"Timing: (\d+):(\d\d)-(\d+):(\d\d)", "\n".join(notes_txt))

    return {
        "report": rep,
        "slide_count": len(prs.slides),
        "pdf_pages": doc.page_count,
        "pdf_size_in": [round(doc[0].rect.width / 72, 3), round(doc[0].rect.height / 72, 3)],
        "banned_found": found,
        "url_slides": url_slides,
        "fieldkit_mention_slides": fieldkit_slides,
        "timings": times,
        "visible": visible,
        "notes": notes_txt,
    }


def main():
    for d in (V2, PNG, REV, GUIDES, RENDER):
        os.makedirs(d, exist_ok=True)
    for d in (PNG, REV, GUIDES):
        shutil.rmtree(d); os.makedirs(d)

    plan = [(n, S.STEPS[n]) for n in range(1, N + 1)]
    cvs = canvases(plan)
    pptx_path = os.path.join(OUT, DECK + ".pptx")
    pdf_path = os.path.join(OUT, DECK + ".pdf")
    render_pptx(cvs, pptx_path)
    html = render_html(cvs, os.path.join(RENDER, "deck_v2.html"),
                       "How I Changed Jobs Without Starting My Career Over (v2)")
    names = ["slide-%02d-%s.png" % (n, S.SLUGS[n]) for n in range(1, N + 1)]
    pngs = shoot(html, PNG, names, pdf_path=pdf_path)
    print("v2 deck: %d slides" % len(pngs))

    rev_plan, rev_names = [], []
    for n in range(1, N + 1):
        if S.STEPS[n] > 1:
            for s in range(1, S.STEPS[n] + 1):
                rev_plan.append((n, s))
                rev_names.append("slide-%02d-build-%d.png" % (n, s))
    rcvs = canvases(rev_plan)
    render_pptx(rcvs, os.path.join(V2, BUILDS + ".pptx"))
    rhtml = render_html(rcvs, os.path.join(RENDER, "builds_v2.html"), "Reveal builds v2")
    shoot(rhtml, REV, rev_names)
    print("v2 reveal builds: %d frames" % len(rev_names))

    sheet(pngs, os.path.join(V2, "contact-sheet-v2.png"), 560,
          "HOW I CHANGED JOBS WITHOUT STARTING MY CAREER OVER",
          "Version 2   ·   10 slides   ·   1920 x 1080   ·   16:9")
    sheet(pngs, os.path.join(V2, "phone-thumbnail-check-v2.png"), 320,
          "PHONE-THUMBNAIL LEGIBILITY CHECK (V2)",
          "Each slide at 320 x 180.", labels=False)
    guide_overlay(pngs, GUIDES)

    result = qa(measure(html), pptx_path, pdf_path)
    with open(os.path.join(V2, "qa-raw.json"), "w") as f:
        json.dump(result, f, indent=1)
    r = result["report"]
    print("slides=%d pdf_pages=%d size=%s" % (result["slide_count"],
                                              result["pdf_pages"], result["pdf_size_in"]))
    for k in ("clipped", "wrapped", "crowded", "safe_area", "endscreen", "bottom_edge"):
        print("%-12s %d %s" % (k, len(r[k]), r[k][:6]))
    print("banned:", result["banned_found"])
    print("url on slides:", result["url_slides"])
    print("timings:", result["timings"])


if __name__ == "__main__":
    main()
