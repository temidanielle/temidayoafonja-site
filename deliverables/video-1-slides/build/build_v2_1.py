"""Build version 2.1 and run the correction-pass QA.

Version 2.1 is a two-line correction of version 2 plus a linear
progressive-reveal recording deck. Nothing else changes.

Outputs (version 2 files are never touched):
  out/Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.1.pptx
  out/Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.1.pdf
  out/Video-1-Reveal-Builds_v2.1.pptx
  out/v2.1/...  supporting renders, sheets, guides and raw QA

    python3 build/build_v2_1.py
"""
import os, sys, json, shutil, hashlib

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck import W, H, SAFE, SAFE_CLEAR, ENDCARD_CLEAR, Canvas, render_pptx, render_html
from build import shoot, font, RENDER
from build_v2 import sheet as _sheet, guide_overlay as _guides, measure, overlaps
import slides_v2_1 as S
from PIL import Image, ImageDraw

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "out")
V = os.path.join(OUT, "v2.1")
PNG = os.path.join(V, "png")
REV = os.path.join(V, "reveals")
GUIDES = os.path.join(V, "guides")
DECK = "Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.1"
RECDECK = "Video-1-Reveal-Builds_v2.1"
N = 10

# The recording deck: every slide in running order, reveal slides expanded into
# sequential duplicate builds. No animations, so the reveals survive export and
# import into any presentation tool.
RECORDING_PLAN = [(n, s) for n in range(1, N + 1) for s in range(1, S.STEPS[n] + 1)]


def canvases(plan):
    out = []
    for n, step in plan:
        cv = Canvas(n, step)
        S.BUILDERS[n](cv, step)
        out.append(cv)
    return out


def texts(cv):
    return [p["text"] for el in cv.els if el["t"] == "text" for p in el["paras"]]


def sha(path):
    return hashlib.sha256(open(path, "rb").read()).hexdigest()


def fingerprint(cv):
    """Exact, comparable description of everything drawn on a canvas."""
    out = []
    for el in cv.els:
        e = dict(el)
        for k in ("fill", "line"):
            if k in e and e[k] is not None:
                e[k] = str(e[k])
        if e["t"] == "text":
            e["paras"] = [dict(p, color=str(p["color"])) for p in e["paras"]]
        out.append(e)
    return json.dumps(out, sort_keys=True)


def pixel_delta(a, b):
    """Antialiasing tolerance check between two renders of the same content."""
    from PIL import ImageChops
    d = ImageChops.difference(Image.open(a).convert("RGB"),
                              Image.open(b).convert("RGB")).convert("L")
    if d.getbbox() is None:
        return {"max": 0, "pixels": 0}
    return {"max": d.getextrema()[1],
            "pixels": sum(1 for v in d.getdata() if v > 0)}


def recording_sheet(paths, plan, out_path):
    cols, tw = 5, 420
    th = int(tw * H / W)
    rows = (len(paths) + cols - 1) // cols
    gx, gy, pad, top = 34, 74, 56, 172
    cw = pad * 2 + cols * tw + (cols - 1) * gx
    ch = top + rows * (th + gy) + pad - 20
    im = Image.new("RGB", (cw, ch), (245, 241, 232))
    d = ImageDraw.Draw(im)
    d.text((pad, 56), "RECORDING DECK, ADVANCE ORDER", font=font(34, True),
           fill=(15, 35, 70))
    d.text((pad, 104), "Video-1-Reveal-Builds_v2.1   ·   %d frames   ·   "
           "sequential duplicate slides, no animations" % len(paths),
           font=font(22), fill=(90, 107, 130))
    d.rectangle([pad, 148, cw - pad, 150], fill=(201, 168, 76))
    for i, p in enumerate(paths):
        n, step = plan[i]
        r, c = divmod(i, cols)
        x, y = pad + c * (tw + gx), top + r * (th + gy)
        im.paste(Image.open(p).resize((tw, th), Image.LANCZOS), (x, y))
        d.rectangle([x, y, x + tw - 1, y + th - 1], outline=(224, 217, 200))
        tag = "%02d  slide %d" % (i + 1, n)
        if S.STEPS[n] > 1:
            tag += ", build %d of %d" % (step, S.STEPS[n])
        d.text((x, y + th + 12), tag, font=font(18, True), fill=(15, 35, 70))
    im.save(out_path)
    return out_path


def main():
    for d in (V, PNG, REV, GUIDES, RENDER):
        os.makedirs(d, exist_ok=True)
    for d in (PNG, REV, GUIDES):
        shutil.rmtree(d); os.makedirs(d)

    # ---------------------------------------------------- main 10-slide deck
    main_plan = [(n, S.STEPS[n]) for n in range(1, N + 1)]
    cvs = canvases(main_plan)
    pptx_path = os.path.join(OUT, DECK + ".pptx")
    pdf_path = os.path.join(OUT, DECK + ".pdf")
    render_pptx(cvs, pptx_path)
    html = render_html(cvs, os.path.join(RENDER, "deck_v2_1.html"),
                       "How I Changed Jobs Without Starting My Career Over (v2.1)")
    names = ["slide-%02d-%s.png" % (n, S.SLUGS[n]) for n in range(1, N + 1)]
    pngs = shoot(html, PNG, names, pdf_path=pdf_path)

    # ------------------------------------------------ recording (build) deck
    rcvs = canvases(RECORDING_PLAN)
    rec_pptx = os.path.join(OUT, RECDECK + ".pptx")
    render_pptx(rcvs, rec_pptx)
    rhtml = render_html(rcvs, os.path.join(RENDER, "recording_v2_1.html"),
                        "Recording deck v2.1")
    rnames = ["frame-%02d-slide-%02d-build-%d.png" % (i + 1, n, s)
              for i, (n, s) in enumerate(RECORDING_PLAN)]
    rpngs = shoot(rhtml, REV, rnames,
                  pdf_path=os.path.join(V, RECDECK + ".pdf"))

    _sheet(pngs, os.path.join(V, "contact-sheet-v2.1.png"), 560,
           "HOW I CHANGED JOBS WITHOUT STARTING MY CAREER OVER",
           "Version 2.1   ·   10 slides   ·   1920 x 1080   ·   16:9")
    recording_sheet(rpngs, RECORDING_PLAN, os.path.join(V, "recording-deck-order-v2.1.png"))
    _guides(pngs, GUIDES)

    # ------------------------------------------------------------------- QA
    from pptx import Presentation
    import pymupdf

    rep = {}
    rep["main_slides"] = len(Presentation(pptx_path).slides)
    rep["recording_slides"] = len(Presentation(rec_pptx).slides)
    rep["pdf_pages"] = pymupdf.open(pdf_path).page_count
    rep["recording_plan"] = RECORDING_PLAN

    # corrections present, old wording gone
    prs = Presentation(pptx_path)
    vis = ["\n".join(sh.text_frame.text for sh in sl.shapes
                     if sh.has_text_frame and sh.text_frame.text.strip())
           for sl in prs.slides]
    nts = [sl.notes_slide.notes_text_frame.text for sl in prs.slides]
    flat = lambda t: " ".join(t.split())
    rep["slide2_corrected"] = ("I carried that way of seeing into cybersecurity and "
                               "later into people strategy." in flat(vis[1]))
    rep["slide2_old_gone"] = "cybersecurity workforce strategy" not in flat(vis[1])
    rep["slide6_corrected"] = ("after the onboarding redesign I led with my team."
                               in flat(nts[5]))
    rep["slide6_old_gone"] = "my team and I led" not in flat(nts[5])
    rep["slide6_visible_unchanged"] = flat(vis[5])

    # reveals never remove anything, and the last build equals the main slide
    growth, matches = [], []
    by_slide = {}
    for i, (n, s) in enumerate(RECORDING_PLAN):
        by_slide.setdefault(n, []).append((s, rcvs[i], rpngs[i]))
    for n, frames in by_slide.items():
        for k in range(1, len(frames)):
            prev, cur = set(texts(frames[k - 1][1])), set(texts(frames[k][1]))
            missing = prev - cur
            growth.append({"slide": n, "step": frames[k][0],
                           "removed": sorted(missing), "added": len(cur - prev)})
        same_content = fingerprint(frames[-1][1]) == fingerprint(cvs[n - 1])
        px = pixel_delta(frames[-1][2], pngs[n - 1])
        matches.append({"slide": n, "content_identical_to_main": same_content,
                        "byte_identical_render": sha(frames[-1][2]) == sha(pngs[n - 1]),
                        "pixel_delta": px})
    rep["reveal_growth"] = growth
    rep["last_build_matches_main"] = matches

    # geometry on every recording frame
    m = measure(rhtml)
    safe_hits, end_hits, clipped = [], [], []
    for fr in m:
        n, s = RECORDING_PLAN[fr["slide"] - 1]
        for it in fr["items"]:
            z = SAFE_CLEAR
            if overlaps(it, {"x": z["x"], "y": z["y"], "w": z["w"], "h": z["h"]}):
                safe_hits.append((n, s, it["text"]))
            if n == N:
                e = ENDCARD_CLEAR
                if overlaps(it, {"x": e["x"], "y": e["y"], "w": e["w"], "h": e["h"]}):
                    end_hits.append((n, s, it["text"]))
            if (it["x"] < -1 or it["y"] < -1 or it["x"] + it["w"] > W + 1
                    or it["y"] + it["h"] > H + 1):
                clipped.append((n, s, it["text"]))
    rep["recording_safe_area_hits"] = safe_hits
    rep["recording_endscreen_hits"] = end_hits
    rep["recording_clipped"] = clipped

    with open(os.path.join(V, "qa-raw-v2.1.json"), "w") as f:
        json.dump(rep, f, indent=1)

    print("main deck slides:      ", rep["main_slides"])
    print("main deck pdf pages:   ", rep["pdf_pages"])
    print("recording deck slides: ", rep["recording_slides"])
    print("slide 2 corrected:     ", rep["slide2_corrected"], "| old gone:", rep["slide2_old_gone"])
    print("slide 6 note corrected:", rep["slide6_corrected"], "| old gone:", rep["slide6_old_gone"])
    print("reveal frames that removed content:",
          sum(1 for g in growth if g["removed"]))
    print("last build content == main slide:",
          all(x["content_identical_to_main"] for x in matches))
    worst = max(matches, key=lambda x: x["pixel_delta"]["pixels"])
    print("  largest render difference: slide %d, %d px of %d, max channel delta %d"
          % (worst["slide"], worst["pixel_delta"]["pixels"], W * H,
             worst["pixel_delta"]["max"]))
    print("recording safe-area hits:", len(safe_hits),
          "| end-screen hits:", len(end_hits), "| clipped:", len(clipped))


if __name__ == "__main__":
    main()
