# -*- coding: utf-8 -*-
"""Apply the authorised Slide 12 CTA correction to both Video 1 decks.

Text only. Layout, box positions, palette and typography family are untouched.
Two type sizes come down because the new strings are materially longer than the
ones they replace and the boxes cannot widen without colliding with the
artifact images to the right. Every value was measured against the real
embedded fonts before applying:

  eyebrow   FREE CAREER EVIDENCE STARTER            3.24in in 8.33in  (13pt kept)
  headline  ONE ACCOMPLISHMENT ->                   5.46in in 6.11in  (32 -> 29pt)
            ONE PORTABLE PROOF LINE                 6.06in in 6.11in
            two lines end at 4.57in, clear of the body copy at 4.81in
  body      the approved Promise, wrapped to three lines exactly as the copy it
            replaces, 15pt kept
  URL       temidayoafonja.com/career-evidence-starter  4.51in in the 5.00in
            button (20 -> 14pt); at 20pt it would be 6.50in and burst the button

At 32pt no two-line split of the headline fits (the shortest overflowing line
is 6.66in), and a three-line block would collide with the body copy.
"""
import copy
from pptx import Presentation
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def q(t): return "{%s}%s" % (A, t)

EDITS = [
    ("THE CAPABILITY FORMATION FIELD KIT",
     ["FREE CAREER EVIDENCE STARTER"], None),
    ("Is your job still\x0bbuilding you?",
     ["ONE ACCOMPLISHMENT →", "ONE PORTABLE PROOF LINE"], 2900),
    ("Complete a private, evidence-led career position\x0bassessment using the "
     "last 90 days of your\x0bactual work.",
     ["Turn one accomplishment into proof you can use in a",
      "performance review, interview, internal move or",
      "career pivot."], None),
    ("temidayoafonja.com/fieldkit",
     ["temidayoafonja.com/career-evidence-starter"], 1400),
]
TARGETS = {"Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pptx": [12],
           "Video-1-Reveal-Builds_v2.4.pptx": [21]}

def set_lines(shape, lines, size=None):
    body = shape.text_frame._txBody
    ps = body.findall(q("p"))
    assert len(ps) == 1, "expected a single paragraph"
    p = ps[0]
    runs = p.findall(q("r"))
    assert runs, "expected a run to clone"
    template = runs[0]
    for child in list(p):
        if etree.QName(child).localname in ("r", "br"):
            p.remove(child)
    for i, line in enumerate(lines):
        if i:
            etree.SubElement(p, q("br"))
        r = copy.deepcopy(template)
        r.find(q("t")).text = line
        if size is not None:
            r.find(q("rPr")).set("sz", str(size))
        p.append(r)

total = 0
for path, slides in TARGETS.items():
    prs = Presentation(path)
    for n in slides:
        for sh in prs.slides[n - 1].shapes:
            if not sh.has_text_frame:
                continue
            for frm, to, size in EDITS:
                if sh.text_frame.text == frm:
                    set_lines(sh, to, size)
                    total += 1
                    print("  %-8s slide %-3d %r -> %r"
                          % (path[:8], n, frm[:34], to[0][:38]))
    prs.save(path)
assert total == 8, "expected 8 edits (4 per deck), got %d" % total
print("edits applied:", total)
