# -*- coding: utf-8 -*-
"""Replace the Field Kit page images on Video 1's CTA slide with the real
Career Evidence Starter artifact.

The Starter cover (PDF page 1) goes to the FRONT picture and the Portable
Proof Line page (PDF page 5) to the one BEHIND it, matching the approved
composition. Both renders are 1870x2420 -- byte-for-byte the same pixel
dimensions as the images they replace -- so every shape position, crop and
z-order is untouched.

The swap is done at the package level: only the two media parts change. Every
other part, including all other slide XML, is copied through byte-identically.
"""
import hashlib, shutil, zipfile
from pptx import Presentation

COVER = "/tmp/starter/starter_cover.png"        # front
PROOF = "/tmp/starter/starter_proofline.png"    # behind

FIELDKIT_COVER_SHA = "623b8e97866e"   # Read Your Position. Know Your Exposure.
FIELDKIT_PAGE_SHA  = "010204141a74"   # Part Two: Optionality

TARGETS = [("Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pptx", 12),
           ("Video-1-Reveal-Builds_v2.4.pptx", 21)]

def media_partnames(path, slide_no):
    """Map each picture on the slide to its media part name inside the package."""
    prs = Presentation(path)
    out = {}
    for sh in prs.slides[slide_no - 1].shapes:
        if sh.shape_type == 13:
            part = sh.part.related_part(sh._element.blip_rId)
            out[sh.image.sha1[:12]] = str(part.partname)
    return out

new = {}
for src, key in ((COVER, "cover"), (PROOF, "proof")):
    new[key] = open(src, "rb").read()
    print("%-6s %s  %d bytes" % (key, src.split("/")[-1], len(new[key])))

for path, slide_no in TARGETS:
    m = media_partnames(path, slide_no)
    assert FIELDKIT_COVER_SHA in m and FIELDKIT_PAGE_SHA in m, \
        "expected both Field Kit images on %s slide %d" % (path, slide_no)
    swap = {m[FIELDKIT_COVER_SHA].lstrip("/"): new["cover"],
            m[FIELDKIT_PAGE_SHA].lstrip("/"): new["proof"]}
    print("\n%s slide %d" % (path, slide_no))
    for part, blob in swap.items():
        print("   %-24s <- %d bytes" % (part, len(blob)))

    tmp = "/tmp/starter/_out.pptx"
    zin = zipfile.ZipFile(path)
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = swap.get(item.filename, zin.read(item.filename))
            zout.writestr(item, data)
    zin.close()
    shutil.move(tmp, path)
print("\nartifact swap complete")
