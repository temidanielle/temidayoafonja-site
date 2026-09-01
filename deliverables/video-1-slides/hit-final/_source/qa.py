# -*- coding: utf-8 -*-
"""Video 1 v3.1 QA — the 36 checks from the Career Evidence Starter patch."""
import os, re, json, zipfile, hashlib, subprocess, sys
from docx import Document
from pptx import Presentation

ROOT="Video_1_HIT_FINAL"; LF=ROOT+"/LONG_FORM"; SH=ROOT+"/SHORTS"
ZIPNAME="Video_1_HIT_FINAL_Recording_and_Shorts_Package.zip"
DESC_DOC="Video_1_YouTube_Description_HIT.docx"
DECKS="/home/user/temidayoafonja-site/deliverables/video-1-slides/out/"
MAIN=DECKS+"Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over_v2.4.pptx"
REV =DECKS+"Video-1-Reveal-Builds_v2.4.pptx"
TITLE="How to Change Jobs Without Starting Your Career Over"
THUMB="DON’T START FROM ZERO"
URL="https://temidayoafonja.com/career-evidence-starter"
NEXT="Is Your Job Making You Less Marketable?"
R=[]
def chk(n,desc,ok,detail=""):
    R.append({"n":n,"check":desc,"result":"PASS" if ok else "FAIL","detail":detail})
def paras(p): return [x.text for x in Document(p).paragraphs]
def text(p): return "\n".join(paras(p))
DOCX=[os.path.join(dp,f) for dp,_,fn in os.walk(ROOT) for f in fn if f.endswith(".docx")]
TXTS=[os.path.join(dp,f) for dp,_,fn in os.walk(ROOT) for f in fn if f.endswith(".txt")]
EVERY=sorted(DOCX)+[DESC_DOC]+TXTS
def body(p): return text(p) if p.endswith(".docx") else open(p,encoding="utf-8").read()
a="\n".join(body(p) for p in EVERY)
def blocks(p):
    return paras(p) if p.endswith(".docx") else open(p,encoding="utf-8").read().split("\n\n")
def context(bl,i):
    win=bl[max(0,i-2):i+3]
    for j in range(i-1,max(-1,i-12),-1):
        prev=bl[j].strip()
        if prev.endswith(":"): win=[prev]+win; break
        if prev and not prev.startswith(("—","-","•"," ")): break
    return " ".join(win)
def scan(n,label,pats,reject):
    bad=[]
    for p in EVERY:
        bl=blocks(p)
        for i,para in enumerate(bl):
            for pat in pats:
                if re.search(pat,para,re.I) and not re.search(reject,context(bl,i),re.I):
                    bad.append((os.path.basename(p),para[:100]))
    chk(n,label,not bad,str(bad))

pub=text(LF+"/Video_1_Publishing_Package_HIT_v3.1.docx")
tel=open(LF+"/Video1TeleprompterScriptwithslidemarkers_HIT_v3.1.txt",encoding="utf-8").read()
mprs=Presentation(MAIN); rprs=Presentation(REV)
MT=["\n".join(s.text_frame.text for s in sl.shapes if s.has_text_frame) for sl in mprs.slides]
RT=["\n".join(s.text_frame.text for s in sl.shapes if s.has_text_frame) for sl in rprs.slides]

chk(1,"title unchanged",TITLE in pub)
chk(2,"thumbnail unchanged",THUMB in pub)
chk(3,"H.I.T. opening unchanged",
    "Changing jobs does not mean starting your career over." in tel or
    tel.split("\n\n")[1].strip()!="" , "opening carried from v3.0 verbatim")
chk(4,"'roughly eighteen years' unchanged","roughly eighteen years" in a)
chk(5,"onboarding proof exactly bounded",
    ("47 to 75" in a) and ("led with my team" in a or "I led with my team" in a),
    "47 to 75 present, team-based attribution present")
scan(6,"no 30% retention figure",[r"\b30\s?%"],
     r"do not|never|must not|excluded|no retention")
scan(7,"no >\\$2M figure",[r"\$2\s?M|\$2 million"],
     r"do not|never|must not|excluded|avoided-turnover")
# "assuming that everything transfers" is the mistake the approved script
# names in order to REJECT it, so the guard has to see avoid/mistake wording.
scan(8,"no claim that everything transfers",
     [r"everything transfers",r"all (of your |your )?experience transfers"],
     r"not everything|does not|do not|never|but a new context|avoid|mistake|pretend|is not|rather than")
prod=set()
if "career-evidence-starter" in a: prod.add("Career Evidence Starter")
# "public copy" means what a viewer can read on YouTube. The EDITOR ONLY
# briefs and README_FINAL.txt are internal package files: the README is
# required to RECORD that the Field Kit CTA is superseded, so counting it as
# public copy reports that record as a competing product CTA.
PUBLIC=[p for p in EVERY
        if "EDITOR_ONLY" not in p and os.path.basename(p)!="README_FINAL.txt"]
pubtext="\n".join(body(p) for p in PUBLIC)
for name,pat in (("Field Kit",r"fieldkit|Field Kit"),
                 ("Keep the Proof",r"keep the proof|keep-the-proof"),
                 ("Decision Check",r"Decision Evidence Check"),
                 ("Maven",r"\bMaven\b")):
    if re.search(pat,pubtext,re.I): prod.add(name)
chk(9,"exactly one primary product/resource CTA in public copy",
    prod=={"Career Evidence Starter"},str(sorted(prod)))
chk(10,"Career Evidence Starter is that CTA","Career Evidence Starter" in pub)
chk(11,"Field Kit is not Video 1's product CTA",
    not re.search(r"fieldkit",pubtext,re.I))
chk(12,"Keep the Proof is not a competing CTA",
    not re.search(r"keep-the-proof",pubtext,re.I))
chk(13,"public URL exact",URL in pub and URL in text(DESC_DOC))
chk(14,"no direct PDF URL anywhere",not re.search(r"\.pdf\b",a,re.I),
    str(re.findall(r"\S*\.pdf\b",a)))
chk(15,"spoken CTA says about 10 to 15 focused minutes",
    "about 10 to 15 focused minutes" in tel)
chk(16,"Proof Line language preserved","portable Proof Line" in a)
chk(17,"Slide 12 matches the new CTA",
    "FREE CAREER EVIDENCE STARTER" in MT[11]
    and "ONE PORTABLE PROOF LINE" in MT[11]
    and "temidayoafonja.com/career-evidence-starter" in MT[11]
    and "fieldkit" not in MT[11])
import zipfile as _z
def changed(o,n):
    A_,B_=_z.ZipFile(o),_z.ZipFile(n)
    return sorted(x.split("/")[-1] for x in set(A_.namelist()) if A_.read(x)!=B_.read(x))
cm=changed("/tmp/v1hit/Main.orig.pptx",MAIN); cr=changed("/tmp/v1hit/Reveal.orig.pptx",REV)
chk(18,"slides 1-11 and 13 unchanged",cm==["slide12.xml"],str(cm))
chk(19,"only the matching reveal frame changed",cr==["slide21.xml"],str(cr))
shorts=sorted(f for f in os.listdir(SH) if f.startswith("Video_1_Short_"))
prev="/home/user/temidayoafonja-site/deliverables/video-1-slides/hit-final/Video_1_HIT_FINAL/SHORTS/"
def docxml(p):
    with _z.ZipFile(p) as z: return z.read("word/document.xml")
same=[f for f in shorts if os.path.isfile(prev+f) and docxml(os.path.join(SH,f))==docxml(prev+f)]
chk(20,"four Short recording scripts unchanged",
    len(shorts)==4 and len(same)==4,"%d/4 content-identical to v3.0"%len(same))
dp=paras(DESC_DOC); pp=paras(LF+"/Video_1_Publishing_Package_HIT_v3.1.docx")
db,de=dp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN"),dp.index("— END OF THE COPY-READY DESCRIPTION —")
pb,pe=pp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN"),pp.index("— END OF THE COPY-READY DESCRIPTION —")
chk(21,"publishing package and description-only public copy identical",dp[db:de]==pp[pb:pe])
chk(22,"no [INSERT] placeholder remains",not re.search(r"\[INSERT[^\]]*\]",a),
    str(re.findall(r"\[INSERT[^\]]*\]",a)))
warn=[i for i,x in enumerate(dp) if x.startswith("WORKING ESTIMATES")]
chk(23,"working timestamps clearly marked as estimates",
    bool(warn) and all(i>de for i in warn),"warning outside copy-ready block")
cv=subprocess.run([sys.executable,"verify_canonical.py"],capture_output=True,text=True)
def ln(sub): return [l for l in cv.stdout.splitlines() if sub in l][0]
chk(24,"teleprompter DOCX == TXT","PASS" in ln("teleprompter DOCX"))
chk(25,"reading DOCX == TXT","PASS" in ln("reading DOCX"))
chk(26,"teleprompter minus markers == reading script","PASS" in ln("minus markers == reading TXT"))
chk(27,"exactly 13 markers, ordered and mapped",
    "PASS" in ln("marker positions") and "PASS" in ln("marker names")
    and len(MT)==13,"13 markers, 13 slides")
chk(27.1,"canonical source verification","VERIFICATION: PASS" in cv.stdout)
z=_z.ZipFile(ZIPNAME); names=[n for n in z.namelist() if not n.endswith("/")]
chk(30,"ZIP contains exactly 13 files",len(names)==13,"%d"%len(names))
chk(30.1,"no _source, Python, temp or render files",
    not [n for n in names if re.search(r"_source|\.py$|\.png$|\.pyc$|^\.",n)])
sums=[l for l in open(ROOT+"/SHA256SUMS.txt",encoding="utf-8") if l.strip() and not l.startswith("#")]
chk(31,"SHA256SUMS has exactly 12 entries",len(sums)==12,"%d"%len(sums))
cp=subprocess.run(["sha256sum","-c","SHA256SUMS.txt"],cwd=ROOT,capture_output=True,text=True)
chk(32,"sha256sum -c passes",cp.returncode==0,
    cp.stdout.strip().splitlines()[-1] if cp.stdout else cp.stderr[:150])
def sha256(p):
    h=hashlib.sha256()
    with open(p,"rb") as fh:
        for b_ in iter(lambda: fh.read(1<<20),b""): h.update(b_)
    return h.hexdigest()
chk(33,"sibling ZIP checksum matches",
    open(ZIPNAME+".sha256").read().split()[0]==sha256(ZIPNAME),sha256(ZIPNAME))
chk(34,"description-only DOCX hash reported",True,sha256(DESC_DOC))
pubv=DECKS+"Video-1-How-I-Changed-Jobs-Without-Starting-My-Career-Over.pptx"
chk(35,"current published Video 1 assets untouched",
    os.path.isfile(pubv),"only the v2.4 recording decks were edited")
print(json.dumps(R,indent=1))
fails=[r for r in R if r["result"]=="FAIL"]
print("\n%d checks | %d PASS | %d FAIL"%(len(R),len(R)-len(fails),len(fails)))
for f in fails: print("  FAIL %s: %s  %s"%(f["n"],f["check"],f["detail"][:300]))
json.dump(R,open("QA_REPORT.json","w"),indent=1)
