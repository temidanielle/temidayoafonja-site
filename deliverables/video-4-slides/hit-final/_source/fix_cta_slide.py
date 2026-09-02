# -*- coding: utf-8 -*-
"""Apply the authorised Video 4 CTA slide correction to both decks.

Text only. Slide and reveal counts, design, typography family, colours, layout
and the Slide 11 Watch Next card are untouched. One type size comes down,
measured against the real embedded font first:

  headline  FREE CAREER EVIDENCE STARTER  11.47in at 46pt in a 10.42in box, so
            46pt -> 41pt gives 10.21in on one line, ending 0.21in clear of the
            gold rule. Two lines at 46pt would run 1.53in and cross the rule.
  sub-copy  ONE ACCOMPLISHMENT -> ONE PORTABLE PROOF LINE  8.66in in 10.42in
            at the existing 25pt: unchanged.
  URL       temidayoafonja.com/career-evidence-starter  10.07in in 10.83in at
            the existing 31pt: unchanged.
"""
import copy
from pptx import Presentation
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def q(t): return "{%s}%s" % (A, t)

EDITS = [
    ("KEEP THE PROOF", ["FREE CAREER EVIDENCE STARTER"], 4100),
    ("A 60-Minute Career Evidence System",
     ["ONE ACCOMPLISHMENT → ONE PORTABLE PROOF LINE"], None),
    ("temidayoafonja.com/keep-the-proof",
     ["temidayoafonja.com/career-evidence-starter"], None),
]
TARGETS = {"Video_4_Main_Slides.pptx": [10], "Video_4_Reveal_Builds.pptx": [25]}

def set_lines(shape, lines, size=None):
    body = shape.text_frame._txBody
    ps = body.findall(q("p"))
    assert len(ps) == 1, "expected a single paragraph"
    p = ps[0]
    runs = p.findall(q("r"))
    assert runs, "expected a run to clone"
    template = runs[0]
    for child in list(p):
        if etree.QName(child).localname in ("r", "br"):
            p.remove(child)
    for i, line in enumerate(lines):
        if i:
            etree.SubElement(p, q("br"))
        r = copy.deepcopy(template)
        r.find(q("t")).text = line
        if size is not None:
            r.find(q("rPr")).set("sz", str(size))
        p.append(r)

total = 0
for path, slides in TARGETS.items():
    prs = Presentation(path)
    for n in slides:
        for sh in prs.slides[n - 1].shapes:
            if not sh.has_text_frame:
                continue
            for frm, to, size in EDITS:
                if sh.text_frame.text == frm:
                    set_lines(sh, to, size)
                    total += 1
                    print("  %-26s slide %-3d %r -> %r"
                          % (path, n, frm[:30], to[0][:36]))
    prs.save(path)
assert total == 6, "expected 6 edits (3 per deck), got %d" % total
print("edits applied:", total)
