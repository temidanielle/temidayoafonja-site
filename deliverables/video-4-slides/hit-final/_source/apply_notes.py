# -*- coding: utf-8 -*-
"""Rewrite the speaker notes of one video's two decks.

Only notesSlideN.xml parts may change. Everything else -- slide XML, media,
rels, theme, presentation.xml -- must be byte-identical afterwards."""
import sys, zipfile, hashlib, os, copy, importlib
from pptx import Presentation
from pptx.oxml.ns import qn

def parts(path):
    with zipfile.ZipFile(path) as z:
        return {i.filename: hashlib.sha256(z.read(i.filename)).hexdigest()
                for i in z.infolist()}

def set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    src = None
    for p in tf.paragraphs:
        for r in p.runs:
            src = r._r.find(qn('a:rPr')); break
        if src is not None: break
    body = tf._txBody
    for p in body.findall(qn('a:p')): body.remove(p)
    for line in text.split("\n\n"):
        p = body.makeelement(qn('a:p'), {}); body.append(p)
        r = p.makeelement(qn('a:r'), {}); p.append(r)
        if src is not None: r.append(copy.deepcopy(src))
        t = p.makeelement(qn('a:t'), {}); t.text = line.replace("\n", " ")
        r.append(t)

def run(path, notes):
    before = parts(path)
    pres = Presentation(path)
    assert len(pres.slides) == len(notes), (path, len(pres.slides), len(notes))
    for s, n in zip(pres.slides, notes): set_notes(s, n)
    pres.save(path)
    after = parts(path)
    changed = sorted(k for k in after if before.get(k) != after[k])
    bad = [c for c in changed if not c.startswith("ppt/notesSlides/notesSlide")]
    return changed, sorted(set(after)-set(before)), sorted(set(before)-set(after)), bad

if __name__ == "__main__":
    mod, main_pptx, rev_pptx = sys.argv[1], sys.argv[2], sys.argv[3]
    sys.path.insert(0, os.path.dirname(os.path.abspath(mod)))
    m = importlib.import_module(os.path.basename(mod)[:-3])
    out={}
    for path, notes, key in ((main_pptx, m.NOTES, "main"),
                             (rev_pptx, m.reveal_notes(), "reveal")):
        ch, ad, rm, bad = run(path, notes)
        out[key]={"changed":ch,"added":ad,"removed":rm,"nonnotes":bad}
        print("\n==", os.path.basename(path))
        print("  changed parts: %d"%len(ch))
        print("  added:", ad, " removed:", rm)
        print("  NON-NOTES PARTS CHANGED:", bad)
    import json
    json.dump({"main_nonnotes":out["main"]["nonnotes"],
               "reveal_nonnotes":out["reveal"]["nonnotes"],
               "main_changed":out["main"]["changed"],
               "reveal_changed":out["reveal"]["changed"]},
              open(os.path.join(os.path.dirname(os.path.abspath(mod)),
                                "_partdiff.json"),"w"), indent=1)
