# -*- coding: utf-8 -*-
"""QA for the Videos 4, 6 and 8 patch pass."""
import os, re, json, zipfile, hashlib, subprocess, sys
from docx import Document
from pptx import Presentation

DECKS="/home/user/temidayoafonja-site/deliverables/video-%s-slides/out/"
R=[]
def chk(v,n,desc,ok,detail=""):
    R.append({"video":v,"n":n,"check":desc,"result":"PASS" if ok else "FAIL","detail":detail})
def paras(p): return [x.text for x in Document(p).paragraphs]
def text(p): return "\n".join(paras(p))
def body(p): return text(p) if p.endswith(".docx") else open(p,encoding="utf-8").read()

def files_of(root, descdoc):
    pkg=[os.path.join(dp,f) for dp,_,fn in os.walk(root) for f in fn
         if f.endswith((".docx",".txt"))]
    return sorted(pkg)+[descdoc]

def blocks(p):
    return paras(p) if p.endswith(".docx") else open(p,encoding="utf-8").read().split("\n\n")
def context(bl,i):
    win=bl[max(0,i-2):i+3]
    for j in range(i-1,max(-1,i-12),-1):
        prev=bl[j].strip()
        if prev.endswith(":"): win=[prev]+win; break
        if prev and not prev.startswith(("—","-","•"," ")): break
    return " ".join(win)
def scan(v,n,label,files,pats,reject):
    bad=[]
    for p in files:
        bl=blocks(p)
        for i,para in enumerate(bl):
            for pat in pats:
                if re.search(pat,para,re.I) and not re.search(reject,context(bl,i),re.I):
                    bad.append((os.path.basename(p),para[:90]))
    chk(v,n,label,not bad,str(bad))

def pkg_checks(v, root, descdoc, zipname, manifest_n=13):
    F=files_of(root,descdoc)
    z=zipfile.ZipFile(zipname); names=[x for x in z.namelist() if not x.endswith("/")]
    chk(v,"zip","ZIP contains exactly 13 files",len(names)==manifest_n,"%d"%len(names))
    chk(v,"clean","no _source, Python, temp or render files in ZIP",
        not [x for x in names if re.search(r"_source|\.py$|\.png$|\.pyc$|^\.",x)])
    sums=[l for l in open(root+"/SHA256SUMS.txt",encoding="utf-8")
          if l.strip() and not l.startswith("#")]
    chk(v,"sums","SHA256SUMS has exactly 12 entries",len(sums)==12,"%d"%len(sums))
    cp=subprocess.run(["sha256sum","-c","SHA256SUMS.txt"],cwd=root,capture_output=True,text=True)
    chk(v,"sumsc","sha256sum -c passes",cp.returncode==0,
        cp.stdout.strip().splitlines()[-1] if cp.stdout else cp.stderr[:120])
    def sha(p):
        h=hashlib.sha256()
        with open(p,"rb") as fh:
            for b in iter(lambda: fh.read(1<<20),b""): h.update(b)
        return h.hexdigest()
    chk(v,"sib","sibling ZIP checksum matches",
        open(zipname+".sha256").read().split()[0]==sha(zipname),sha(zipname))
    readme=open(root+"/README_FINAL.txt",encoding="utf-8").read()
    listed=[x.split("/",1)[1] for x in names]
    chk(v,"readme","README matches archive",
        not [f for f in listed if os.path.basename(f) not in readme])
    return F, sha(zipname), sha(descdoc)

def copy_identical(v, pub, descdoc):
    dp,pp=paras(descdoc),paras(pub)
    db=dp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
    de=dp.index("— END OF THE COPY-READY DESCRIPTION —")
    pb=pp.index("COPY-READY YOUTUBE DESCRIPTION — BEGIN")
    pe=pp.index("— END OF THE COPY-READY DESCRIPTION —")
    chk(v,"copy","publishing and description-only public copy identical",dp[db:de]==pp[pb:pe])
    chk(v,"warn","working-estimates warning outside the copy-ready block",
        all(i>de for i,x in enumerate(dp) if x.startswith("WORKING ESTIMATES")))

def deck_parts(orig, now):
    a,b=zipfile.ZipFile(orig),zipfile.ZipFile(now)
    return sorted(x.split("/")[-1] for x in set(a.namelist()) if a.read(x)!=b.read(x))

# ============================================================ VIDEO 4
root="/tmp/v4p/Video_4_HIT_FINAL"; dd="/tmp/v4p/Video_4_YouTube_Description_HIT.docx"
zp="/tmp/v4p/Video_4_HIT_FINAL_Recording_and_Shorts_Package.zip"
F4,z4,d4=pkg_checks("4",root,dd,zp)
pub4=root+"/LONG_FORM/Video_4_Publishing_Package_HIT_v2.1.docx"
pub4t=text(pub4)   # the TEXT; earlier revisions tested the path string
a4="\n".join(body(p) for p in F4)
tel4=open(root+"/LONG_FORM/Video4TeleprompterScriptwithslidemarkers_HIT_v2.1.txt",encoding="utf-8").read()
PUB4=[p for p in F4 if "EDITOR_ONLY" not in p and os.path.basename(p)!="README_FINAL.txt"]
pt4="\n".join(body(p) for p in PUB4)
chk("4",1,"title unchanged","How to Explain Your Career Change" in pub4t)
chk("4",2,"thumbnail unchanged","YOUR CAREER MAKES SENSE" in pub4t)
chk("4",3,"Career Evidence Starter is sole primary CTA",
    "career-evidence-starter" in pt4 and not re.search(r"keep-the-proof|fieldkit|Decision Evidence Check",pt4,re.I),
    "")
chk("4",4,"Keep the Proof no longer Video 4's CTA",
    not re.search(r"keep-the-proof",pt4,re.I))
chk("4",5,"no direct PDF URL",not re.search(r"\.pdf\b",a4,re.I))
chk("4",6,"three-sentence framework unchanged",
    "My career has moved across" in a4 or "three-sentence" in a4.lower())
chk("4",7,"spoken CTA says 10 to 15 focused minutes",
    "about 10 to 15 focused minutes" in tel4)
chk("4",8,"Proof Line language present","portable Proof Line" in a4)
scan("4",9,"no 30% retention figure",F4,[r"\b30\s?%"],r"do not|never|must not|excluded|\bno\b")
scan("4",10,"no >$2M figure",F4,[r"\$2\s?M|\$2 million"],r"do not|never|must not|excluded|\bno\b|neither")
chk("4",11,"no [INSERT] placeholder",not re.search(r"\[INSERT[^\]]*\]",a4),
    str(re.findall(r"\[INSERT[^\]]*\]",a4)))
chk("4",12,"exactly four Short recording scripts",
    len([f for f in os.listdir(root+"/SHORTS") if f.startswith("Video_4_Short_")])==4)
mp=deck_parts("/tmp/v4p/Main.orig.pptx",DECKS%"4"+"Video_4_Main_Slides.pptx")
rp=deck_parts("/tmp/v4p/Reveal.orig.pptx",DECKS%"4"+"Video_4_Reveal_Builds.pptx")
chk("4",13,"only the CTA slide and its reveal frame changed",
    mp==["slide10.xml"] and rp==["slide25.xml"],"main %s | reveal %s"%(mp,rp))
m4=Presentation(DECKS%"4"+"Video_4_Main_Slides.pptx"); r4=Presentation(DECKS%"4"+"Video_4_Reveal_Builds.pptx")
MT=["\n".join(s.text_frame.text for s in sl.shapes if s.has_text_frame) for sl in m4.slides]
chk("4",14,"slide counts preserved",
    len(m4.slides._sldIdLst)==11 and len(r4.slides._sldIdLst)==26,
    "%d slides, %d frames"%(len(m4.slides._sldIdLst),len(r4.slides._sldIdLst)))
chk("4",15,"CTA slide carries the new CTA",
    "FREE CAREER EVIDENCE STARTER" in MT[9] and "career-evidence-starter" in MT[9]
    and "KEEP THE PROOF" not in MT[9])
chk("4",16,"Slide 11 Watch Next unchanged","WATCH NEXT" in MT[10])
copy_identical("4",pub4,dd)

# ============================================================ VIDEO 6
root="/tmp/v6p/Video_6_HIT_FINAL"; dd="/tmp/v6p/Video_6_YouTube_Description_HIT.docx"
zp="/tmp/v6p/Video_6_HIT_FINAL_Recording_and_Shorts_Package.zip"
F6,z6,d6=pkg_checks("6",root,dd,zp)
pub6=root+"/LONG_FORM/Video_6_Publishing_Package_HIT_v2.1.docx"; pub6t=text(pub6)
a6="\n".join(body(p) for p in F6)
tel6=open(root+"/LONG_FORM/Video6TeleprompterScriptwithslidemarkers_HIT_v2.1.txt",encoding="utf-8").read()
chk("6",1,"title unchanged","Are You Growing—or Just Being Given More Work?" in pub6t)
chk("6",2,"thumbnail unchanged","MORE WORK ≠ GROWTH" in pub6t)
chk("6",3,"Field Kit remains sole CTA",
    "fieldkit" in a6 and not re.search(r"keep-the-proof|career-evidence-starter",a6,re.I))
chk("6",4,"CAR named 2-3 times in long-form",
    2 <= tel6.count("CAR test") <= 3, "%d occurrences"%tel6.count("CAR test"))
chk("6",5,"CAR always means Complexity, Authority, Return",
    "CAR test: Complexity, Authority and Return" in tel6)
chk("6",6,"no second mnemonic introduced",
    not re.search(r"\b3 Cs\b|three Cs|Capture, Clarify",a6))
chk("6",7,"Short 3 carries the one authorised line",
    "I run it through the CAR test." in text(root+"/SHORTS/Video_6_Short_3_More_Scope_Was_It_Growth.docx"))
for f,frag in (("Video_6_Short_1_Workload_Grows_Faster.docx","Your workload can grow faster"),
               ("Video_6_Short_2_Accountability_Not_Authority.docx","Accountability is not the same"),
               ("Video_6_Short_4_What_Comes_Off_Your_Plate.docx","comes off your plate")):
    chk("6","s"+f[13],"Short %s unchanged"%f[13],
        frag in text(root+"/SHORTS/"+f) and "CAR test" not in text(root+"/SHORTS/"+f))
chk("6",8,"working chapter title updated","The CAR Test for Real Growth" in a6
    and "The 3 Tests for Real Growth" not in a6)
chk("6",9,"public description names the CAR test",
    "I use the CAR test—Complexity, Authority and Return—to" in pub6t)
chk("6",10,"no [INSERT] placeholder",not re.search(r"\[INSERT[^\]]*\]",a6))
chk("6",11,"exactly four Short recording scripts",
    len([f for f in os.listdir(root+"/SHORTS") if f.startswith("Video_6_Short_")])==4)
D6=DECKS%"6"
import glob
chk("6",12,"no PowerPoint or reveal deck changed",
    subprocess.run(["git","status","--porcelain","deliverables/video-6-slides/out"],
        cwd="/home/user/temidayoafonja-site",capture_output=True,text=True).stdout.strip()=="",
    "git reports no change")
chk("6",13,"editor overlay direction present",
    "THE CAR TEST" in text(root+"/LONG_FORM/Video_6_EDITOR_ONLY_HIT_Brief_v2.1.docx"))
copy_identical("6",pub6,dd)

# ============================================================ VIDEO 8
root="/tmp/v8p/Video_8_HIT_FINAL"; dd="/tmp/v8p/Video_8_YouTube_Description_HIT.docx"
zp="/tmp/v8p/Video_8_HIT_FINAL_Recording_and_Shorts_Package.zip"
F8,z8,d8=pkg_checks("8",root,dd,zp)
pub8=root+"/LONG_FORM/Video_8_Publishing_Package_HIT_v2.1.docx"; pub8t=text(pub8)
a8="\n".join(body(p) for p in F8)
tel8=open(root+"/LONG_FORM/Video8TeleprompterScriptwithslidemarkers_HIT_v2.1.txt",encoding="utf-8").read()
chk("8",1,"title unchanged","How to Switch Industries Without Starting Over" in pub8t)
chk("8",2,"thumbnail unchanged","YOUR EXPERIENCE STILL COUNTS" in pub8t)
chk("8",3,"Field Kit remains sole CTA",
    "fieldkit" in a8 and not re.search(r"keep-the-proof|career-evidence-starter",a8,re.I))
chk("8",4,"three Cs named 2-3 times in long-form",
    2 <= tel8.count("three Cs") <= 3, "%d occurrences"%tel8.count("three Cs"))
chk("8",5,"the Cs always mean Capability, Context, Credential",
    "three Cs of an industry change: Capability, Context and Credential" in tel8)
chk("8",6,"no second mnemonic introduced",not re.search(r"CAR test|Capture, Clarify",a8))
chk("8",7,"Short 2 carries the one authorised line",
    "Before you change industries, use the three Cs:" in
    text(root+"/SHORTS/Video_8_Short_2_Stop_Calling_Everything_Transferable.docx"))
chk("8",8,"CISM first-attempt non-pass wording unchanged",
    "didn’t pass the first time" in a8)
scan("8",9,"no claim that all experience transfers",F8,
     [r"everything transfers",r"all (of your |your )?experience transfers"],
     r"not everything|does not|do not|never|but a new context|avoid|mistake|pretend|is not|rather than")
chk("8",10,"working chapter title updated",
    "The 3 Cs: Capability, Context and Credential" in a8)
chk("8",11,"public description names the three Cs",
    "I use the three Cs of an industry change—Capability, Context and Credential—to" in pub8t)
chk("8",12,"no [INSERT] placeholder",not re.search(r"\[INSERT[^\]]*\]",a8))
chk("8",13,"exactly four Short recording scripts",
    len([f for f in os.listdir(root+"/SHORTS") if f.startswith("Video_8_Short_")])==4)
chk("8",14,"no PowerPoint or reveal deck changed",
    subprocess.run(["git","status","--porcelain","deliverables/video-8-slides/out"],
        cwd="/home/user/temidayoafonja-site",capture_output=True,text=True).stdout.strip()=="",
    "git reports no change")
chk("8",15,"editor cue present",
    "THE 3 Cs OF AN INDUSTRY CHANGE" in text(root+"/LONG_FORM/Video_8_EDITOR_ONLY_HIT_Brief_v2.1.docx"))
copy_identical("8",pub8,dd)

json.dump(R,open("/tmp/patchpass/QA_PATCH_PASS.json","w"),indent=1)
fails=[r for r in R if r["result"]=="FAIL"]
print("\n%d checks | %d PASS | %d FAIL"%(len(R),len(R)-len(fails),len(fails)))
for f in fails: print("  FAIL v%s %s: %s | %s"%(f["video"],f["n"],f["check"],f["detail"][:220]))
print("\nZIP  v4 %s\nZIP  v6 %s\nZIP  v8 %s"%(z4,z6,z8))
print("DESC v4 %s\nDESC v6 %s\nDESC v8 %s"%(d4,d6,d8))
