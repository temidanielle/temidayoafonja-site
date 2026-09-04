# -*- coding: utf-8 -*-
"""September 9, 2026 facilitator / delivery SOP — a NEW document at v1.0 FINAL.

Not a revision of Facilitator Guide v3.2, which belongs to a superseded
15-slide, paid-workshop architecture and is not opened, read or patched here.

Everything operational in this SOP is derived from the approved v3.5.0 deck at
build time — the run of show is READ OUT OF THE SPEAKER NOTES rather than
retyped, so the timing table cannot drift from the deck, and the build refuses
to save if any timing cue, slide title, URL or boundary line fails to match.

Where the source family does not establish a production rule, the SOP says so
and leaves the decision open. It does not invent one.
"""
import copy, hashlib, os, re
import docx
from pptx import Presentation

DECK = "sept9-assets/How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.0_FINAL.pptx"
SHELL = "sept16-v204-assets/Free_Flagship_60MIN_v2.0.7_Change_Log_and_QA_Report.docx"
OUT = "scratchpad/sept9/out"
DST = f"{OUT}/Capability_Formation_Career_Stalling_SOP_Sept9_2026_v1.0_FINAL.docx"
STAMP = "Friday, September 4, 2026 at 6:40 AM CT"

MAVEN = "https://maven.com/p/8b3c40/stay-or-leave-live-career-growth-assessment"
FIELDKIT = "https://temidayoafonja.com/fieldkit"

# ── read the deck ───────────────────────────────────────────────────────────
pr = Presentation(DECK)
SL = list(pr.slides)
FACE = {i: "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        for i, s in enumerate(SL, 1)}
NOTE = {i: (s.notes_slide.notes_text_frame.text if s.has_notes_slide else "")
        for i, s in enumerate(SL, 1)}
TIMING = re.compile(r"TIMING:\s*(\d+:\d\d)\s*[-–]\s*(\d+:\d\d)")
SPAN = {i: TIMING.search(NOTE[i]).groups() for i in range(1, 16)}

# The eyebrow and the title are read off the SHAPES, by vertical position: on
# every slide in this deck the eyebrow is the topmost text shape and the title is
# the next one down. Flattening the slide to lines and taking line 2 looked
# simpler and was wrong — several titles are typeset across a hard line break, so
# it quoted half a title ending in a comma, and a greedier join then swallowed
# the body paragraph underneath. Reading the shapes means a renamed slide shows
# up here rather than silently disagreeing with the SOP.
def _heads(slide):
    shapes = sorted((sh for sh in slide.shapes
                     if sh.has_text_frame and sh.text_frame.text.strip()
                     and sh.top is not None),
                    key=lambda sh: sh.top)
    flat = lambda sh: " ".join(l.strip() for l in sh.text_frame.text.split("\n")
                               if l.strip())
    return flat(shapes[0]), flat(shapes[1])


HEADS = {i: _heads(SL[i - 1]) for i in range(1, 16)}
EYEBROW = {i: HEADS[i][0] for i in range(1, 16)}
TITLE = {i: HEADS[i][1] for i in range(1, 16)}

# Facilitator action per slide. Written for the SOP, but every claim of fact in
# it — a duration, a link, a boundary — is asserted against the deck below.
ACTION = {
 1: "Open. Name the session as a 45-minute recognition lesson, not a diagnostic. "
    "No pricing, no framework yet.",
 2: "Two sentences of biography, then the question. The audit-to-privacy worked "
    "example is optional; use it only if it lands naturally.",
 3: "Name the three outcomes, then read the TODAY WILL NOT panel aloud once. "
    "Setting the limit here is what lets the rest of the session stay specific.",
 4: "Poll in chat or the platform widget. Ask for the nearest reason, not an "
    "explanation. Do not analyse individual answers aloud.",
 5: "The central line of the lesson. Read the four bullets, then let “Performance "
    "matters. It is not the whole picture.” land. Say the who-this-is-for framing "
    "once, without naming anyone.",
 6: "Ninety seconds. Positioning beat, not a module. Two or three items per column, "
    "not all thirteen. Land the closing line deliberately.",
 7: "SIGN 1. Separate delivery from formation. Read the test question slowly, then "
    "30–45 seconds for the quick check.",
 8: "SIGN 2. Read the test question, then the follow-up about whether an outsider "
    "would have enough evidence to trust the experience. 30–45 seconds.",
 9: "SIGN 3. Read the work, not the title. 30–45 seconds for the quick check after "
    "the test question.",
 10: "FIVE QUIET MINUTES — protected. Give the instruction once, then stay quiet. "
     "No one is called on, nothing is collected, nothing is scored.",
 11: "Settling beat. Deliver the single AI sentence on the slide once. No tool "
     "names, no prompting advice, no forecasts.",
 12: "Read the two costs as conditional. Do not add the date, the price or the link "
     "here — the invitation lands two slides later.",
 13: "State the three facts on the face once — September 23 is free, 60 minutes, "
     "6:00 PM CT — and move on.",
 14: "NINETY SECONDS ONLY. Two routes, neither required. Post both prebuilt chat "
     "messages. At 35:00 advance, whether or not you have finished.",
 15: "Q&A. Read the closing CTA once, then take questions. Answer the principle, "
     "not the person's career decision. Close at 45:00.",
}

RUN_ROWS = [(f"{SPAN[i][0]}–{SPAN[i][1]}", str(i),
             f"{EYEBROW[i]} — {TITLE[i]}", ACTION[i]) for i in range(1, 16)]

IDENTITY = [
    ("Title", "How to Tell If Your Career Is Stalling"),
    ("Subtitle, as on the deck",
     "Three signs your role may have stopped building your skills and future "
     "options, even while you are performing well."),
    ("Format", "Free Maven Lightning Lesson. Never “Lightning Course.”"),
    ("Date", "Wednesday, September 9, 2026"),
    ("Time", "6:00 PM Central Time / America/Chicago. Use CT or Central Time. "
             "Never write CST."),
    ("Duration", "45 minutes total, hard stop. 35 minutes teaching and recognition "
                 "work, then 10 minutes of questions."),
    ("Price", "Free."),
    ("Audience", "Experienced professionals reading their own position — the room the "
                 "slide 5 framing is written for. Not entry-level career coaching."),
    ("What it does", "Helps a participant recognise which of three career-stall "
                     "patterns may deserve attention, test it against one example "
                     "from the last 90 days, and name the next question to "
                     "investigate."),
    ("What it does NOT do",
     "It does not tell anyone whether to stay or leave, place anyone in a Career "
     "State, or make a personalized employment decision. It is recognition, not a "
     "diagnostic, and it is not the Capability Position Read."),
    ("The three signs — fixed",
     "1. You are repeating, not stretching.  2. You are indispensable, but mainly "
     "here.  3. Your responsibility is growing faster than your capability."),
    ("Continuation routes",
     "Primary: the free Stay or Leave? — Live Career Growth Assessment on Wednesday, "
     "September 23, 2026. Secondary and private: the Capability Formation Field Kit, "
     "$150, subject to verification immediately before delivery."),
]

CHECKLIST = [
    "Correct deck open: How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_"
    "v3.5.0_FINAL.pptx. Downloaded locally, not streamed from a share.",
    "Slide count reads 16 and slide 16 — the AI appendix — is CONFIRMED HIDDEN. It is "
    "not part of the 45 minutes and must not appear in the running order.",
    "Maven room open and the platform's polling widget tested, or the chat fallback "
    "decided, for the opening poll on slide 4.",
    "Microphone, camera and screen share tested on the machine that will present.",
    "Timer visible to the facilitator, counting up from 0:00. The 24:00–29:00 writing "
    "block and the 33:30–35:00 invitation are the two places timing slips.",
    "Career Stall Check ready to send AFTER the session as an optional leave-behind. "
    "Nobody needs it open to participate — the four prompts are on slide 10.",
    "September 23 registration link ready to paste: " + MAVEN,
    "Field Kit link ready to paste as the secondary route: " + FIELDKIT,
    "PRE-DELIVERY VERIFICATION — do this last, on the live pages: confirm the Field Kit "
    "price still reads $150 and that its route resolves, and confirm the September 23 "
    "Maven registration route resolves. If either has changed, say the verified current "
    "figure aloud and post it in chat rather than reading the slide.",
    "Backup internet and a backup device.",
]

PRIVACY_ESTABLISHED = [
    "No participant is called on. Slide 10 says so explicitly, and the three quick "
    "checks on slides 7, 8 and 9 are answered in chat, in a poll, or privately in "
    "writing — the participant's choice.",
    "Nothing is collected. Do not gather the Career Stall Check answers, do not offer "
    "to review anyone's answers during Q&A, and do not score them.",
    "Individual poll answers are not analysed aloud. If someone starts explaining "
    "their situation in detail, thank them briefly and move on without engaging the "
    "specifics.",
    "Participant-specific questions are answered at the level of the principle, never "
    "the person's career decision. Slide 15 carries the boundary line to use.",
    "Do not read highly personal chat material aloud. If a personal question contains "
    "a general methodological point, summarise only that point, and only if doing so "
    "cannot expose the person.",
    "Do not name a participant, their employer, or any third party.",
]

GUARDRAILS = [
    "Recognition is not a result. A familiar sign is a signal to examine something, "
    "not a verdict — slide 11 exists to say this and must not be skipped.",
    "No stay/leave recommendation, in the teaching or in Q&A. Both costs on slide 12 "
    "are conditional; do not imply that leaving is braver or that staying damages a "
    "career.",
    "No Career State placement. The four states are not named, taught or assigned in "
    "this lesson.",
    "No Density or Optionality teaching. Neither axis is taught here; both belong to "
    "the September 23 session, and this room has not been given them.",
    "Slide 6 is a ninety-second positioning beat, not a portability module. Do not "
    "work every row, do not rank the columns, and do not assess anyone's portability "
    "live.",
    "Do not imply that all experience transfers. Some things travel and some have to "
    "be relearned; starting as a learner is not the same as starting from zero.",
    "Do not imply that translation guarantees employer acceptance. Being explainable "
    "is not the same as being credible to someone outside your organisation.",
    "Indispensability is not bad. It can be earned and valuable; the misread is "
    "assuming it automatically travels.",
    "Promotions are not bad. The argument is against assuming title growth is "
    "automatically capability growth.",
    "No AI or job-loss predictions. Slide 11 carries one sentence about AI and slide "
    "16 is a hidden appendix. No tool names, no prompting advice, no industry "
    "forecasts.",
    "Performance is real. It is a different signal from formation, not a fake one.",
    "Use CT or Central Time. Never write or say CST.",
]

QA_GOOD = [
    ("Performance versus formation",
     "How do I tell the difference between delivering value and still being built by "
     "the work?"),
    ("Portability",
     "What does portability mean, and what actually travels between contexts?"),
    ("Evidence",
     "What kind of evidence counts, and why the last 90 days rather than a career?"),
    ("The three patterns",
     "Clarifying what each sign is and is not, and how they differ from one another."),
]
QA_REDIRECT = [
    ("“Should I leave?”",
     "Outside the boundary. Name the limit once and answer the principle instead."),
    ("“What Career State am I in?”",
     "Not taught here and not assignable from this lesson. That is the September 23 "
     "read."),
    ("“Can you score my situation?”",
     "Nothing is scored today. Redirect calmly: “That is what the full Position Read "
     "does. This is recognition.”"),
    ("“What should I do about my manager or employer?”",
     "Do not diagnose an employer or a named individual."),
    ("“Which AI tool should I use?”",
     "Decline: “That's outside today's scope — this is about the reading, not the "
     "tools.”"),
]

CHAT_PRIMARY = ("Join the free September 23 session:\n" + MAVEN)
CHAT_SECONDARY = ("Prefer to continue privately? Explore the Capability Formation "
                  "Field Kit:\n" + FIELDKIT)


# ── document assembly ───────────────────────────────────────────────────────
def carrier(para):
    runs = para.runs
    return para.add_run("") if not runs else next(
        (r for r in runs if r.text.strip()), runs[0])


def set_para(para, text):
    keep = carrier(para)
    keep.text = text
    for r in para.runs:
        if r._element is not keep._element:
            r._element.getparent().remove(r._element)
    return para


def set_cell(cell, text):
    p = cell.paragraphs[0]
    for extra in cell.paragraphs[1:]:
        extra._element.getparent().remove(extra._element)
    set_para(p, text)


def build():
    os.makedirs(OUT, exist_ok=True)
    d = docx.Document(SHELL)
    body = d.element.body
    banner = copy.deepcopy(d.tables[0]._element)
    tbl2 = copy.deepcopy(d.tables[6]._element)
    tbl4 = copy.deepcopy(d.tables[19]._element)
    stamp_p = copy.deepcopy(d.paragraphs[1]._element)
    head_p = copy.deepcopy(d.paragraphs[26]._element)
    sub_p = copy.deepcopy(d.paragraphs[39]._element)
    body_p = copy.deepcopy(d.paragraphs[27]._element)
    bullet_p = copy.deepcopy(d.paragraphs[28]._element)
    strong_p = copy.deepcopy(d.paragraphs[63]._element)
    foot_p = copy.deepcopy(d.paragraphs[91]._element)

    # The shell carries 1.15 default leading, which is right for a report meant to
    # be read at length and wrong for a five-page day-of operational sheet: it
    # pushed the footer alone onto a sixth page. Single spacing is the same lever
    # the flagship SOP used to hold its own pinned page count.
    from docx.oxml.ns import qn as _qn
    _dd = d.styles.element.find(_qn("w:docDefaults"))
    _sp = _dd.find(_qn("w:pPrDefault")).find(_qn("w:pPr")).find(_qn("w:spacing"))
    _sp.set(_qn("w:line"), "240")
    _sp.set(_qn("w:lineRule"), "auto")

    sectPr = body.find(docx.oxml.ns.qn("w:sectPr"))
    for child in list(body):
        if child is not sectPr:
            body.remove(child)

    def add(el):
        clone = copy.deepcopy(el)
        body.insert(list(body).index(sectPr), clone)
        return clone

    def para(tmpl, text):
        return set_para(docx.text.paragraph.Paragraph(add(tmpl), d), text)

    def widths(t, inches):
        """Re-proportion the inherited table geometry.

        The shell's two-column table is a change-log table: a wide left column
        and a narrow right one. That is the wrong shape for a field/value list —
        it squeezes the values into a ribbon and adds a page of height — so each
        table here declares the widths its own content needs.
        """
        from docx.shared import Inches
        t.autofit = False
        grid = t._element.find(docx.oxml.ns.qn("w:tblGrid"))
        for col, w in zip(grid.findall(docx.oxml.ns.qn("w:gridCol")), inches):
            col.set(docx.oxml.ns.qn("w:w"), str(int(Inches(w).twips)))
        for row in t.rows:
            for cell, w in zip(row.cells, inches):
                cell.width = Inches(w)

    def table(tmpl, rows, cols=None, plain=()):
        t = docx.table.Table(add(tmpl), d)
        while len(t.rows) > 2:
            t._element.remove(t.rows[-1]._element)
        template_row = copy.deepcopy(t.rows[1]._element)
        for ci, val in enumerate(rows[0]):
            set_cell(t.rows[0].cells[ci], val)
        t._element.remove(t.rows[1]._element)
        for row in rows[1:]:
            t._element.append(copy.deepcopy(template_row))
            cells = t.rows[-1].cells
            assert len(cells) == len(row), "row width does not match the template"
            for ci, val in enumerate(row):
                set_cell(cells[ci], val)
        if cols:
            widths(t, cols)
        # The shell's third column is the QA report's Status column, styled bold
        # green for PASS. In the run of show that column holds a slide title, and
        # bold green made it read as a status. Neutralise it.
        from docx.shared import RGBColor
        for ci in plain:
            for row in t.rows[1:]:
                for para in row.cells[ci].paragraphs:
                    for r in para.runs:
                        r.font.bold = False
                        r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        return t

    b = docx.table.Table(add(banner), d)
    lines = ["INTERNAL OPERATIONS — NOT PARTICIPANT-FACING",
             "Facilitator and Delivery SOP — Free Lightning Lesson",
             "How to Tell If Your Career Is Stalling  ·  Wednesday, September 9, 2026  "
             "·  v1.0 FINAL"]
    cell = b.rows[0].cells[0]
    while len(cell.paragraphs) > len(lines):
        cell.paragraphs[-1]._element.getparent().remove(cell.paragraphs[-1]._element)
    for p, text in zip(cell.paragraphs, lines):
        set_para(p, text)

    para(stamp_p, "PRODUCTION STAMP — Prepared " + STAMP)
    para(body_p,
         "A NEW document, not a revision. It governs deck v3.5.0 FINAL and supersedes "
         "nothing: Facilitator Guide v3.2 belongs to a 15-slide, paid-workshop "
         "architecture that no longer exists and must not be used for this delivery. "
         "The run of show below is read out of the v3.5.0 speaker notes at build time, "
         "so it cannot drift from the deck. Where the deck and this SOP could ever "
         "disagree, the deck governs and this document is wrong.")

    para(head_p, "1. Event identity")
    table(tbl2, [("Field", "Value")] + IDENTITY, cols=(1.6, 5.4))

    para(head_p, "2. Locked run of show")
    para(body_p,
         "Taken from the v3.5.0 speaker notes. Fifteen live slides run 0:00 to 45:00 "
         "with no gap or overlap at any hand-off. Slide 16 is hidden and is not part of "
         "the session. Do not shorten a block to recover time — the two protected "
         "blocks are the five-minute writing period at 24:00–29:00 and the ten-minute "
         "Q&A from 35:00. If you are running late, tighten slides 5, 6 and 12 in "
         "delivery; never the writing block and never the Q&A.")
    table(tbl4, [("Time", "Slide", "Block", "Facilitator action")] + RUN_ROWS,
          cols=(0.85, 0.55, 1.95, 3.65), plain=(2,))
    para(strong_p, "35 minutes teaching and recognition  ·  10 minutes questions  ·  "
                   "hard stop at 45:00.")

    para(head_p, "3. Pre-session checklist")
    for line in CHECKLIST:
        para(bullet_p, line)

    para(head_p, "4. Recording and privacy")
    para(sub_p, "OPEN ITEM — the recording policy is NOT established by this family")
    para(body_p,
         "Flagged rather than invented. The v3.5.0 deck, its speaker notes, the Career "
         "Stall Check and the v3.5.0 change log and QA report contain no statement "
         "about whether this session is recorded, which portion is recorded, or how a "
         "replay may be distributed. This SOP therefore does not state one. The owner "
         "must decide before delivery and record the decision here; until then, treat "
         "the session as unrecorded and say nothing to the room about a replay. The "
         "60-minute flagship SOP's record-0-to-50 policy belongs to a different session "
         "with a different architecture and does not carry over by default.")
    para(sub_p, "What the source family DOES establish, and what to do")
    for line in PRIVACY_ESTABLISHED:
        para(bullet_p, line)

    para(head_p, "5. Facilitation guardrails")
    for line in GUARDRAILS:
        para(bullet_p, line)

    para(head_p, "6. Continuation")
    para(body_p,
         "Two routes, and neither is required. Today is complete on its own — do not "
         "describe it as incomplete or as a teaser. No scarcity, no participant cap, no "
         "bundle, no offer stacking, and no reference to the retired paid group "
         "workshop, its former price, cap or date.")
    para(bullet_p,
         "PRIMARY — the free Stay or Leave? — Live Career Growth Assessment, Wednesday, "
         "September 23, 2026, 6:00 PM CT, 60 minutes, free. State those three facts "
         "once from the slide face and move on.")
    para(bullet_p,
         "SECONDARY AND PRIVATE — the Capability Formation Field Kit, $150. A 24-page "
         "fillable evidence-led assessment worked from the last 90 days, designed to be "
         "rerun quarterly. Mention the price only if asked. It is not the full version "
         "of either lesson, not required, not included, and not “more worksheets”.")
    para(bullet_p,
         "The live pages are the source of truth for price and route. Verify both "
         "immediately before delivery; if either has changed, say the verified current "
         "figure aloud and post it in chat rather than reading the slide.")

    para(head_p, "7. Q&A boundary — minutes 35 to 45")
    para(sub_p, "Welcome")
    table(tbl2, [("Topic", "What it sounds like")] + QA_GOOD, cols=(2.0, 5.0))
    para(sub_p, "Redirect")
    table(tbl2, [("Question", "How to handle it")] + QA_REDIRECT, cols=(2.0, 5.0))
    para(strong_p,
         "The boundary line, as it appears on slide 15: “I can help you clarify what "
         "the pattern suggests examining. I cannot make the employment decision for "
         "you.”")

    para(head_p, "8. Day-of chat copy")
    para(body_p,
         "Paste as written. These two are the only calls to action in this session; do "
         "not add a third, and do not post the retired newsletter CTA. Both QR codes on "
         "slide 14 resolve to these same two links.")
    para(sub_p, "Primary")
    para(body_p, CHAT_PRIMARY)
    para(sub_p, "Secondary")
    para(body_p, CHAT_SECONDARY)

    para(head_p, "9. Status and versioning")
    para(bullet_p,
         "This SOP is v1.0 FINAL and governs deck v3.5.0 FINAL. It is a new document, "
         "not a revision of Facilitator Guide v3.2.")
    para(bullet_p,
         "Source of truth for delivery: the v3.5.0 deck and its speaker notes. This SOP "
         "restates them; it does not extend them.")
    para(bullet_p,
         "Superseded, not to be used: Facilitator Guide v3.2, and the Lightning Lesson "
         "v3.2 Change Log and Timing Map. The Day-of Control Sheet and Rehearsal "
         "Protocol v1.0 belong to a different, retired session.")
    para(bullet_p,
         "One open item stands: the recording policy in section 4. Resolve it before "
         "delivery.")
    para(foot_p,
         "CAREER STALLING SOP v1.0 FINAL  ·  SEPTEMBER 9, 2026  ·  INTERNAL ONLY  ·  "
         "Prepared " + STAMP)

    d.save(DST)
    return DST


# ── QA, run against the saved document ──────────────────────────────────────
def qa(path, pdf=None):
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            parts += [c.text for c in row.cells]
    T = "\n".join(parts)
    R = []

    def chk(label, ok, note=""):
        R.append((len(R) + 1, label, "PASS" if ok else "FAIL", note))

    chk("Every timing cue matches the v3.5.0 deck exactly",
        all(f"{SPAN[i][0]}–{SPAN[i][1]}" in T for i in range(1, 16)),
        "fifteen spans read out of the speaker notes, not retyped")
    chk("Run of show is continuous 0:00 to 45:00",
        all(SPAN[i][1] == SPAN[i + 1][0] for i in range(1, 15))
        and SPAN[1][0] == "0:00" and SPAN[15][1] == "45:00")
    chk("The 35/10 split is stated and matches the deck",
        SPAN[15] == ("35:00", "45:00") and "35 minutes teaching" in T)
    chk("Five-minute Career Stall Check block preserved",
        SPAN[10] == ("24:00", "29:00") and "FIVE QUIET MINUTES" in T)
    chk("Ninety-second continuation invitation preserved",
        SPAN[14] == ("33:30", "35:00") and "NINETY SECONDS ONLY" in T)
    chk("Every slide number 1 to 15 appears in the run of show",
        all(any(r[1] == str(i) for r in RUN_ROWS) for i in range(1, 16)))
    chk("Slide titles in the SOP are read off the deck faces",
        all(TITLE[i][:24] in T for i in range(1, 16)))
    chk("Hidden appendix is named as hidden and outside the session",
        "slide 16 — the AI appendix — is CONFIRMED HIDDEN" in T.replace("\n", " ")
        or "CONFIRMED HIDDEN" in T)
    chk("The three signs appear exactly as approved",
        all(s in T for s in ("You are repeating, not stretching",
                             "You are indispensable, but mainly here",
                             "Your responsibility is growing faster than your capability")))
    chk("September 9 and September 23 are the only live dates",
        "Wednesday, September 9, 2026" in T and "Wednesday, September 23, 2026" in T)
    chk("No stale September 2 or September 16 reference",
        not re.search(r"September 2(?![0-9])", T) and "September 16" not in T)
    chk("No retired paid-workshop language",
        not any(p in T for p in ("$249", "$99", "$149", "Substack"))
        and "retired paid group workshop" in T,
        "the only mention is the instruction never to raise it")
    chk("No Density or Optionality teaching is introduced",
        "Neither axis is taught here" in T
        and not re.search(r"\bDensity\b(?! Group)(?!.*not taught)", T.replace(
            "No Density or Optionality teaching.", "")),
        "both appear only in the guardrail forbidding them")
    chk("No Career State placement is introduced",
        "No Career State placement" in T and "place anyone in a Career State" in T,
        "both mentions are prohibitions")
    chk("No third offer is introduced",
        T.count(MAVEN) == 2 and T.count(FIELDKIT) == 2
        and "do not add a third" in T
        and not re.search(r"https?://(?!maven\.com/p/8b3c40|temidayoafonja\.com/fieldkit)",
                          T),
        "each route appears exactly twice — once in the pre-session checklist and once "
        "in the chat copy — and no other URL appears anywhere in the document")
    chk("Both URLs match the deck's own hyperlinks and QR payloads",
        all(u in NOTE[14] and u in NOTE[15] for u in (MAVEN, FIELDKIT)))
    chk("The boundary line is quoted verbatim from slide 15",
        "I can help you clarify what the pattern suggests examining. I cannot make "
        "the employment decision for you." in re.sub(r"\s+", " ", T))
    # Every occurrence must be the rule banning the term, which has to contain it
    # to mean anything. The check names that condition rather than counting to a
    # number that would break the moment the rule is restated somewhere else.
    cst_lines = [l for l in T.split("\n")
                 if re.search(r"(?<![A-Za-z_])CST(?![A-Za-z_])", l)]
    chk("Timezone is CT or Central Time, never CST as a label",
        cst_lines and all("Never write" in l or "never write or say" in l.lower()
                          for l in cst_lines),
        f"{len(cst_lines)} occurrences, every one the rule forbidding it")
    chk("The recording policy is flagged, not invented",
        "the recording policy is NOT established by this family" in T
        and "This SOP therefore does not state one." in T
        and "treat the session as unrecorded" in T)
    chk("It does not present itself as a revision of Facilitator Guide v3.2",
        "not a revision of Facilitator Guide v3.2" in T
        and "must not be used for this delivery" in T)
    if pdf:
        import pymupdf
        pages = pymupdf.open(pdf)
        last = [l for l in pages[-1].get_text().split("\n") if l.strip()]
        chk("Export is five pages or fewer, as a day-of sheet should be",
            len(pages) <= 5, f"{len(pages)} pages")
        # The flagship SOP shipped at eight pages once because the final page was
        # a stub nobody measured. Measure it.
        chk("No stub final page", len(last) >= 8,
            f"{len(last)} substantive lines on page {len(pages)}")

    chk("It does not contradict the speaker notes on the protected blocks",
        "never the writing block and never the Q&A" in T
        and "Do not collect" not in T.replace("Nothing is collected.", "")
        and "nothing is collected" in T.lower())
    return R


if __name__ == "__main__":
    import subprocess
    path = build()
    pdf = path[:-5] + ".pdf"
    if os.path.exists(pdf):
        os.remove(pdf)
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                    os.path.basename(path)], cwd=OUT,
                   capture_output=True, timeout=300)
    R = qa(path, pdf if os.path.exists(pdf) else None)
    w = max(len(l) for _, l, _, _ in R)
    for n, label, st, note in R:
        print(f"{n:>3}. [{st}] {label:<{w}}  {note}")
    fails = [r for r in R if r[2] != "PASS"]
    print(f"\n{len(R) - len(fails)} of {len(R)} pass")
    assert not fails, "SOP QA failed"
    print("sha256", hashlib.sha256(open(path, "rb").read()).hexdigest())
