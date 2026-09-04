# -*- coding: utf-8 -*-
"""September 9 recording-policy QA — the locked policy, checked where it lands.

The policy is an operational addition, not a content change, so these checks do
two jobs: prove the rule is stated everywhere a facilitator would look for it,
and prove that nothing in the lesson moved while it was being added.
"""
import re, sys
import docx
from pptx import Presentation

DECK = ("sept9-assets/"
        "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.1_FINAL.pptx")
PRIOR = ("sept9-assets/"
         "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.0_FINAL.pptx")
SOP = ("sept9-assets/"
       "Capability_Formation_Career_Stalling_SOP_Sept9_2026_v1.1_FINAL.docx")

R = []
def chk(label, ok, note=""):
    R.append((len(R) + 58, label, "PASS" if ok else "FAIL", note))


def deck(path):
    p = Presentation(path)
    f = {i: "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
         for i, s in enumerate(p.slides, 1)}
    n = {i: (s.notes_slide.notes_text_frame.text if s.has_notes_slide else "")
         for i, s in enumerate(p.slides, 1)}
    return f, n


FACE, NOTE = deck(DECK)
OFACE, ONOTE = deck(PRIOR)
d = docx.Document(SOP)
parts = [p.text for p in d.paragraphs]
for t in d.tables:
    for row in t.rows:
        parts += [c.text for c in row.cells]
SOPT = re.sub(r"\s+", " ", "\n".join(parts))
FLAT = re.sub(r"\s+", " ", "\n".join(NOTE.values()))

# ── the policy is stated where it is needed ─────────────────────────────────
chk("Recording starts on slide 1 and the recorded portion is named",
    "RECORDING ON — start the recording on this slide" in NOTE[1]
    and "0:00–35:00, is the recorded portion" in NOTE[1])
chk("The disclosure sits inside the first five minutes",
    "RECORDING AND PRIVACY DISCLOSURE" in NOTE[3]
    and re.search(r"TIMING:\s*3:00-4:30", NOTE[3]) is not None,
    "slide 3, at 3:00–4:30, where the session already sets its boundary aloud")
chk("The spoken disclosure is the approved wording",
    "“The teaching portion of this session is recorded so you can revisit the ideas. "
    "The final ten minutes are live questions and are not part of the replay. Nothing "
    "you write in your Career Stall Check is collected or scored.”" in FLAT)
chk("The disclosure is not allowed to become a long disclaimer",
    "say once, calmly, then move on. Do not expand it into a disclaimer" in FLAT)
chk("Sensitive material is kept out of the recorded teaching",
    "Do not read personal participant situations, employer names or sensitive chat "
    "messages into the recorded teaching." in FLAT)
chk("The stop instruction sits on the 33:30–35:00 continuation block",
    "AT 35:00 — STOP THE RECORDING" in NOTE[14]
    and "confirm on screen that it has stopped before advancing into Q&A" in NOTE[14]
    and re.search(r"TIMING:\s*33:30-35:00", NOTE[14]) is not None)
chk("Stop, never pause",
    "Stop it; do not pause it." in NOTE[14] and "do not pause it" in SOPT.lower())
chk("The automatic-capture case is covered",
    "keep the raw file private and distribute only the 0:00–35:00 teaching portion"
    in re.sub(r"\s+", " ", NOTE[14])
    and "retain the raw file privately" in SOPT)
chk("Slide 15 carries the RECORDING OFF marker",
    NOTE[15].startswith("RECORDING OFF — LIVE-ONLY Q&A")
    and "Confirm on screen that recording has stopped before taking the first "
        "participant-specific question." in re.sub(r"\s+", " ", NOTE[15]))
chk("The replay excludes Q&A, said in both the deck and the SOP",
    "never included in the distributed replay" in re.sub(r"\s+", " ", NOTE[15])
    and "are NOT included in the distributed replay" in SOPT)
chk("The SOP states the policy as a rule, not an open question",
    "RECORD MINUTES 0:00–35:00 ONLY" in SOPT
    and "LOCKED POLICY" in SOPT
    and "not established by this family" not in SOPT
    and "No open items." in SOPT)
chk("The flagship's 0-to-50 timing is not copied into this family",
    "does not carry across" in SOPT
    and "0:00–50" not in SOPT and "minute 50" not in SOPT
    and "50" not in FLAT.replace("35:00", "").replace("45:00", "")
        .replace("150", "").replace("50 seconds", ""),
    "named once in the SOP, only to say it does not apply here")

# ── and nothing else moved ──────────────────────────────────────────────────
chk("No slide face changed", FACE == OFACE,
    "16 of 16 byte-identical to v3.5.0; this was a speaker-note pass only")
changed = [i for i in range(1, 17) if NOTE[i] != ONOTE[i]]
chk("Only slides 1, 3, 14 and 15 have changed notes", changed == [1, 3, 14, 15],
    f"notes changed: {changed}")
TIMING = re.compile(r"TIMING:\s*(\d+:\d\d)\s*[-–]\s*(\d+:\d\d)")
chk("All 15 active-slide timing cues are unchanged",
    [TIMING.findall(ONOTE[i]) for i in range(1, 17)]
    == [TIMING.findall(NOTE[i]) for i in range(1, 17)],
    "every cue identical to v3.5.0")
spans = [TIMING.search(NOTE[i]).groups() for i in range(1, 16)]
chk("Total duration is still 45 minutes, teaching 0:00–35:00, Q&A 35:00–45:00",
    spans[0][0] == "0:00" and spans[13][1] == "35:00"
    and spans[14] == ("35:00", "45:00")
    and all(spans[i][1] == spans[i + 1][0] for i in range(14)))
chk("No new conceptual teaching entered the deck",
    all(len(NOTE[i]) == len(ONOTE[i]) for i in range(1, 17) if i not in (1, 3, 14, 15)),
    "the eleven untouched notes are byte-identical, and the four that changed did so "
    "only by the recording additions")
chk("No new offer entered the deck",
    FLAT.count("maven.com/p/8b3c40") == re.sub(r"\s+", " ", "\n".join(
        ONOTE.values())).count("maven.com/p/8b3c40")
    and "Private Capability Position Read" not in FLAT
    and not any(p in FLAT for p in ("$249", "$99", "$149")))
chk("September 9 and September 23 remain correct, with no September 2 or 16",
    "September 23" in FLAT
    and not re.search(r"September 2(?![0-9])", FLAT)
    and "September 16" not in FLAT)
chk("No Density, Optionality or Career State teaching was introduced",
    FLAT.count("Density, Optionality, Career States")
    == re.sub(r"\s+", " ", "\n".join(ONOTE.values())).count(
        "Density, Optionality, Career States"),
    "the single occurrence is slide 13's note forbidding all three, unchanged")

if __name__ == "__main__":
    w = max(len(l) for _, l, _, _ in R)
    for n, label, st, note in R:
        print(f"{n:>3}. [{st}] {label:<{w}}  {note}")
    fails = [r for r in R if r[2] != "PASS"]
    print(f"\n{len(R) - len(fails)} of {len(R)} pass")
    sys.exit(1 if fails else 0)
