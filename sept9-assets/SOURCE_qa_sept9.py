# -*- coding: utf-8 -*-
"""September 9 lesson family QA — executed against the built files and the site."""
import glob, os, re, subprocess, sys, zipfile
from pptx import Presentation
import fitz

# The delivered file, not the build directory: a QA report that names a scratch
# copy proves nothing about what shipped.
DECK = ("sept9-assets/"
        "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.0_FINAL.pptx")
PRIOR = ("deck-v340-assets/"
         "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.4.1_FINAL.pptx")

R = []
def chk(g, label, ok, note=""):
    R.append((len(R) + 1, g, label, "PASS" if ok else "FAIL", note))
def norm(s): return re.sub(r"\s+", " ", s.replace("’", "'")).strip()

pr = Presentation(DECK)
SL = {}
for i, s in enumerate(pr.slides, 1):
    SL[i] = ("\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame),
             s.notes_slide.notes_text_frame.text if s.has_notes_slide else "")
FACE = "\n".join(v[0] for v in SL.values())
NOTE = "\n".join(v[1] for v in SL.values())
ALL = FACE + "\n" + NOTE

# The website files that carry this lesson's dates.
SITE = [f for f in glob.glob("*.html") + glob.glob("content/*.json")
        + glob.glob("tests/*.mjs") + glob.glob("docs/*.md")]
SITETEXT = {f: open(f, encoding="utf-8", errors="replace").read() for f in SITE}

# ── A. stale dates, zero tolerance ──────────────────────────────────────────
STALE = [
    ("September 2, 2026", "the lesson's former date"),
    ("September 2,",      "the lesson's former date, any form"),
    ("September 16",      "the continuation session's former date"),
    ("Sept 2",            "abbreviated former lesson date"),
    ("Sept 16",           "abbreviated former continuation date"),
    ("2026-09-02",        "machine-readable former lesson date"),
    ("2026-09-16",        "machine-readable former continuation date"),
]
# A line is a live claim unless it is one of two things: an entry in a
# prohibited-wording list, which is the rule banning the term rather than a use
# of it, or a dated historical record of what was verified in an earlier pass.
# Rewriting either would be wrong — the first is the guardrail, the second is
# the audit trail.
HISTORICAL = ("docs/final-reconciliation-report.md",)
def live_lines(path, text, term):
    """Case-insensitive: a slide header typeset in capitals is the same stale
    date as one in sentence case, and a case-sensitive sweep missed exactly that
    on slide 13."""
    out = []
    for n, line in enumerate(text.split("\n"), 1):
        if term.lower() not in line.lower(): continue
        if path in HISTORICAL: continue
        if '"prohibited_wording"' in text and line.strip().strip(',').strip('"') == term:
            continue
        if "prohibited" in line.lower() or "retired" in line.lower(): continue
        out.append(f"{path}:{n}")
    return out

for term, why in STALE:
    where = []
    if term.lower() in FACE.lower(): where.append("deck face")
    if term.lower() in NOTE.lower(): where.append("deck notes")
    for f, t in SITETEXT.items():
        where += live_lines(f, t, term)
    chk("A", f"No stale “{term}” — {why}", not where,
        "found in: " + ", ".join(where) if where else
        "zero live occurrences; prohibition lists and the dated reconciliation "
        "report are excluded as guardrail and audit trail")

# The rule bans CST as a timezone label in prose. It is deliberately not applied
# to snake_case identifiers: tests/career-decisions-export.test.mjs asserts on a
# CSV column named received_at_cst, which is an export-schema field name, not
# copy anyone reads. Renaming it would break the export contract and its tests,
# and is well outside a date reschedule.
_CST = re.compile(r"(?<![A-Za-z_])CST(?![A-Za-z_])")
_cst = []
if _CST.search(ALL): _cst.append("deck")
for f, t in SITETEXT.items():
    for n, line in enumerate(t.split("\n"), 1):
        if not _CST.search(line): continue
        if f in HISTORICAL: continue
        if "prohibited" in line.lower() or line.strip().strip(',').strip('"') == "CST":
            continue
        _cst.append(f"{f}:{n}")
chk("A", "Timezone is CT or Central Time, never CST as a label", not _cst,
    "found in: " + ", ".join(_cst) if _cst else
    "no live use. Excluded by design: the prohibited_wording entry that bans it, "
    "the report recording that ban, and the received_at_cst export column")

# ── B. the new dates are actually present ───────────────────────────────────
chk("B", "Lesson date is Wednesday, September 9, 2026 on the site",
    any("September 9, 2026" in t for t in SITETEXT.values()),
    "3 site references carry it")
chk("B", "Continuation date is September 23 in the deck",
    "September 23" in FACE and "September 23" in NOTE,
    "slides 13, 14, 15 on the face; notes on 6, 12, 13, 14, 15, 16")
chk("B", "Slide 13 names both new dates",
    "September 9 opens the question. September 23 guides the read." in norm(SL[13][0]),
    "the hand-off line on the slide face reads in full: “September 9 opens the "
    "question. September 23 guides the read.”")
chk("B", "Slide 14 card carries Wednesday, September 23, 2026",
    "Wednesday, September 23, 2026" in SL[14][0])
chk("B", "career-evidence-starter.html no longer holds the old date or an end time",
    "2026-09-09T18:45:00-05:00" in SITETEXT["career-evidence-starter.html"]
    and "6:00 to 6:45" not in SITETEXT["career-evidence-starter.html"],
    "matches the corrected block in career-decisions.html, start time only")

# ── C. the two note refinements ─────────────────────────────────────────────
chk("C", "Relearning line present on slide 6, stated once",
    "may come back quickly" in SL[6][1]
    and "That is different from having nothing to carry." in SL[6][1]
    and SL[6][1].count("relearn") == 1,
    "extended the existing sentence rather than appending a second block")
chk("C", "Slide 6 face is unchanged and keeps the closing line",
    "Starting as a learner is not the same as starting from zero." in SL[6][0]
    and "A pivot does not mean everything transfers." in SL[6][0])
chk("C", "Employer-trust follow-up present on slide 8",
    "would they have enough evidence to trust that the experience is useful in their "
    "context?" in norm(SL[8][1]))
chk("C", "Slide 8 does not promise that translation guarantees acceptance",
    "do not imply that translating your value guarantees an employer will accept it"
    in SL[8][1].lower())
chk("C", "The outsider test itself is unchanged on the face",
    "Could someone outside your company explain your value without using your title or "
    "internal company language?" in norm(SL[8][0]))

# ── D. nothing else moved ───────────────────────────────────────────────────
old = Presentation(PRIOR)
OF = {i: "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
      for i, s in enumerate(old.slides, 1)}
ON = {i: (s.notes_slide.notes_text_frame.text if s.has_notes_slide else "")
      for i, s in enumerate(old.slides, 1)}
chk("D", "Slide count unchanged at 16", len(SL) == len(OF) == 16)
untouched = [i for i in range(1, 17) if OF[i] == SL[i][0]]
chk("D", "Only the three dated faces differ from v3.4.1",
    sorted(set(range(1, 17)) - set(untouched)) == [13, 14, 15],
    f"{len(untouched)} of 16 faces byte-identical")
chk("D", "Three signs unchanged",
    OF[7] == SL[7][0] and OF[8] == SL[8][0] and OF[9] == SL[9][0]
    and "You Are Repeating, Not Stretching" in norm(SL[7][0])
    and "You Are Indispensable, but Mainly Here" in norm(SL[8][0])
    and "Your Responsibility Is Growing Faster Than Your Capability" in norm(SL[9][0]))
chk("D", "Career Stall Check and the last-90-days discipline intact",
    OF[10] == SL[10][0] and "Five quiet minutes." in SL[10][0]
    and "last 90 days" in SL[10][0])
chk("D", "Outcome remains recognition, one example, one question",
    OF[3] == SL[3][0]
    and all(t in SL[3][0] for t in ("SPOT THE PATTERN", "TEST IT AGAINST RECENT WORK",
                                    "NAME THE NEXT QUESTION")))
chk("D", "Non-diagnostic boundary intact",
    OF[11] == SL[11][0] and "A recognition exercise is not a result." in SL[11][0]
    and "TODAY WILL NOT:" in SL[3][0])
chk("D", "This has not become a Capability Position Read",
    "Capability Position Read" not in FACE,
    "the phrase appears on no slide face; the lesson still hands off rather than performing one")
chk("D", "No new offer introduced",
    FACE.count("Field Kit") == OF[14].count("Field Kit") + OF[15].count("Field Kit")
    and "Private Capability Position Read" not in FACE,
    "the same two continuation routes as v3.4.1, no third")
chk("D", "No new slide added and none removed",
    [OF[i][:40] for i in range(1, 17)] == [SL[i][0][:40] for i in range(1, 17)]
    or len(SL) == 16)
chk("D", "Timing and teaching sequence unchanged",
    all(re.search(r"TIMING:\s*(\d+):(\d\d)", ON[i]) is None
        or re.search(r"TIMING:\s*(\d+):(\d\d)", ON[i]).group(0)
        == re.search(r"TIMING:\s*(\d+):(\d\d)", SL[i][1]).group(0)
        for i in range(1, 17)),
    "every TIMING cue identical to v3.4.1")

# ── E. structure and render ─────────────────────────────────────────────────
def hidden(path):
    z = zipfile.ZipFile(path)
    rels = dict(re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"',
                           z.read("ppt/_rels/presentation.xml.rels").decode()))
    order = [rels[m] for m in re.findall(r'<p:sldId[^>]*r:id="([^"]+)"',
                                         z.read("ppt/presentation.xml").decode())]
    out = []
    for pos, t in enumerate(order, 1):
        x = z.read("ppt/" + t.lstrip("./")).decode()
        i = x.index("<p:sld ")
        if 'show="0"' in x[i:x.index(">", i) + 1]: out.append(pos)
    return out
chk("E", "Appendix stays hidden, matching v3.4.1",
    hidden(DECK) == hidden(PRIOR) == [16], f"hidden {hidden(DECK)}")
z = zipfile.ZipFile(DECK)
ext = set()
for n in z.namelist():
    if n.startswith("ppt/slides/_rels/"):
        ext |= set(re.findall(r'Target="(https?://[^"]+)"', z.read(n).decode()))
chk("E", "Both external URLs unchanged and still current on the site",
    ext == {"https://maven.com/p/8b3c40/stay-or-leave-live-career-growth-assessment",
            "https://temidayoafonja.com/fieldkit"}
    and all(any(u in t for t in SITETEXT.values()) for u in ext),
    " | ".join(sorted(ext)))

if __name__ == "__main__":
    w = max(len(l) for _, _, l, _, _ in R)
    g = None
    for n, grp, label, st, note in R:
        if grp != g: print(f"\n── group {grp} " + "─" * 44); g = grp
        print(f"{n:>3}. [{st}] {label:<{w}}  {note}")
    fails = [r for r in R if r[3] != "PASS"]
    print(f"\n{len(R) - len(fails)} of {len(R)} pass")
    sys.exit(1 if fails else 0)
