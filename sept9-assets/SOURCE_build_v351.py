# -*- coding: utf-8 -*-
"""September 9 Lightning Lesson — v3.5.0 FINAL -> v3.5.1 FINAL.

One change only: the recording policy is locked. It was the single open item the
v1.0 facilitator SOP flagged rather than invented, and the owner has now decided
it. This is an operational policy addition, not a content revision.

  RECORD 0:00-35:00 ONLY. The teaching and recognition portion is recorded. At
  35:00 the recording is STOPPED — not paused — and confirmed stopped on screen
  before the first participant-specific question. Minutes 35:00-45:00 are
  live-only Q&A and are excluded from the distributed replay.

This is a September 9 policy. The 60-minute flagship's record-0-to-50 rule
belongs to a different architecture and is deliberately not copied across.

Four speaker notes gain a marker or a disclosure. No slide face changes, no
slide is added, removed or reordered, no timing cue moves, and no teaching
content is touched.
"""
import hashlib, os, shutil
from pptx import Presentation

SRC = ("sept9-assets/"
       "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.0_FINAL.pptx")
OUT = "scratchpad/sept9/out"
DST = (f"{OUT}/"
       "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.1_FINAL.pptx")

# ── slide 1: when the recording starts ──────────────────────────────────────
# The policy is worthless if the facilitator does not know where to begin, and
# the deck previously carried no recording marker at all.
S1_ADD = (
    "\n\nRECORDING ON — start the recording on this slide. The teaching and "
    "recognition portion, 0:00–35:00, is the recorded portion. Minutes 35:00–45:00 "
    "are live-only and are not part of the distributed replay.")

# ── slide 3: the early disclosure, inside the first five minutes ────────────
# Slide 3 is where the session already sets its boundary out loud, at 3:00–4:30,
# so the disclosure lands with the limits rather than as a separate beat.
S3_ADD = (
    "\n\nRECORDING AND PRIVACY DISCLOSURE — say once, calmly, then move on. Do not "
    "expand it into a disclaimer:\n"
    "“The teaching portion of this session is recorded so you can revisit the ideas. "
    "The final ten minutes are live questions and are not part of the replay. Nothing "
    "you write in your Career Stall Check is collected or scored.”\n"
    "Do not read personal participant situations, employer names or sensitive chat "
    "messages into the recorded teaching.")

# ── slide 14: the production instruction at the hand-off ────────────────────
S14_ADD = (
    "\n\nAT 35:00 — STOP THE RECORDING and confirm on screen that it has stopped "
    "before advancing into Q&A. Stop it; do not pause it. If the platform captures the "
    "whole meeting automatically, keep the raw file private and distribute only the "
    "0:00–35:00 teaching portion.")

# ── slide 15: the marker on the live-only block ─────────────────────────────
S15_ADD = (
    "RECORDING OFF — LIVE-ONLY Q&A. Confirm on screen that recording has stopped "
    "before taking the first participant-specific question. Minutes 35:00–45:00 are "
    "never included in the distributed replay.\n\n")

ADDITIONS = {1: ("append", S1_ADD), 3: ("append", S3_ADD),
             14: ("append", S14_ADD), 15: ("prepend", S15_ADD)}


def faces(pres):
    return ["\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for s in pres.slides]


def notes(pres):
    return [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
            for s in pres.slides]


def build():
    os.makedirs(OUT, exist_ok=True)
    shutil.copyfile(SRC, DST)
    p = Presentation(DST)
    sl = list(p.slides)

    # Anchor each edit to something on the slide, so a reordered deck fails here
    # rather than putting a recording instruction on the wrong beat.
    anchors = {
        1: "How to Tell If Your",
        3: "TODAY WILL NOT:",
        14: "TWO WAYS TO CONTINUE",
        15: "Questions & Applications",
    }
    for i, anchor in anchors.items():
        face = "\n".join(sh.text_frame.text for sh in sl[i - 1].shapes
                         if sh.has_text_frame)
        assert anchor in face, f"slide {i} is not the slide this edit belongs on"

    for i, (how, text) in ADDITIONS.items():
        tf = sl[i - 1].notes_slide.notes_text_frame
        assert "RECORDING" not in tf.text, f"slide {i} already carries a recording marker"
        tf.text = (tf.text + text) if how == "append" else (text + tf.text)

    p.save(DST)

    # ── prove the revision was controlled ───────────────────────────────────
    before, after = Presentation(SRC), Presentation(DST)
    bf, af, bn, an = faces(before), faces(after), notes(before), notes(after)
    assert len(af) == len(bf) == 16, "slide count moved"
    assert bf == af, "a slide face changed; this pass touches speaker notes only"

    changed = [i + 1 for i in range(16) if bn[i] != an[i]]
    for i in changed:
        how, text = ADDITIONS[i]
        expect = (bn[i - 1] + text) if how == "append" else (text + bn[i - 1])
        assert expect == an[i - 1], \
            f"slide {i}: a note changed by something other than its sanctioned addition"
    assert changed == sorted(ADDITIONS), f"unexpected note changes: {changed}"

    print("built", os.path.basename(DST))
    print(f"  faces changed: none   notes changed: {changed}")
    print("  sha256", hashlib.sha256(open(DST, "rb").read()).hexdigest())
    return DST


if __name__ == "__main__":
    build()
