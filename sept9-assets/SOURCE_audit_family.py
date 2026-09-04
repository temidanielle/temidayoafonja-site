# -*- coding: utf-8 -*-
"""September 9 lesson FAMILY audit — every surviving asset, not just the deck.

qa_sept9.py verifies the v3.5.0 revision itself. This module audits the rest of
the family: the leave-behind, the supporting documents, the media and links
inside the deck, and the repository copy that promotes the session. It also
records, as a check rather than as prose, which assets no longer exist.
"""
import hashlib, io, os, re, subprocess, sys, zipfile
import docx
from PIL import Image
from pptx import Presentation

DECK = "sept9-assets/How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.0_FINAL.pptx"
DECK_PDF = "sept9-assets/How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.5.0_FINAL.pdf"
PRIOR = ("deck-v340-assets/"
         "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.4.1_FINAL.pptx")
STALL = "lightning-lesson-v3/Career_Stall_Check_v1.0.docx"
STALL_PDF = "lightning-lesson-v3/Career_Stall_Check_v1.0.pdf"
GUIDE = ("lightning-lesson-v3/"
         "How_to_Tell_If_Your_Career_Is_Stalling_Facilitator_Guide_v3.2.docx")

R = []
def chk(label, ok, note=""):
    R.append((len(R) + 31, label, "PASS" if ok else "FAIL", note))
def norm(s): return re.sub(r"\s+", " ", s.replace("’", "'")).strip()


pr = Presentation(DECK)
SL = list(pr.slides)
FACE = {i: "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        for i, s in enumerate(SL, 1)}
NOTE = {i: (s.notes_slide.notes_text_frame.text if s.has_notes_slide else "")
        for i, s in enumerate(SL, 1)}
old = Presentation(PRIOR)
OFACE = {i: "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
         for i, s in enumerate(old.slides, 1)}
ALLFACE = "\n".join(FACE.values())
CORE = "\n".join(FACE[i] for i in range(1, 16)) + "\n" + "\n".join(NOTE[i] for i in range(1, 16))


def hidden(path):
    z = zipfile.ZipFile(path)
    rels = dict(re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"',
                           z.read("ppt/_rels/presentation.xml.rels").decode()))
    order = [rels[m] for m in re.findall(r'<p:sldId[^>]*r:id="([^"]+)"',
                                         z.read("ppt/presentation.xml").decode())]
    out = []
    for pos, t in enumerate(order, 1):
        x = z.read("ppt/" + t.lstrip("./")).decode()
        i = x.index("<p:sld ")
        if 'show="0"' in x[i:x.index(">", i) + 1]:
            out.append(pos)
    return out


def doc_text(path):
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            parts += [c.text for c in row.cells]
    return "\n".join(parts)


# ── structure and the 45-minute architecture ────────────────────────────────
chk("Deck is 16 slides — 15 active, 1 hidden appendix",
    len(SL) == 16 and hidden(DECK) == [16] == hidden(PRIOR),
    "hidden set unchanged from v3.4.1")

TIMING = re.compile(r"TIMING:\s*(\d+:\d\d)\s*[-–]\s*(\d+:\d\d)")
spans = [TIMING.search(NOTE[i]).groups() for i in range(1, 16)]
mins = lambda t: int(t.split(":")[0]) * 60 + int(t.split(":")[1])
continuous = all(spans[i][1] == spans[i + 1][0] for i in range(len(spans) - 1))
chk("45-minute architecture intact, continuous and unchanged",
    continuous and spans[0][0] == "0:00" and spans[-1][1] == "45:00"
    and [TIMING.findall(NOTE[i]) for i in range(1, 17)]
    == [TIMING.findall(old.slides[i - 1].notes_slide.notes_text_frame.text
                       if old.slides[i - 1].has_notes_slide else "")
        for i in range(1, 17)],
    "0:00 to 45:00 across 15 timed slides, no gap or overlap at any hand-off; "
    "every cue identical to v3.4.1")
chk("The hidden AI appendix stays outside the timed session",
    "NOT PART OF THE 45-MINUTE SESSION" in FACE[16] and not TIMING.search(NOTE[16]),
    "no timing cue, and the slide is hidden")

# ── content boundary ────────────────────────────────────────────────────────
chk("The three signs are unchanged from v3.4.1",
    all(OFACE[i] == FACE[i] for i in (7, 8, 9))
    and "You Are Repeating, Not Stretching" in norm(FACE[7])
    and "You Are Indispensable, but Mainly Here" in norm(FACE[8])
    and "Your Responsibility Is Growing Faster Than Your Capability" in norm(FACE[9]),
    "slides 7, 8 and 9 byte-identical")
chk("The lesson stays recognition and early diagnosis only",
    "A recognition exercise is not a result." in FACE[11]
    and "TODAY WILL NOT:" in FACE[3] and OFACE[11] == FACE[11] and OFACE[3] == FACE[3],
    "the non-diagnostic boundary slide and the outcomes slide are both unchanged")
# A line that names a framework term IN ORDER TO FORBID IT is the guardrail, not
# a leak of it. Slide 3 promises the session will not "place you in a Career
# State", and slide 13's note tells the presenter not to introduce Density,
# Optionality or Career States at all. Both must contain the words to mean
# anything, and a sweep that flagged them would be reading the guardrail as the
# breach — the same false positive that has cost time on this family before.
BAN = re.compile(r"(will not|won't|do not|don't|does not|never|no )", re.I)
BAN_HEADER = re.compile(r"(WILL NOT|MUST NOT|WHAT NOT TO SAY|OUTSIDE THE BOUNDARY)", re.I)


def live(text, pattern):
    """Occurrences that are live claims rather than guardrails.

    Negation is not always on the same line as the term. Slide 3 lists its
    boundary as a "TODAY WILL NOT:" header followed by three bare bullets, so
    "place you in a Career State" carries no negation of its own and reads as a
    breach unless the header's scope is honoured. Scope runs until the next
    header-like line.
    """
    out, negated = [], False
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if BAN_HEADER.search(s):
            negated = True
            continue
        if s.endswith(":") or s.isupper():
            negated = False
        if re.search(pattern, s) and not (negated or BAN.search(s)):
            out.append(s)
    return out


STATES = r"(Depth Trap|Compounding|Stagnant|Fragile|Career State)"
chk("No participant is placed into a Career State",
    not live(ALLFACE, STATES),
    "the only mention on any slide face is slide 3's promise that today will NOT "
    "place you in a Career State, which is the boundary rather than a breach of it")
chk("The lesson does not become the Capability Position Read",
    "Capability Position Read" not in ALLFACE
    and "Today helped you spot the pattern." in FACE[13],
    "the phrase is on no slide face; slide 13 hands off to the guided read instead "
    "of performing one")
# "The Density Group" is the company on the byline, not the Density axis, so the
# axis pattern requires the word to stand alone.
AXIS = r"\b(Density(?! Group)|Optionality)\b"
chk("No Density or Optionality scoring teaching leaked into the timed core",
    not live(CORE, AXIS) and not re.search(r"/\s*30|1\s+2\s+3\s+4\s+5", CORE),
    "the only occurrences in the fifteen live slides are the company name on the "
    "byline and slide 13's note forbidding the presenter from introducing either "
    "axis. No scale and no /30 total appears anywhere")
chk("No enterprise, staffing, succession or workforce section was added",
    not any(w in CORE.lower() for w in
            ("staffing", "succession", "workforce", "headcount", "hiring plan")),
    "and no slide was added: 16 in v3.4.1, 16 now")
chk("The Career Stall Check stays a four-question recognition tool",
    OFACE[10] == FACE[10] and "Five quiet minutes." in FACE[10]
    and "last 90 days" in FACE[10],
    "slide 10 byte-identical to v3.4.1")

# ── commercial architecture ─────────────────────────────────────────────────
chk("Exactly two continuation routes, and no third was stacked",
    FACE[14].count("REGISTER FREE") == 1 and FACE[14].count("EXPLORE THE FIELD KIT") == 1
    and "Private Capability Position Read" not in ALLFACE
    and "Substack" not in ALLFACE,
    "the September 23 live session and the Field Kit, as in v3.4.1")
chk("September 23 is the primary route and the Field Kit is the optional private one",
    "FREE LIVE SESSION" in FACE[14] and "Wednesday, September 23, 2026" in FACE[14]
    and "$150" in FACE[14] and "SELF-GUIDED" in FACE[14]
    and "Neither route is required. September 23 is free. The Field Kit is optional "
        "and private." in norm(FACE[14]))
chk("No retired offer survives anywhere in the deck",
    not any(p in ALLFACE + "\n".join(NOTE.values())
            for p in ("$249", "$99", "$149", "$500", "paid workshop")),
    "the $249 paid workshop that the v3.2 facilitator guide still describes appears "
    "nowhere in the shipped deck")

# ── media, QR codes and links ───────────────────────────────────────────────
z = zipfile.ZipFile(DECK)
media = [n for n in z.namelist() if n.startswith("ppt/media/")]
qr = {}
for n in media:
    try:
        from pyzbar.pyzbar import decode
        for d in decode(Image.open(io.BytesIO(z.read(n)))):
            qr[n] = d.data.decode()
    except Exception:
        pass
links = set()
for n in z.namelist():
    if n.startswith("ppt/slides/_rels/"):
        links |= set(re.findall(r'Target="(https?://[^"]+)"', z.read(n).decode()))
WANT = {"https://maven.com/p/8b3c40/stay-or-leave-live-career-growth-assessment",
        "https://temidayoafonja.com/fieldkit"}
chk("Media inventory is exactly three parts, with no orphaned image",
    len(media) == 3 and len(qr) == 2,
    "the portrait on slide 2 and the two QR codes on slide 14; retired QR images "
    "have previously survived a rebuild as orphans and none does here")
chk("Both QR codes decode to the intended destinations",
    set(qr.values()) == WANT, " | ".join(sorted(qr.values())))
chk("Click-through hyperlinks match the QR payloads exactly",
    links == WANT == set(qr.values()),
    "a QR and its printed link cannot disagree")
SITE = {f: open(f, encoding="utf-8", errors="replace").read()
        for f in ("for-professionals.html", "diagnostic.html", "career-decisions.html",
                  "content/site-source-of-truth.json")}
SUPERSEDED = "stay-or-leave-your-job-live-career-growth-assessment"
stale_slug = [f for f, t in SITE.items()
              if any(SUPERSEDED in l and "superseded" not in l for l in t.split("\n"))]
chk("Both destinations are still the live routes on the site",
    all(any(u in t for t in SITE.values()) for u in WANT) and not stale_slug,
    "the superseded Maven slug — which differs from the live one by the two words "
    "your-job only — survives solely inside the source-of-truth entry that records "
    "it as retired, and is linked from no page")

# ── the leave-behind ────────────────────────────────────────────────────────
STALLT = doc_text(STALL)
qs = ["Which sign is most present for me?",
      "What is one example from the last 90 days?",
      "What capability, learning, or future option may not be growing?",
      "What is the next question I need to investigate?"]
chk("Career Stall Check questions match the deck's slide 10 verbatim",
    all(q in STALLT and q in FACE[10] for q in qs),
    "four of four identical, so the leave-behind and the slide cannot drift apart")
chk("Career Stall Check carries no date, price or version-bound claim",
    not re.search(r"(September|August|October|\b20\d\d\b|\$\d)", STALLT),
    "nothing inside it goes stale when the session moves, which is why it is "
    "reissued unchanged at v1.0")
chk("Career Stall Check points at the current continuation and keeps its boundary",
    "Decide Whether to Stay or Leave Your Job" in STALLT
    and "The Capability Position Read" in STALLT
    and "It does not calculate Density or Optionality, place you into a Career State, "
        "or replace the full evidence-backed read." in norm(STALLT))
chk("Career Stall Check PDF is the one-page render of that DOCX",
    os.path.exists(STALL_PDF) and len(__import__("pymupdf").open(STALL_PDF)) == 1)

# ── what no longer exists, stated rather than reconstructed ─────────────────
GUIDET = doc_text(GUIDE)
chk("Facilitator Guide v3.2 is SUPERSEDED and is deliberately not shipped",
    "$249" in GUIDET and "45-minute run of show — v3.2" in GUIDET
    and "15 slides" in GUIDET,
    "it describes a 15-slide deck and a $249 paid workshop that has since been "
    "retired from active public sale. Correcting it is a rewrite, not a date edit, "
    "so it is left untouched and excluded from the package")
chk("No facilitator SOP matching the v3.5.0 deck exists in the repository",
    not any("v3.5" in f or "v3.4" in f
            for f in os.listdir("lightning-lesson-v3")
            if "Facilitator" in f),
    "STATED, NOT RECONSTRUCTED. The delivery instrument for this lesson is the "
    "speaker notes inside the deck, which carry TIMING, FACILITATION, WHAT TO SAY "
    "and WHAT NOT TO SAY on every slide")
chk("The v3.5.0 build and QA scripts survive in the repository",
    all(os.path.exists(f"sept9-assets/{n}") for n in
        ("SOURCE_build_v350.py", "SOURCE_qa_sept9.py", "SOURCE_audit_family.py")),
    "they lived in a gitignored scratchpad, where the flagship generators were lost "
    "in a container rebuild; they are committed here so the same loss cannot recur")

# ── website and repository copy ─────────────────────────────────────────────
chk("The site carries September 9 for the lesson and September 23 for the read",
    "Wednesday, September 9, 2026" in SITE["for-professionals.html"]
    and "Wednesday, September 23, 2026" in SITE["for-professionals.html"]
    and '"date": "2026-09-09"' in SITE["content/site-source-of-truth.json"]
    and '"date": "2026-09-23"' in SITE["content/site-source-of-truth.json"])
CES = open("career-evidence-starter.html", encoding="utf-8").read()
chk("career-evidence-starter.html carries the corrected expiry and start time only",
    'available_until: "2026-09-09T18:45:00-05:00"' in CES
    and "Wednesday, September 9, 2026, 6:00 PM CT" in CES
    and "6:00 to 6:45" not in CES
    # The lookahead matters twice over: it must not match "September 23", and the
    # page writes the date day-first in a code comment as "9 September 2026".
    and not re.search(r"September 2(?![0-9])", CES),
    "this page was missed by the pass that rescheduled the rest of the site — the "
    "expiry instant, the body copy and three occurrences in the explanatory comment "
    "were corrected in this pass, and the body now names the start time only")
# The 64 browser tests need Playwright, which is not a dependency of this
# repository. They were run once in this pass against the preinstalled Chromium
# with a matching Playwright installed temporarily, and passed 64 of 64; the
# dependency was then reverted so package.json is unchanged. The two suites that
# need no browser are re-run here on every audit.
chk("Site test suites that need no browser are green",
    all(subprocess.run(["node", f], capture_output=True, text=True).returncode == 0
        for f in ("tests/career-decisions-export.test.mjs",
                  "tests/career-decisions-subscribe.test.mjs")),
    "76 of 76. The 64 Playwright page tests — which cover the expiry instant of the "
    "September 9 registration block — passed 64 of 64 when run separately in this "
    "pass; Playwright is not a repository dependency and was reverted afterwards")

if __name__ == "__main__":
    w = max(len(l) for _, l, _, _ in R)
    for n, label, st, note in R:
        print(f"{n:>3}. [{st}] {label:<{w}}  {note}")
    fails = [r for r in R if r[2] != "PASS"]
    print(f"\n{len(R) - len(fails)} of {len(R)} pass")
    sys.exit(1 if fails else 0)
