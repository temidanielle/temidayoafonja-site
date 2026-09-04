# -*- coding: utf-8 -*-
"""September 9 Lightning Lesson — v3.4.1 FINAL -> v3.5.0 FINAL.

A controlled revision, not a rebuild. Two things change:

  DATES     The session moves from Wednesday, September 2 to Wednesday,
            September 9, 2026, and the continuation session it points at moves
            from September 16 to Wednesday, September 23, 2026.

  NOTES     Two short speaker-note additions, no new slides and no new offer:
            slide 6 gains a line on relearning in a new context; slide 8 gains
            a spoken follow-up about whether an outsider would have enough
            evidence to trust the experience.

Method follows the v3.3.x/v3.4.x lineage: open the prior approved PPTX and
rewrite specific runs under script control. The generators for this deck lived
in a gitignored scratchpad and did not survive the container rebuild, but this
deck has always been maintained by scripted patching of the approved file, so
the method is unchanged rather than a fallback.

Every teaching line is left alone. The three signs, the Career Stall Check, the
last-90-days discipline, the portability distinction and "Starting as a learner
is not the same as starting from zero." are untouched.
"""
import hashlib, os, re, shutil
from pptx import Presentation

SRC = ("deck-v340-assets/"
       "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_v3.4.1_FINAL.pptx")
OUT = "scratchpad/sept9/out"
os.makedirs(OUT, exist_ok=True)
DST = os.path.join(OUT, "How_to_Tell_If_Your_Career_Is_Stalling_Lightning_Lesson_"
                        "v3.5.0_FINAL.pptx")

# ── date substitutions, applied to faces and notes alike ────────────────────
# Ordered longest-first so "September 16, 2026" is rewritten before the bare
# "September 16" pattern can match part of it.
DATES = [
    # The slide-13 panel header is typeset in capitals, so it needs its own
    # entry: the mixed-case patterns below cannot match "SEPTEMBER 16".
    ("SEPTEMBER 16", "SEPTEMBER 23"),
    ("SEPTEMBER 2,", "SEPTEMBER 9,"),
    ("Wednesday, September 16, 2026", "Wednesday, September 23, 2026"),
    ("September 16, 2026",            "September 23, 2026"),
    ("September 16",                  "September 23"),
    ("September 2, 2026",             "September 9, 2026"),
    # The lesson's own date appears on slide 13 as a bare "September 2". The
    # negative lookahead stops it matching the "2" of September 23 after the
    # continuation date has already been rewritten above.
    (re.compile(r"September 2(?![0-9])"), "September 9"),
]

# ── the two note additions ──────────────────────────────────────────────────
# Slide 6's WHAT TO SAY already said "some have to be relearned in the new
# context", so the relearning point is extended in place rather than appended as
# a second block. Appending would have said the same thing twice and crowded a
# ninety-second beat.
RELEARN_OLD = ("Some things travel with you and some have to be relearned in the new "
               "context.")
RELEARN_NEW = ("Some things travel with you and some have to be relearned. Some of this "
               "may come back quickly, but it still has to be learned in the new "
               "context. That is different from having nothing to carry.")

# Slide 8 has no equivalent line, so the follow-up is appended to WHAT TO SAY,
# immediately after the instruction to read the test question.
TRUST_OLD = ("Read the test question, then give the room 30-45 seconds for the quick "
             "check.")
TRUST_NEW = ("Read the test question, then ask the follow-up: and would they have enough "
             "evidence to trust that the experience is useful in their context? Then "
             "give the room 30-45 seconds for the quick check. Being explainable is not "
             "the same as being credible to someone outside your organisation — do not "
             "imply that translating your value guarantees an employer will accept it.")


def patch_runs(container, fn):
    """Rewrite run text in place, so every run keeps its own formatting."""
    hits = 0
    for para in container.paragraphs:
        for run in para.runs:
            new = fn(run.text)
            if new != run.text:
                run.text = new; hits += 1
    return hits


def apply_dates(text):
    for pat, rep in DATES:
        text = pat.sub(rep, text) if hasattr(pat, "sub") else text.replace(pat, rep)
    return text


def build():
    shutil.copyfile(SRC, DST)
    p = Presentation(DST)
    sl = list(p.slides)

    face_hits = note_hits = 0
    for s in sl:
        for sh in s.shapes:
            if sh.has_text_frame:
                face_hits += patch_runs(sh.text_frame, apply_dates)
        if s.has_notes_slide:
            note_hits += patch_runs(s.notes_slide.notes_text_frame, apply_dates)

    # ── slide 6: relearning ─────────────────────────────────────────────────
    n6 = sl[5].notes_slide.notes_text_frame
    assert "A pivot does not mean everything transfers." in \
        "\n".join(sh.text_frame.text for sh in sl[5].shapes if sh.has_text_frame), \
        "slide 6 is not the career-portability slide"
    assert RELEARN_OLD in n6.text, "slide 6: the relearning sentence was not found"
    n6.text = n6.text.replace(RELEARN_OLD, RELEARN_NEW)

    # ── slide 8: employer legibility and trust ──────────────────────────────
    n8 = sl[7].notes_slide.notes_text_frame
    assert "Could someone outside your company explain your value" in \
        "\n".join(sh.text_frame.text for sh in sl[7].shapes if sh.has_text_frame), \
        "slide 8 is not the sign-2 slide carrying the outsider test"
    assert TRUST_OLD in n8.text, "slide 8: the test-question instruction was not found"
    assert "enough evidence to trust" not in n8.text, "slide 8 already carries the follow-up"
    n8.text = n8.text.replace(TRUST_OLD, TRUST_NEW)

    p.save(DST)

    # ── prove the revision was controlled ───────────────────────────────────
    before, after = Presentation(SRC), Presentation(DST)
    bf = ["\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
          for s in before.slides]
    af = ["\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
          for s in after.slides]
    bn = [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
          for s in before.slides]
    an = [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
          for s in after.slides]

    assert len(af) == len(bf) == 16, "slide count moved"
    changed_faces = [i + 1 for i in range(16) if bf[i] != af[i]]
    changed_notes = [i + 1 for i in range(16) if bn[i] != an[i]]

    # Rather than hardcode which slides may change, prove that each difference is
    # one the brief authorised. A face may differ only by a date rewrite. A note
    # may differ only by a date rewrite, or by being slide 6 or 8 whose text
    # matches the date-rewritten original with exactly the sanctioned edit applied.
    for i in changed_faces:
        assert apply_dates(bf[i - 1]) == af[i - 1], \
            f"slide {i}: a face changed by something other than a date"
    for i in changed_notes:
        expect = apply_dates(bn[i - 1])
        if i == 6:
            expect = expect.replace(apply_dates(RELEARN_OLD), apply_dates(RELEARN_NEW))
        elif i == 8:
            expect = expect.replace(apply_dates(TRUST_OLD), apply_dates(TRUST_NEW))
        assert expect == an[i - 1], \
            f"slide {i}: a note changed by something other than a date or its sanctioned edit"
    assert 6 in changed_notes and 8 in changed_notes, "the two note additions did not land"

    print("built", os.path.basename(DST))
    print(f"  date rewrites: {face_hits} on faces, {note_hits} in notes")
    print(f"  faces changed: {changed_faces}  notes changed: {changed_notes}")
    print("  sha256", hashlib.sha256(open(DST, "rb").read()).hexdigest())
    return DST


if __name__ == "__main__":
    build()
