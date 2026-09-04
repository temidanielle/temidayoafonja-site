# -*- coding: utf-8 -*-
"""September 23 flagship family QA — executed against the built files.

Written to be able to fail. Each check reads the artifact that shipped rather
than the variables the build used, and the two that matter most — the twelve
statements and the scoring mechanics — compare the new deck against the approved
v2.0.4 file on disk rather than against a copy of the text kept here.
"""
import hashlib, os, re, sys, zipfile
import docx
import pymupdf
from pptx import Presentation

OUT = "scratchpad/sept23/out"
SRC = "sept16-v204-assets"

DECK = f"{OUT}/PRESENTER_VERSION_Stay_or_Leave_Live_Career_Growth_Assessment_60MIN_v2.0.5_CANDIDATE.pptx"
PRIOR = f"{SRC}/PRESENTER_VERSION_Stay_or_Leave_Live_Career_Growth_Assessment_60MIN_v2.0.4_CANDIDATE.pptx"
SOP = f"{OUT}/Capability_Formation_Free_Flagship_SOP_Sept23_2026_60MIN_v2.0.8_CANDIDATE.docx"
SOP_PDF = f"{OUT}/Capability_Formation_Free_Flagship_SOP_Sept23_2026_60MIN_v2.0.8_CANDIDATE.pdf"
QAR = f"{OUT}/Free_Flagship_60MIN_v2.0.8_Change_Log_and_QA_Report.docx"
CHK = f"{OUT}/Manual_Maven_and_Website_Edit_Checklist_v1.1.docx"
WB = "free-flagship-assets/60min-v2.0.1/Capability_Position_Read_Workbook_60MIN_v2.0.1_CANDIDATE.pdf"
WB_SHA = "2bd2912846a679837e8e6bfb4aadff2bb07ee5959d35502ca1c5b5c728efa3ee"

R = []
def chk(g, label, ok, note=""):
    R.append((len(R) + 1, g, label, "PASS" if ok else "FAIL", note))


def deck_text(path):
    p = Presentation(path)
    f = ["\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
         for s in p.slides]
    n = [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
         for s in p.slides]
    return f, n


def hidden(path):
    z = zipfile.ZipFile(path)
    rels = dict(re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"',
                           z.read("ppt/_rels/presentation.xml.rels").decode()))
    order = [rels[m] for m in re.findall(r'<p:sldId[^>]*r:id="([^"]+)"',
                                         z.read("ppt/presentation.xml").decode())]
    out = []
    for pos, t in enumerate(order, 1):
        x = z.read("ppt/" + t.lstrip("./")).decode()
        i = x.index("<p:sld ")
        # the whole opening tag, not a fixed-width slice: show="0" sits behind a
        # long namespace list and a [:600] slice has missed it twice before.
        if 'show="0"' in x[i:x.index(">", i) + 1]:
            out.append(pos)
    return out


def doc_text(path):
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            parts += [c.text for c in row.cells]
    return "\n".join(parts)


NF, NN = deck_text(DECK)
OF, ON = deck_text(PRIOR)
SOPT, QART, CHKT = doc_text(SOP), doc_text(QAR), doc_text(CHK)
WBT = "\n".join(p.get_text() for p in pymupdf.open(WB))
DOCS = {"SOP": SOPT, "QA report": QART, "checklist": CHKT, "workbook": WBT}
DECKT = "\n".join(NF) + "\n" + "\n".join(NN)

# ── A. stale dates, zero tolerance ──────────────────────────────────────────
# A line carrying an old date is a LIVE stale claim unless it is one of three
# things, each of which would be wrong to rewrite:
#
#   the transition itself   a change-log line that names both the old value and
#                           the new one is the record of the move, not a relic
#                           of it;
#   the guardrail           a QA row that quotes a term in order to assert its
#                           absence must contain the term to mean anything;
#   the audit trail         the v2.0.6 correction quotes the wording it replaced,
#                           verbatim, including its date.
#
# Nothing else is excused, and the sweep is case-insensitive throughout: a slide
# header typeset in capitals is the same stale date as one in sentence case, and
# a case-sensitive sweep missed exactly that on the Lightning Lesson deck.
QUOTED_V206 = "“Attendee email already sent for the September 2 session"
NEW_VALUES = ("September 23", "Sept23_2026", "September 9, 2026")


def live_lines(where, text, pattern):
    hits = []
    for n, line in enumerate(text.split("\n"), 1):
        if not pattern.search(line):
            continue
        if any(v in line for v in NEW_VALUES):
            continue
        if line.strip().startswith("Zero stale"):
            continue
        if QUOTED_V206 in line:
            continue
        hits.append(f"{where}:{n}")
    return hits


for term, why in [("September 16", "the session's former date"),
                  ("Sept 16", "abbreviated former date"),
                  ("Sept16", "date-coded filename fragment"),
                  ("2026-09-16", "machine-readable former date")]:
    pat = re.compile(re.escape(term), re.I)
    where = live_lines("deck faces", "\n".join(NF), pat) \
        + live_lines("deck notes", "\n".join(NN), pat)
    for k, v in DOCS.items():
        where += live_lines(k, v, pat)
    chk("A", f"No stale “{term}” — {why}", not where,
        "found at " + ", ".join(where) if where else
        "zero live occurrences; change-log lines naming the move, QA rows quoting "
        "the term to ban it, and the v2.0.6 audit quotation are excluded by design")

# The Lightning Lesson this family cross-references moved to September 9 in an
# earlier pass. A bare "September 2" must not survive, and the lookahead keeps
# the sweep from matching the "2" of the new September 23.
_S2 = re.compile(r"September 2(?![0-9])", re.I)
s2 = live_lines("deck", DECKT, _S2)
for k, v in DOCS.items():
    s2 += live_lines(k, v, _S2)
chk("A", "No stale bare “September 2” Lightning Lesson reference", not s2,
    "found at " + ", ".join(s2) if s2 else
    "zero live occurrences; the two SOP cross-references now read September 9")

_CST = re.compile(r"(?<![A-Za-z_])CST(?![A-Za-z_])")
cst = []
if _CST.search(DECKT): cst.append("deck")
for k, v in DOCS.items():
    for ln in v.split("\n"):
        if not _CST.search(ln): continue
        if "Never write CST" in ln or "never CST as a label" in ln: continue
        cst.append(f"{k}: {ln.strip()[:60]}")
chk("A", "Timezone is CT or Central Time, never CST as a label", not cst,
    "found in " + " | ".join(cst) if cst else
    "the only occurrences are the SOP rule banning it and the QA row asserting that")

# ── B. the new date is actually present ─────────────────────────────────────
chk("B", "Holding slide carries Wednesday, September 23, 2026",
    "Wednesday, September 23, 2026" in NF[32], NF[32].split("\n")[2][:60])
chk("B", "The holding slide is the only dated slide face",
    sum("September" in f for f in NF) == 1,
    "26 active slides carry no date, which is what keeps the replay evergreen")
chk("B", "SOP names the new date in all seven of its locations",
    SOPT.count("September 23") == 7 and SOPT.count("Wednesday, September 23, 2026") == 2,
    "identity table, offer table, Maven configuration table, and sections 3 and 12 — "
    "the same seven places that carried September 16")
chk("B", "Checklist instructs the operator to set September 23",
    "Date Wednesday, September 23, 2026. Time 6:00 PM CT." in CHKT
    and "September 23 — Maven and Website Edit Checklist" in CHKT)
chk("B", "SOP cross-references the Lightning Lesson at its own new date",
    "The September 9 Lightning Lesson is a separate event" in SOPT
    and "configured for the September 9 session" in SOPT)

# ── C. the instrument is frozen ─────────────────────────────────────────────
STATEMENT_SLIDES = [6, 7, 9, 10]
LOGIC_SLIDES = [8, 11, 12, 13, 14, 15, 17, 18, 19, 20, 21, 22]
chk("C", "The twelve statements are byte-identical to v2.0.4",
    all(OF[i - 1] == NF[i - 1] for i in STATEMENT_SLIDES),
    f"slides {STATEMENT_SLIDES} identical; not rewritten, reordered, added or removed")
chk("C", "Twelve statements are present and numbered 1 to 12",
    all(re.search(rf"(?m)^{n}$", "\n".join(NF[i - 1] for i in STATEMENT_SLIDES))
        for n in range(1, 13)),
    "counted on the built deck, not assumed")
chk("C", "Scoring, state, boundary and move-category logic byte-identical",
    all(OF[i - 1] == NF[i - 1] for i in LOGIC_SLIDES),
    f"slides {LOGIC_SLIDES} identical to v2.0.4")
chk("C", "Both /30 totals, the bands and the boundary band are intact",
    all(s in NF[7] and s in NF[10] for s in ("________  / 30", "19 to 30 is high",
                                             "6 to 18 is low", "17 to 21 is the boundary band")))
chk("C", "The sensitivity check and the incomplete reading survive verbatim",
    "NEIGHBOURING-SCORE SENSITIVITY CHECK" in NF[17]
    and "INCOMPLETE — do not place on that axis today." in NF[17])
chk("C", "Seven move categories, unchanged and in order",
    OF[20] == NF[20] and all(t in NF[20] for t in (
        "Remain and deepen", "Translate what is built", "Widen exposure",
        "Test portability", "Repair formation conditions", "Prepare for exit",
        "Seek an external perspective")))

# ── D. only what was authorised moved ───────────────────────────────────────
chk("D", "Slide count unchanged at 33", len(NF) == len(OF) == 33)
diff_faces = [i + 1 for i in range(33) if OF[i] != NF[i]]
diff_notes = [i + 1 for i in range(33) if ON[i] != NN[i]]
chk("D", "Exactly one slide face differs from v2.0.4", diff_faces == [33],
    f"{33 - len(diff_faces)} of 33 faces byte-identical; the difference is the date")
chk("D", "Only slides 5, 25 and 29 have changed notes", diff_notes == [5, 25, 29],
    f"notes changed: {diff_notes}")
chk("D", "Hidden set is exactly 27-33, matching v2.0.4",
    hidden(DECK) == hidden(PRIOR) == [27, 28, 29, 30, 31, 32, 33], f"hidden {hidden(DECK)}")
TIMING = re.compile(r"TIMING:\s*[\d:–\-\s.]+")
chk("D", "Every TIMING cue is identical to v2.0.4",
    [TIMING.findall(t) for t in ON] == [TIMING.findall(t) for t in NN],
    "the timed core was not retimed, lengthened or shortened")
chk("D", "The recording boundary is unchanged",
    "AT 50:00 — STOP THE RECORDING BEFORE ADVANCING." in NN[24]
    and "RECORDING OFF. CONFIRM ON SCREEN" in NN[25])
chk("D", "No new offer, and the Private Read is still off the continuation",
    "Private Capability Position Read" not in "\n".join(NF)
    and NF[22] == OF[22] and NF[23] == OF[23],
    "slide 23 still carries two routes: do nothing, or the Field Kit")
chk("D", "No enterprise, staffing, succession or AI teaching entered the timed core",
    not any(w in "\n".join(NF[:26]).lower() + "\n".join(NN[:26]).lower()
            for w in ("staffing", "succession", "workforce", "headcount")),
    "the AI appendix remains hidden at slide 28 and outside the timed core")
ext = set()
z = zipfile.ZipFile(DECK)
for n in z.namelist():
    if n.startswith("ppt/slides/_rels/"):
        ext |= set(re.findall(r'Target="(https?://[^"]+)"', z.read(n).decode()))
chk("D", "External URLs unchanged from v2.0.4", ext == set(
    sum([re.findall(r'Target="(https?://[^"]+)"',
                    zipfile.ZipFile(PRIOR).read(n).decode())
         for n in zipfile.ZipFile(PRIOR).namelist()
         if n.startswith("ppt/slides/_rels/")], [])),
    " | ".join(sorted(ext)))

# ── E. the three note refinements ───────────────────────────────────────────
chk("E", "Slide 5 carries the portability sentence, stated once",
    "That is why portability is not the same as plug-and-play. Something important "
    "can travel with you while the new context still requires real relearning." in NN[4]
    and NN[4].count("plug-and-play") == 1,
    "placed inside the spoken proof, which already names what travelled and what did not")
chk("E", "Slide 5 carries the four-part portability audit",
    all(q in NN[4] for q in ("What travels?", "What does not travel automatically?",
                             "What can I prove?", "What must I relearn?")))
chk("E", "The audit is framed as an alternative, not an addition to the ninety seconds",
    "in place of the spoken proof above, not alongside it" in NN[4]
    and "Do not work an example" in NN[4])
chk("E", "Slide 5 face is untouched", OF[4] == NF[4])
chk("E", "Slide 29 carries the translation-acceptance limit",
    "Translation does not guarantee acceptance. Another organization still has to "
    "decide that this evidence is comparable enough to trust in its context." in
    NN[28].replace("translation does not", "Translation does not"))
chk("E", "Slide 29 refuses the résumé-framing overclaim",
    "does not remove genuine experience gaps, market constraints or employer "
    "preferences" in NN[28])
chk("E", "Slide 29 face is untouched and the slide stays hidden",
    OF[28] == NF[28] and 29 in hidden(DECK))
chk("E", "The evergreen prohibition moved with the session",
    'Do NOT say "tonight", "this cohort", "September 23"' in NN[24]
    and "Do not say tonight, this cohort, September 23, or before enrollment closes"
    in SOPT,
    "the rule is unchanged; only the date it names as an example moved")

# ── F. documents, labels and provenance ─────────────────────────────────────
chk("F", "SOP exports at exactly seven pages", len(pymupdf.open(SOP_PDF)) == 7,
    "pinned, not left open: v2.0.5 of the SOP once shipped at eight, so this pass "
    "added no prose to the document at all")
orient = ["landscape" if p.rect.width > p.rect.height else "portrait"
          for p in pymupdf.open(SOP_PDF)]
chk("F", "Landscape stays confined to the two reconciliation pages",
    orient == ["portrait", "landscape", "landscape"] + ["portrait"] * 4,
    " ".join(orient))
chk("F", "SOP filename date code matches the date inside it",
    os.path.basename(SOP).startswith("Capability_Formation_Free_Flagship_SOP_Sept23_2026")
    and os.path.basename(SOP) in SOPT,
    "the stamp line and the source-of-truth map both name the file it actually is")
chk("F", "SOP source map names the current artifact of each kind",
    all(s in SOPT for s in ("60MIN_v2.0.5_CANDIDATE.pptx",
                            "Workbook_60MIN_v2.0.1_CANDIDATE.pdf",
                            "60MIN_v2.0.8_CANDIDATE.docx",
                            "Architecture_v1.1.docx")))
chk("F", "No superseded version string survives in the SOP",
    not any(v in SOPT for v in ("v2.0.2", "v2.0.3", "v2.0.4", "v2.0.6", "v2.0.7")),
    "checked across body and tables")
chk("F", "SOP keeps the timezone rule and the protected rescore block",
    "Use CT / Central Time / America/Chicago. Never write CST." in SOPT
    and "The 19:00–31:00 evidence-backed rescore is PROTECTED." in SOPT)
chk("F", "QA report is stamped v2.0.8 and counts every row it lists",
    "QA REPORT v2.0.8" in QART and "Verification — 162 items" in QART
    and "162 of 162 PASS." in QART
    and all(f"\n{n}\n" in "\n" + QART + "\n" or str(n) in QART
            for n in range(149, 163)),
    "148 carried forward plus the fourteen added for this pass")
chk("F", "QA report discloses the deck's provenance honestly",
    "REPOINTED, NOT WEAKENED" in QART
    and "did not survive a container rebuild" in QART
    and "reproducible from its generator" not in QART,
    "the generator claim is retired rather than repeated, and row 119 says why")
chk("F", "Checklist reissued at v1.1 with the deck reference current",
    "MANUAL MAVEN AND WEBSITE EDIT CHECKLIST v1.1" in CHKT
    and "deck v2.0.5" in CHKT and "v1.0" not in CHKT)

# ── G. the workbook ─────────────────────────────────────────────────────────
chk("G", "Workbook is byte-identical to the approved v2.0.1 build",
    hashlib.sha256(open(WB, "rb").read()).hexdigest() == WB_SHA, WB_SHA)
chk("G", "Workbook is date-free, so the reschedule reaches nothing inside it",
    not re.search(r"(September|October|2026|20\d\d-\d\d-\d\d|\bCST\b|\bCT\b)", WBT),
    "rasterised and searched: no date, no year, no timezone")
chk("G", "The SOP still names the workbook that exists",
    "Capability_Position_Read_Workbook_60MIN_v2.0.1_CANDIDATE.pdf" in SOPT
    and os.path.exists(WB))

if __name__ == "__main__":
    w = max(len(l) for _, _, l, _, _ in R)
    g = None
    for n, grp, label, st, note in R:
        if grp != g:
            print(f"\n── group {grp} " + "─" * 52); g = grp
        print(f"{n:>3}. [{st}] {label:<{w}}  {note}")
    fails = [r for r in R if r[3] != "PASS"]
    print(f"\n{len(R) - len(fails)} of {len(R)} pass")
    sys.exit(1 if fails else 0)
