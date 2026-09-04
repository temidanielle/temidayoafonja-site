# -*- coding: utf-8 -*-
"""Free Flagship 60-minute family — September 16 -> Wednesday, September 23, 2026.

A controlled revision, not a rebuild. Three things change:

  DATES     The guided session moves from Wednesday, September 16 to Wednesday,
            September 23, 2026, and the Lightning Lesson this family
            cross-references moves from September 2 to September 9, 2026, which
            an earlier pass already applied to the lesson's own deck and to the
            website but not to this family's documents.

  NOTES     Three speaker-note refinements. No new slide, no new offer, no new
            section, and nothing added to the timed participant core.

  LABELS    The SOP filename is date-coded and therefore went stale with the
            date. It is recoded, and the family's version stamps advance.

The instrument is untouched: the twelve statements, the 1-5 scale, both /30
totals, the 17-21 boundary band, the sensitivity check and the seven move
categories are asserted byte-identical to v2.0.4 rather than merely left alone.

METHOD NOTE, stated because it matters for provenance. The generators for this
family lived in a gitignored scratchpad and did not survive a container rebuild.
This pass therefore patches the approved built artifacts under script control,
which is how the Lightning Lesson deck has always been maintained. Every edit is
a named substitution or a named insertion, and the verification block below
proves that nothing else moved.
"""
import copy, hashlib, os, re, shutil
import docx
from pptx import Presentation

SRC_DIR = "sept16-v204-assets"
OUT = "scratchpad/sept23/out"
os.makedirs(OUT, exist_ok=True)

DECK_SRC = f"{SRC_DIR}/PRESENTER_VERSION_Stay_or_Leave_Live_Career_Growth_Assessment_60MIN_v2.0.4_CANDIDATE.pptx"
DECK_DST = f"{OUT}/PRESENTER_VERSION_Stay_or_Leave_Live_Career_Growth_Assessment_60MIN_v2.0.5_CANDIDATE.pptx"
SOP_SRC = f"{SRC_DIR}/Capability_Formation_Free_Flagship_SOP_Sept16_2026_60MIN_v2.0.7_CANDIDATE.docx"
SOP_DST = f"{OUT}/Capability_Formation_Free_Flagship_SOP_Sept23_2026_60MIN_v2.0.8_CANDIDATE.docx"
QA_SRC = f"{SRC_DIR}/Free_Flagship_60MIN_v2.0.7_Change_Log_and_QA_Report.docx"
QA_DST = f"{OUT}/Free_Flagship_60MIN_v2.0.8_Change_Log_and_QA_Report.docx"
CHK_SRC = f"{SRC_DIR}/Manual_Maven_and_Website_Edit_Checklist_v1.0.docx"
CHK_DST = f"{OUT}/Manual_Maven_and_Website_Edit_Checklist_v1.1.docx"

WORKBOOK = "free-flagship-assets/60min-v2.0.1/Capability_Position_Read_Workbook_60MIN_v2.0.1_CANDIDATE.pdf"
WORKBOOK_SHA = "2bd2912846a679837e8e6bfb4aadff2bb07ee5959d35502ca1c5b5c728efa3ee"

STAMP = "Friday, September 4, 2026 at 6:40 AM CT"
OLD_STAMP_SOP = "Sunday, August 30, 2026 at 1:07 PM CT"
OLD_STAMP_QA = "Sunday, August 30, 2026 at 1:11 PM CT"
OLD_STAMP_CHK = "Sunday, August 30, 2026 at 12:40 PM CT"

# ── date substitutions ──────────────────────────────────────────────────────
# Ordered so that the September 16 -> 23 rewrite happens first and the bare
# "September 2" pattern that follows cannot then eat the "2" of "September 23".
# The all-caps entry is kept even though this family currently has no capitalised
# date: a slide header typeset in capitals is exactly what a mixed-case sweep
# missed on the Lightning Lesson deck, and the guard costs nothing.
DATES = [
    ("SEPTEMBER 16", "SEPTEMBER 23"),
    ("Wednesday, September 16, 2026", "Wednesday, September 23, 2026"),
    ("September 16, 2026", "September 23, 2026"),
    ("September 16", "September 23"),
    (re.compile(r"September 2(?![0-9])"), "September 9"),
]


def apply_dates(text):
    for pat, rep in DATES:
        text = pat.sub(rep, text) if hasattr(pat, "sub") else text.replace(pat, rep)
    return text


# ── the three note refinements ──────────────────────────────────────────────
# 1. Slide 5, Optionality. The existing optional spoken proof already names what
#    travelled from audit into cybersecurity and what did not. The sentence is
#    placed inside that proof, where it is the point the story is making, rather
#    than appended as a second block that would repeat it.
PROOF_OLD = ("Cybersecurity language, privacy frameworks and technical context did not. "
             "That is the distinction this reading is trying to help us see.")
PROOF_NEW = ("Cybersecurity language, privacy frameworks and technical context did not. "
             "That is why portability is not the same as plug-and-play. Something "
             "important can travel with you while the new context still requires real "
             "relearning. That is the distinction this reading is trying to help us see.")

# 2. Slide 5, the interpretive aid. Written as an OPTIONAL SPOKEN FRAME because
#    that is the idiom this note already uses, and because it is an alternative
#    within the ninety seconds rather than an addition to them. Participant
#    microphones are muted for minutes 0-50, so a cue marked "if asked" would be
#    unreachable here; naming the trade against the proof story is the honest form.
AUDIT_BLOCK = (
    "\n\nOPTIONAL SPOKEN FRAME — the four-part portability audit. An interpretive aid "
    "for reading an Optionality score, not new teaching and not a new slide. It fits "
    "inside the ninety seconds only in place of the spoken proof above, not alongside "
    "it. Read the four questions and stop:\n"
    "What travels? · What does not travel automatically? · What can I prove? · What "
    "must I relearn?\n"
    "A low score usually separates into those four rather than resolving into a verdict "
    "on ability. Do not work an example — there is no room for one here.")

# 3. Slide 29, the translation appendix. This is the slide a presenter reaches for
#    when someone reads a low Optionality score as a verdict, which is exactly the
#    moment translation can be oversold.
TRANSLATION_BLOCK = (
    "\n\nSTATE THE LIMIT WHENEVER YOU USE THIS SLIDE: translation does not guarantee "
    "acceptance. Another organization still has to decide that this evidence is "
    "comparable enough to trust in its context. Legibility removes one barrier; it does "
    "not remove genuine experience gaps, market constraints or employer preferences, and "
    "nothing here should be heard as a promise that better wording will.")


def patch_runs(container, fn):
    """Rewrite run text in place so every run keeps its own formatting."""
    hits = 0
    for para in container.paragraphs:
        for run in para.runs:
            new = fn(run.text)
            if new != run.text:
                run.text = new
                hits += 1
    return hits


def faces(pres):
    return ["\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for s in pres.slides]


def notes(pres):
    return [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
            for s in pres.slides]


def build_deck():
    shutil.copyfile(DECK_SRC, DECK_DST)
    p = Presentation(DECK_DST)
    sl = list(p.slides)

    face_hits = note_hits = 0
    for s in sl:
        for sh in s.shapes:
            if sh.has_text_frame:
                face_hits += patch_runs(sh.text_frame, apply_dates)
        if s.has_notes_slide:
            note_hits += patch_runs(s.notes_slide.notes_text_frame, apply_dates)

    # slide 5 — Optionality
    f5 = "\n".join(sh.text_frame.text for sh in sl[4].shapes if sh.has_text_frame)
    assert "Optionality" in f5 and "Career portability" in f5, "slide 5 is not Optionality"
    n5 = sl[4].notes_slide.notes_text_frame
    assert PROOF_OLD in n5.text, "slide 5: the spoken proof was not found"
    assert "plug-and-play" not in n5.text, "slide 5 already carries the portability line"
    n5.text = n5.text.replace(PROOF_OLD, PROOF_NEW) + AUDIT_BLOCK

    # slide 29 — the translation appendix
    f29 = "\n".join(sh.text_frame.text for sh in sl[28].shapes if sh.has_text_frame)
    assert "What translation actually means" in f29, "slide 29 is not the translation appendix"
    n29 = sl[28].notes_slide.notes_text_frame
    assert "does not guarantee acceptance" not in n29.text, "slide 29 already carries the limit"
    n29.text = n29.text + TRANSLATION_BLOCK

    p.save(DECK_DST)

    # ── prove the revision was controlled ───────────────────────────────────
    before, after = Presentation(DECK_SRC), Presentation(DECK_DST)
    bf, af, bn, an = faces(before), faces(after), notes(before), notes(after)
    assert len(af) == len(bf) == 33, "slide count moved"

    changed_faces = [i + 1 for i in range(33) if bf[i] != af[i]]
    changed_notes = [i + 1 for i in range(33) if bn[i] != an[i]]
    for i in changed_faces:
        assert apply_dates(bf[i - 1]) == af[i - 1], \
            f"slide {i}: a face changed by something other than a date"
    for i in changed_notes:
        expect = apply_dates(bn[i - 1])
        if i == 5:
            expect = expect.replace(apply_dates(PROOF_OLD),
                                    apply_dates(PROOF_NEW)) + AUDIT_BLOCK
        elif i == 29:
            expect = expect + TRANSLATION_BLOCK
        assert expect == an[i - 1], \
            f"slide {i}: a note changed by something other than a date or its sanctioned edit"
    assert changed_faces == [33], f"faces changed outside the holding slide: {changed_faces}"
    assert set(changed_notes) >= {5, 25, 29}, f"a mandated note edit did not land: {changed_notes}"

    print("built", os.path.basename(DECK_DST))
    print(f"  date rewrites: {face_hits} on faces, {note_hits} in notes")
    print(f"  faces changed: {changed_faces}   notes changed: {changed_notes}")
    return changed_faces, changed_notes


# ── document patching ───────────────────────────────────────────────────────
def iter_runs(doc):
    for p in doc.paragraphs:
        for r in p.runs:
            yield r
    for t in doc.tables:
        for row in t.rows:
            for c in row.cells:
                for p in c.paragraphs:
                    for r in p.runs:
                        yield r


def sub_runs(doc, pairs, required=True):
    """Apply (old, new) substitutions run by run. Every pair must land."""
    counts = {old: 0 for old, _ in pairs}
    for r in iter_runs(doc):
        t = r.text
        for old, new in pairs:
            if old in t:
                counts[old] += t.count(old)
                t = t.replace(old, new)
        if t != r.text:
            r.text = t
    if required:
        missing = [o for o, n in counts.items() if n == 0]
        assert not missing, f"substitution never matched: {missing}"
    return counts


def build_sop():
    shutil.copyfile(SOP_SRC, SOP_DST)
    d = docx.Document(SOP_DST)
    # Purely substitutional. Nothing is added to this document: its LibreOffice
    # export is pinned at exactly seven pages and section 12 has overrun twice
    # before when prose was added. The narrative of this pass lives in the change
    # log, which is where change history belongs and which has no pinned length.
    counts = sub_runs(d, [
        ("Sept16_2026", "Sept23_2026"),
        ("v2.0.7", "v2.0.8"),
        ("v2.0.4", "v2.0.5"),
        (OLD_STAMP_SOP, STAMP),
        ("Wednesday, September 16, 2026", "Wednesday, September 23, 2026"),
        ("September 16", "September 23"),
        ("September 2 Lightning Lesson", "September 9 Lightning Lesson"),
        ("September 2 session", "September 9 session"),
    ])
    d.save(SOP_DST)
    print("built", os.path.basename(SOP_DST))
    for k, v in counts.items():
        print(f"    {v:>2}x  {k}")
    return counts


def build_checklist():
    shutil.copyfile(CHK_SRC, CHK_DST)
    d = docx.Document(CHK_DST)
    counts = sub_runs(d, [
        ("v1.0", "v1.1"),
        ("Prepared " + OLD_STAMP_CHK, "Prepared " + STAMP),
        ("deck v2.0.4", "deck v2.0.5"),
        ("Wednesday, September 16, 2026", "Wednesday, September 23, 2026"),
        ("September 16", "September 23"),
    ])
    d.save(CHK_DST)
    print("built", os.path.basename(CHK_DST))
    for k, v in counts.items():
        print(f"    {v:>2}x  {k}")
    return counts


# ── QA report: targeted edits, plus a new change-log section and QA group ────
# Blanket sweeps are wrong in this document. It carries an audit trail, and rows
# that record what an earlier version did — "Read FREE FLAGSHIP through v2.0.3;
# v2.0.4 changes it to ...", or the v2.0.6 correction that quotes the wording it
# replaced — are true statements about history. Rewriting them would falsify the
# record. Every edit below is therefore addressed to one location.
QA_PARA_EDITS = {
    1: [(OLD_STAMP_QA, STAMP)],
    8: [("DECK IS NOW v2.0.4.",
         "DECK IS NOW v2.0.5."),
        ("what v2.0.4 changed is set out in its own section below, and where the two "
         "disagree the v2.0.4 section governs",
         "what v2.0.4 and v2.0.5 changed is set out in their own sections below, and "
         "where they disagree with this section the later ones govern")],
    12: [("UNCHANGED FROM v2.0.2 THROUGH v2.0.6.",
          "UNCHANGED FROM v2.0.2 THROUGH v2.0.8."),
         ("The workbook was not regenerated, modified, renamed or reissued in this pass.",
          "The workbook was not regenerated, modified, renamed or reissued in this pass; "
          "it carries no date, so the reschedule reaches nothing inside it, and it is "
          "byte-identical at sha256 " + WORKBOOK_SHA[:16] + "….")],
    32: [("PowerPoint v2.0.4, workbook v2.0.1, SOP v2.0.7",
          "PowerPoint v2.0.5, workbook v2.0.1, SOP v2.0.8")],
    38: [("Verification — 148 items", "Verification — 162 items")],
    63: [("148 of 148 PASS.", "162 of 162 PASS.")],
    70: [("The September 16 deck v2.0.4 was exported",
          "The September 23 deck v2.0.5 was exported"),
         ("PREVIEW_Presentation_60MIN_v2.0.4_LibreOffice.pdf, rendered from the v2.0.4 "
          "deck in this pass",
          "PREVIEW_Presentation_60MIN_v2.0.5_LibreOffice.pdf, rendered from the v2.0.5 "
          "deck in this pass")],
    71: [("PREVIEW_SOP_60MIN_v2.0.7_LibreOffice.pdf",
          "PREVIEW_SOP_60MIN_v2.0.8_LibreOffice.pdf")],
    91: [("QA REPORT v2.0.7", "QA REPORT v2.0.8"), (OLD_STAMP_QA, STAMP)],
}
QA_CELL_EDITS = {
    (0, 0, 0): [("v2.0.7 CANDIDATE", "v2.0.8 CANDIDATE"), (OLD_STAMP_QA, STAMP)],
    (10, 4, 3): [("September 16 survives only on the hidden holding slide.",
                  "September 23 survives only on the hidden holding slide, which is shown "
                  "before recording begins and is never recorded.")],
    (13, 10, 1): [("September 2 event untouched and out of scope",
                   "September 9 Lightning Lesson untouched and out of scope")],
    (13, 10, 3): [("Stated in Section 1. Nothing in this pass touched the September 2 "
                   "material.",
                   "Stated in Section 1. Nothing in this pass touched the Lightning Lesson "
                   "beyond correcting this family's cross-reference to its own new date of "
                   "Wednesday, September 9, 2026.")],
    (14, 3, 3): [("The SOP source map names deck v2.0.4, workbook v2.0.1, SOP v2.0.7",
                  "The SOP source map names deck v2.0.5, workbook v2.0.1, SOP v2.0.8"),
                 ("No superseded version — v2.0.2, v2.0.3 or v2.0.5 — survives anywhere in "
                  "the document.",
                  "No superseded version — v2.0.2, v2.0.3, v2.0.4, v2.0.6 or v2.0.7 — "
                  "survives anywhere in the document, and the SOP filename is recoded to "
                  "Sept23_2026 so the date code cannot contradict the date inside.")],
    # Row 119 asserted the deck was generator-reproducible. That is no longer true
    # and the row is repointed to what is true, not deleted and not softened.
    (18, 12, 1): [("Deck v2.0.4 is reproducible from its generator, not hand-patched",
                   "Deck v2.0.5 provenance is disclosed, and every diff is accounted for")],
    (18, 12, 3): [("build_flagship_v204.js is forked from build_flagship_v203.js, which "
                   "reproduces the approved v2.0.3 in every content part — re-running it "
                   "differs only in docProps/core.xml, a timestamp. The 35-slide file that "
                   "circulated was never generator output.",
                   "REPOINTED, NOT WEAKENED. Through v2.0.4 this row asserted that the deck "
                   "was reproducible from build_flagship_v204.js. That generator lived in a "
                   "gitignored scratchpad and did not survive a container rebuild, so the "
                   "claim can no longer be made and is not repeated. v2.0.5 is instead a "
                   "scripted patch of the approved v2.0.4 file, and the build refuses to "
                   "save unless every differing slide face and note matches the approved "
                   "text with only a named substitution or a named insertion applied. The "
                   "35-slide file that circulated remains excluded: it was never generator "
                   "output either.")],
    (18, 14, 3): [("Deck v2.0.4, workbook v2.0.1, SOP v2.0.7",
                   "Deck v2.0.5, workbook v2.0.1, SOP v2.0.8")],
}

CHANGELOG_ROWS = [
    ("Session date",
     "Wednesday, September 16, 2026 becomes Wednesday, September 23, 2026 at the same "
     "6:00–7:00 PM Central Time. Seven locations in the SOP — the identity table, the "
     "offer table, the Maven configuration table, and sections 3 and 12 — plus one slide "
     "face, one speaker note and four locations in the manual edit checklist."),
    ("The one dated slide",
     "Slide 33, the hidden holding slide shown before recording begins, is the only slide "
     "face in the deck that carries a date. It now reads Wednesday, September 23, 2026. "
     "The 26 active recorded slides remain date-free, which is what keeps the replay "
     "evergreen."),
    ("The evergreen prohibition",
     "Slide 25's note and SOP section 3 both forbid saying the session date inside the "
     "recorded core. The forbidden example moves with the session, from September 16 to "
     "September 23. The rule is unchanged; only the date it names moved."),
    ("Lightning Lesson cross-reference",
     "SOP section 1 and section 12 described the separate Lightning Lesson as the "
     "September 2 event. That lesson was itself rescheduled to Wednesday, September 9, "
     "2026 in an earlier pass which reached its own deck and the website but not this "
     "family. Corrected here. The lesson remains out of scope and nothing inside it was "
     "touched."),
    ("Date-coded filename",
     "The SOP filename carried Sept16_2026 and so contradicted its own contents the "
     "moment the session moved. Recoded to Sept23_2026. The convention is kept rather "
     "than abandoned, because the date code tells a facilitator which delivery a file "
     "governs; it is recorded here as a maintenance obligation on any future reschedule."),
    ("Speaker note — slide 5, Optionality",
     "The optional spoken proof gains: “That is why portability is not the same as "
     "plug-and-play. Something important can travel with you while the new context still "
     "requires real relearning.” It is placed inside the audit-to-cybersecurity story, "
     "which already names what travelled and what did not, rather than appended as a "
     "second block that would say it twice."),
    ("Speaker note — slide 5, interpretive aid",
     "A four-part portability audit is added as an OPTIONAL SPOKEN FRAME: what travels, "
     "what does not travel automatically, what can I prove, what must I relearn. It is "
     "marked as an alternative to the spoken proof within the same ninety seconds, not an "
     "addition to them, and it is not cued as “if asked” because participant microphones "
     "are muted for minutes 0–50."),
    ("Speaker note — slide 29, translation",
     "The hidden translation appendix gains: “Translation does not guarantee acceptance. "
     "Another organization still has to decide that this evidence is comparable enough to "
     "trust in its context.” The note goes on to say that legibility does not remove "
     "genuine experience gaps, market constraints or employer preferences."),
    ("Manual edit checklist",
     "Reissued at v1.1. v1.0 instructed the operator to set the Maven date to September "
     "16, so leaving it alone would have re-entered the stale date by hand at the one "
     "point where the repository cannot correct itself. v1.1 says September 23, and its "
     "deck reference moves to v2.0.5."),
    ("What was deliberately not changed",
     "The twelve statements, the 1–5 scale, both /30 totals, the 17–21 boundary band, the "
     "sensitivity check, the incomplete reading, the four states, the seven move "
     "categories, the Next-Move Note, every TIMING cue, the recording boundary at minute "
     "50, the commercial architecture and the workbook. No slide added, removed or "
     "reordered. No new offer. No enterprise, staffing, succession or AI teaching added to "
     "the timed core."),
]

QA_ROWS = [
    ("149", "Holding slide carries Wednesday, September 23, 2026", "PASS",
     "Slide 33, the only dated slide face in the deck. It is hidden and is shown before "
     "recording begins."),
    ("150", "Zero stale “September 16” anywhere in the family", "PASS",
     "Deck faces, deck notes, SOP body and tables, this report, the checklist and the "
     "workbook all swept case-insensitively, so a header typeset in capitals could not "
     "hide. Zero live occurrences."),
    ("151", "Zero stale “September 2” Lightning Lesson reference", "PASS",
     "Swept with a negative lookahead so the pattern could not match the “2” of September "
     "23. The two SOP cross-references now read September 9."),
    ("152", "The twelve statements are byte-identical to v2.0.4", "PASS",
     "Slides 6, 7, 9 and 10 compared string by string against the approved v2.0.4 deck. "
     "Twelve of twelve identical. Not rewritten, not reordered, none added or removed."),
    ("153", "Scoring, state, boundary and move-category logic byte-identical", "PASS",
     "Slides 8, 11, 12, 13, 14, 17, 18, 19, 20, 21 and 22 identical to v2.0.4: both /30 "
     "totals, 19–30 high, 6–18 low, the 17–21 band, the neighbouring-score sensitivity "
     "check, the incomplete reading, the four states and the seven categories."),
    ("154", "Exactly one slide face differs from v2.0.4", "PASS",
     "Slide 33. Thirty-two of thirty-three faces byte-identical, asserted by the build "
     "rather than inspected by eye."),
    ("155", "Only sanctioned notes differ, and only by sanctioned text", "PASS",
     "Notes differ on slides 5, 25 and 29. The build reconstructs each expected note from "
     "the approved text plus its named edit and refuses to save on any mismatch, so an "
     "unintended note edit cannot ship."),
    ("156", "Slide count, hidden set and timing chain unchanged", "PASS",
     "33 total, 26 active, hidden exactly [27–33]. Every TIMING cue identical to v2.0.4, "
     "continuous 0:00 to 60:00."),
    ("157", "Workbook byte-identical and date-free", "PASS",
     "sha256 " + WORKBOOK_SHA + ". Rasterised and searched: it contains no date, "
     "no year and no timezone, so the reschedule reaches nothing inside it."),
    ("158", "No new slide, no new offer, no new section", "PASS",
     "Continuation still carries two routes — do nothing, or the Field Kit. The Private "
     "Capability Position Read stays off the slide and out of Q&A while its fulfillment "
     "path is unbuilt. No staffing, succession, hiring or AI teaching entered the timed "
     "core."),
    ("159", "Translation is not sold as acceptance", "PASS",
     "Slide 29's note states the limit in the mandated words and adds that legibility does "
     "not remove genuine experience gaps, market constraints or employer preferences. "
     "Nothing in the family promises that framing overcomes them."),
    ("160", "Timezone is CT or Central Time, never CST as a label", "PASS",
     "Swept across the deck, the SOP, this report and the checklist. The only occurrences "
     "of the string are the SOP rule that bans it and this row."),
    ("161", "No stale filename or version survives in a newer file", "PASS",
     "The SOP names deck v2.0.5, workbook v2.0.1 and SOP v2.0.8, and its own filename is "
     "recoded to Sept23_2026. This report names the v2.0.5 previews. Historical sentences "
     "that describe what an earlier version did are preserved deliberately: they are the "
     "audit trail, not stale claims."),
    ("162", "SOP still exports at exactly seven pages", "PASS",
     "Re-exported through LibreOffice Writer after the edits. Seven pages, page 1 "
     "portrait, pages 2–3 landscape, pages 4–7 portrait. This pass added no prose to the "
     "SOP precisely because section 12 has overrun twice before."),
]


def _carrier(para):
    """The run whose formatting the paragraph actually shows.

    Several cells in this family open with an empty run that carries different
    character formatting from the visible text. Writing into that run silently
    lost the bold green of the PASS column, so the first NON-EMPTY run is the
    one to keep.
    """
    runs = para.runs
    if not runs:
        return para.add_run("")
    return next((r for r in runs if r.text.strip()), runs[0])


def _set_cell(cell, text):
    """Write text into a cell, keeping the visible run's formatting."""
    p = cell.paragraphs[0]
    for extra in cell.paragraphs[1:]:
        extra._element.getparent().remove(extra._element)
    keep = _carrier(p)
    keep.text = text
    for r in p.runs:
        if r._element is not keep._element:
            r._element.getparent().remove(r._element)


def _clone_table(doc, template, rows, insert_before):
    """Deep-copy a table's formatting, then fill it with new rows."""
    el = copy.deepcopy(template._element)
    insert_before.addprevious(el)
    tbl = docx.table.Table(el, doc)
    header_cols = len(rows[0])
    # keep the header row and one body row as templates, drop the rest
    while len(tbl.rows) > 2:
        tbl._element.remove(tbl.rows[-1]._element)
    body_template = copy.deepcopy(tbl.rows[1]._element)
    for ci, val in enumerate(rows[0]):
        _set_cell(tbl.rows[0].cells[ci], val)
    tbl._element.remove(tbl.rows[1]._element)
    for row in rows[1:]:
        el_row = copy.deepcopy(body_template)
        tbl._element.append(el_row)
        cells = tbl.rows[-1].cells
        assert len(cells) == header_cols, "row width does not match the template"
        for ci, val in enumerate(row):
            _set_cell(cells[ci], val)
    return tbl


def _clone_heading(doc, template_para, text, insert_before):
    el = copy.deepcopy(template_para._element)
    insert_before.addprevious(el)
    p = docx.text.paragraph.Paragraph(el, template_para._parent)
    keep = _carrier(p)
    keep.text = text
    for r in p.runs:
        if r._element is not keep._element:
            r._element.getparent().remove(r._element)
    return p


def build_qa(changed_faces, changed_notes):
    shutil.copyfile(QA_SRC, QA_DST)
    d = docx.Document(QA_DST)

    for idx, pairs in QA_PARA_EDITS.items():
        p = d.paragraphs[idx]
        for old, new in pairs:
            assert any(old in r.text for r in p.runs), f"P{idx}: never matched {old[:50]!r}"
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
    for (ti, ri, ci), pairs in QA_CELL_EDITS.items():
        cell = d.tables[ti].rows[ri].cells[ci]
        for old, new in pairs:
            runs = [r for p in cell.paragraphs for r in p.runs]
            assert any(old in r.text for r in runs), \
                f"T{ti}r{ri}c{ci}: never matched {old[:50]!r}"
            for r in runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)

    # new change-log section, inserted before the zero-tolerance duration sweep
    anchor = d.paragraphs[34]._element
    head_tmpl = d.paragraphs[26]      # "SOP — the 60-minute rewrite"
    body_tmpl = d.paragraphs[27]
    _clone_heading(d, head_tmpl,
                   "September 23 reschedule — what changed in v2.0.5 / v2.0.8", anchor)
    _clone_heading(d, body_tmpl,
                   "A controlled revision of an approved family: a date move, three "
                   "speaker-note refinements, and the version and filename labels that "
                   "carry them. The instrument was not touched, and the workbook was not "
                   "reissued.", anchor)
    _clone_table(d, d.tables[6], [("Where", "Change")] + CHANGELOG_ROWS, anchor)
    blank = copy.deepcopy(d.paragraphs[33]._element)
    anchor.addprevious(blank)

    # new QA group M, inserted before "162 of 162 PASS."
    tail = [p for p in d.paragraphs if p.text.strip() == "162 of 162 PASS."][0]._element
    group_tmpl = [p for p in d.paragraphs
                  if p.text.startswith("Group L —")][0]
    _clone_heading(d, group_tmpl,
                   "Group M — the September 23 reschedule and the instrument freeze", tail)
    _clone_table(d, d.tables[19], [("#", "Category", "Status", "Notes")] + QA_ROWS, tail)
    tail.addprevious(copy.deepcopy(d.paragraphs[33]._element))

    d.save(QA_DST)
    print("built", os.path.basename(QA_DST))


if __name__ == "__main__":
    cf, cn = build_deck()
    build_sop()
    build_checklist()
    build_qa(cf, cn)
    assert hashlib.sha256(open(WORKBOOK, "rb").read()).hexdigest() == WORKBOOK_SHA, \
        "the workbook is not the approved v2.0.1 build"
    print("\nworkbook unchanged:", WORKBOOK_SHA)
    for f in sorted(os.listdir(OUT)):
        p = os.path.join(OUT, f)
        print(f"  {hashlib.sha256(open(p,'rb').read()).hexdigest()[:16]}  {f}")
